from flask import Flask, g, render_template, jsonify, request, send_file
import json
import os
import time
import datetime
from pathlib import Path
from gartner_app.api.health import health_blueprint
from gartner_app.api.research import research_blueprint
from gartner_app.couchdb.client import CouchDBConflict
from gartner_app.repositories.datasets import (
    MigrationWriteBlocked,
    RevisionPreconditionRequired,
    build_dataset_repository,
)
from gartner_app.repositories.json_backend import LegacyJsonRepository

app = Flask(__name__)
app.config['TEMPLATES_AUTO_RELOAD'] = True
app.config['MAX_CONTENT_LENGTH'] = int(
    os.getenv('MAX_REQUEST_BYTES', str(64 * 1024 * 1024))
)
app.register_blueprint(health_blueprint)
app.register_blueprint(research_blueprint)
legacy_json_repository = LegacyJsonRepository(Path(__file__).resolve().parent)
dataset_repository = build_dataset_repository(
    Path(__file__).resolve().parent,
    legacy_json_repository,
)


@app.before_request
def validate_api_json_content_type():
    """Reject non-JSON bodies on state-changing API requests."""
    if (
        request.path.startswith('/api/')
        and request.method in {'POST', 'PUT', 'PATCH'}
        and (request.content_length or 0) > 0
        and not request.is_json
    ):
        return jsonify({
            'error': 'Request body must use application/json',
            'code': 'unsupported_media_type',
        }), 415


def read_dataset(source_path):
    """Read through the selected backend and expose its current revision."""
    value = dataset_repository.read_document(source_path)
    if request.method == 'GET':
        g.dataset_revision = dataset_repository.revision(source_path)
    return value


def persist_dataset(source_path, value):
    """Persist through the selected backend and map migration write states."""
    try:
        result = dataset_repository.write_document(
            source_path,
            value,
            expected_revision=request.headers.get('If-Match'),
        )
        revision = result.get('revision')
        if revision:
            g.dataset_revision = revision
    except MigrationWriteBlocked as exc:
        return jsonify({
            'error': str(exc),
            'code': 'compare_mode_read_only',
        }), 503
    except CouchDBConflict:
        return jsonify({
            'error': 'The document changed while it was being edited. Reload and retry.',
            'code': 'revision_conflict',
        }), 409
    except RevisionPreconditionRequired as exc:
        return jsonify({
            'error': str(exc),
            'code': 'if_match_required',
        }), 428
    return None


def dataset_json_response(value, source_path):
    """Return JSON with the current source revision as a strong ETag."""
    response = jsonify(value)
    response.set_etag(dataset_repository.revision(source_path))
    return response


@app.after_request
def add_write_revision_etag(response):
    revision = getattr(g, 'dataset_revision', None)
    if revision and not response.get_etag()[0]:
        response.set_etag(revision)
    return response

# Remove injected debug code from responses
@app.after_request
def remove_debug_code(response):
    if response.mimetype == 'text/html':
        data = response.get_data(as_text=True)
        # Remove debug script and panel
        data = data.replace(
            '''    <script>
        // Capture console.log output and display on page
        const consoleLogs = [];
        const originalLog = console.log;
        const originalError = console.error;
        
        console.log = function(...args) {
            originalLog.apply(console, args);
            consoleLogs.push(args.join(' '));
            // Also update the page if debug panel exists
            const debugPanel = document.getElementById('debug-panel-logs');
            if (debugPanel) {
                debugPanel.innerHTML = consoleLogs.slice(-20).join('<br>');
                debugPanel.scrollTop = debugPanel.scrollHeight;
            }
        };
        
        console.error = function(...args) {
            originalError.apply(console, args);
            consoleLogs.push('❌ ' + args.join(' '));
        };
        
        console.log('📄 HTML body loaded - inline test');
        window.addEventListener('error', (e) => {
            console.error('❌ Global error:', e.message, e.filename, e.lineno);
        });
    </script>''',
            ''
        )
        response.set_data(data)
    return response

# Global state for vendor file selection
class AppState:
    current_vendor_file = 'Vendor 3-7.json'
    current_schema_file = 'schema3-3.json'

app_state = AppState()

with (Path(__file__).resolve().parent / 'migration' / 'canonical_sources.json').open(
    'r', encoding='utf-8'
) as _catalog_file:
    _CANONICAL_CATALOG = json.load(_catalog_file)['sources']


def canonical_paths(*kinds):
    """Return active paths from the source-controlled dataset catalog."""
    return sorted(
        entry['path']
        for entry in _CANONICAL_CATALOG
        if entry.get('status') == 'active' and entry.get('kind') in kinds
    )

# ── Schema registry ──────────────────────────────────────────────────
# Maps schema filenames to their top-level JSON key and internal structure type.
# "nested" = v3.2 style (pillars → sub_capabilities list)
# "flat"   = v4.0 / v5.0 style (separate sub_pillars dict keyed by ID)
SCHEMA_REGISTRY = {
    'schema3-3.json':          {'top_key': 'dfir_capability_taxonomy_v3.2',        'structure': 'nested'},
    'schema4-0_enhanced.json': {'top_key': 'dfir_capability_taxonomy_v4.0_enhanced', 'structure': 'flat'},
    'schema5-0_ai.json':       {'top_key': 'dfir_capability_taxonomy_v5.0_ai',     'structure': 'flat'},
    'AI TriSM Schema 1_0.json': {'top_key': 'ai_trism_taxonomy_v1.0',              'structure': 'flat'},
    'AI TriSM Schema 1_1.json': {'top_key': 'ai_trism_taxonomy_v1.1',              'structure': 'flat'},
    'Preemptive_Cybersecurity_Schema.json': {'top_key': 'preemptive_cybersecurity_taxonomy_v1.0', 'structure': 'flat'},
    'Preemptive_Cybersecurity_Schema_v2.json': {'top_key': 'preemptive_cybersecurity_taxonomy_v2.0', 'structure': 'flat'},
    'Preemptive_Cybersecurity_Schema_v3.json': {'top_key': 'preemptive_cybersecurity_taxonomy_v3.0', 'structure': 'flat'},
    'Secure_by_Design_AI_Controls_Schema.json': {'top_key': 'secure_by_design_ai_controls_v1.0', 'structure': 'flat'},
    'Secure_by_Design_AI_Controls_Schema_v2.json': {'top_key': 'secure_by_design_ai_controls_v2.0', 'structure': 'flat'},
    'MDR_Services_Schema.json': {'top_key': 'mdr_services_taxonomy_v1.0', 'structure': 'flat'},
    'Offensive_Security_Schema.json': {'top_key': 'offensive_security_taxonomy_v1.0', 'structure': 'flat'},
    'Product Market Readiness Schema 1_0.json': {'top_key': 'product_market_readiness_taxonomy_v1.0', 'structure': 'flat'},
    'MDR_MQ_Gap_Schema_App.json': {'top_key': 'mq_gap_taxonomy_v1.0', 'structure': 'flat'},
    'CNAPP_Schema.json': {'top_key': 'cnapp_taxonomy_v1.1', 'structure': 'flat'},
    'CNAPP_MQ_Gap_Schema_App.json': {'top_key': 'cnapp_mq_gap_taxonomy_v1.0', 'structure': 'flat'},
    'Schema_Template_Capability.json': {'top_key': 'capability_schema_template_v1.0', 'structure': 'flat'},
    'Schema_Template_MQ_Gap.json': {'top_key': 'mq_gap_schema_template_v1.0', 'structure': 'flat'},
    'agentic_soc_framework_v1.json': {'top_key': None, 'structure': 'asmf'},
    'agentic_enterprise_operations_framework_v1.json': {'top_key': None, 'structure': 'asmf'},
    'AI_platform_ecosystem_framework_v1.json': {'top_key': None, 'structure': 'asmf'},
}

# Schema display metadata: maps schema filename to title, abbreviation, subtitle
SCHEMA_DISPLAY = {
    'schema3-3.json':              {'title': 'DFIR Vendor Marketplace Analysis 2026', 'abbr': 'DFIR', 'subtitle': 'Filter and analyze incident response vendors by capabilities and specializations'},
    'schema4-0_enhanced.json':     {'title': 'DFIR Vendor Marketplace Analysis 2026', 'abbr': 'DFIR', 'subtitle': 'Filter and analyze incident response vendors by capabilities and specializations'},
    'schema5-0_ai.json':           {'title': 'DFIR Vendor Marketplace Analysis 2026', 'abbr': 'DFIR', 'subtitle': 'Filter and analyze incident response vendors by capabilities and specializations'},
    'AI TriSM Schema 1_0.json':    {'title': 'AI TRiSM Vendor Capability Analysis 2026', 'abbr': 'TRiSM', 'subtitle': 'Evaluate AI Trust, Risk, and Security Management capabilities across vendors'},
    'AI TriSM Schema 1_1.json':    {'title': 'AI TRiSM Vendor Capability Analysis 2026', 'abbr': 'TRiSM', 'subtitle': 'Evaluate AI Trust, Risk, and Security Management capabilities across vendors'},
    'Preemptive_Cybersecurity_Schema.json': {'title': 'Preemptive Cybersecurity Vendor Analysis 2026', 'abbr': 'PreCyber', 'subtitle': 'Evaluate preemptive cybersecurity capabilities — proactive defense that prevents threats before exploitation'},
    'Preemptive_Cybersecurity_Schema_v2.json': {'title': 'Preemptive Cybersecurity Vendor Analysis 2026 v2', 'abbr': 'PreCyber', 'subtitle': 'Evaluate preemptive cybersecurity capabilities with Services Maturity & Pricing evaluation'},
    'Secure_by_Design_AI_Controls_Schema.json': {'title': 'Secure-by-Design AI Controls Maturity Assessment 2026', 'abbr': 'SbD-AI', 'subtitle': 'Self-assessment maturity framework for secure AI product and service capabilities'},
    'Secure_by_Design_AI_Controls_Schema_v2.json': {'title': 'SbD-AI + AIUC-1 Controls Maturity Assessment 2026', 'abbr': 'SbD-AI v2', 'subtitle': 'Extended maturity framework with full AIUC-1 compliance mapping — 7 pillars, 40 sub-pillars'},
    'MDR_Services_Schema.json': {'title': 'MDR Services Vendor Capability & Pricing Analysis 2026', 'abbr': 'MDR', 'subtitle': 'Evaluate MDR service capabilities and composable pricing model maturity across providers'},
    'Offensive_Security_Schema.json': {'title': 'Offensive Security Vendor Analysis 2026', 'abbr': 'OffSec', 'subtitle': 'Evaluate offensive security, CTEM, vulnerability management, and DevSecOps capabilities'},
    'Product Market Readiness Schema 1_0.json': {'title': 'Product Market Readiness — Credibility Gap Analysis 2026', 'abbr': 'PMR', 'subtitle': 'Evaluate vendor go-to-market claims vs. proof of execution across 208 cybersecurity vendors'},
    'MDR_MQ_Gap_Schema_App.json': {'title': 'MDR Magic Quadrant — Gap Criteria Analysis 2026', 'abbr': 'MDR-MQ', 'subtitle': 'Evaluate MDR vendors against 7 MQ criteria not covered by capability or pricing schemas'},
    'CNAPP_Schema.json': {'title': 'CNAPP Vendor Capability Analysis 2026', 'abbr': 'CNAPP', 'subtitle': 'Evaluate cloud-native application protection platform vendors across CSPM, CWPP, CIEM, DevSecOps, CDR, DSPM, and fringe differentiators'},
    'CNAPP_MQ_Gap_Schema_App.json': {'title': 'CNAPP Magic Quadrant — Gap Criteria Analysis 2026', 'abbr': 'CNAPP-MQ', 'subtitle': 'Evaluate CNAPP vendors against the 7 Magic Quadrant criteria not covered by the CNAPP capability schema'},
    'Schema_Template_Capability.json': {'title': 'Schema Template — Capability Assessment', 'abbr': 'TEMPLATE-CAP', 'subtitle': 'Blank capability schema template with annotated structure — use as the starting point for a new market capability assessment schema'},
    'Schema_Template_MQ_Gap.json': {'title': 'Schema Template — MQ Gap Criteria', 'abbr': 'TEMPLATE-MQ', 'subtitle': 'Blank MQ Gap schema template with annotated structure — use as the starting point for Magic Quadrant supplemental criteria scoring'},
    'agentic_soc_framework_v1.json': {'title': 'Agentic Security Operations Adoption Framework 2026', 'abbr': 'ASAF', 'subtitle': 'Vendor-neutral adoption framework for autonomous security operations — 11 dimensions, 44 sub-dimensions, 6 stages'},
    'agentic_enterprise_operations_framework_v1.json': {'title': 'Agentic Enterprise Operations Framework 2026', 'abbr': 'AEOF', 'subtitle': 'Enterprise operations framework for agentic business governance, orchestration, and risk assurance.'},
        'AI_platform_ecosystem_framework_v1.json': {'title': 'AI Platform Ecosystem Framework 2026', 'abbr': 'APEF', 'subtitle': 'Compare seven major AI platform providers across the enterprise AI value chain.'},
}

def discover_schema_files():
    """Return list of available schema JSON files."""
    return [
        path for path in canonical_paths('schema', 'framework')
        if path in SCHEMA_REGISTRY
    ]

def _framework_capabilities(schema_file):
    """Describe optional report features declared by a framework schema."""
    body = load_schema_data(schema_file)
    dimensions = body.get('dimensions', {}) if isinstance(body, dict) else {}
    return {
        'maturity_stages': bool(body.get('maturity_stages')),
        'stage_descriptors': any(
            isinstance(dim, dict) and any(
                isinstance(sd, dict) and sd.get('stage_descriptors')
                for sd in (dim.get('sub_dimensions', {}) or {}).values()
            ) for dim in dimensions.values()
        ) if isinstance(dimensions, dict) else False,
        'transformation_journey': bool(body.get('transformation_journey')),
        'weights': any(isinstance(dim, dict) and isinstance(dim.get('weight'), (int, float)) for dim in dimensions.values()) if isinstance(dimensions, dict) else False,
        'relationships': bool(
            body.get('relationships', {}).get('edges')
            if isinstance(body.get('relationships'), dict)
            else body.get('relationships')
        ),
    }

# Load schema data for sub-pillar definitions
def load_schema_data(schema_file=None):
    """Load schema body (the dict under the top-level key) for any version."""
    if schema_file is None:
        schema_file = app_state.current_schema_file

    try:
        data = dataset_repository.read_schema(schema_file)

        # Find top-level key from registry or auto-detect
        reg = SCHEMA_REGISTRY.get(schema_file)
        if reg:
            top_key = reg['top_key']
            if top_key is None:
                return _strip_schema_notes(data)
        else:
            # Auto-detect: first key that looks like a taxonomy key
            top_key = None
            for k in data:
                if k.startswith('dfir_capability_taxonomy') or k.startswith('ai_trism_taxonomy') or k.startswith('preemptive_cybersecurity_taxonomy') or k.startswith('mdr_services_taxonomy') or k.startswith('offensive_security_taxonomy') or k.startswith('cnapp_taxonomy') or k.startswith('cnapp_mq_gap_taxonomy') or k.startswith('mq_gap_taxonomy'):
                    top_key = k
                    break
            if top_key is None:
                return _strip_schema_notes(data)  # return as-is

        return _strip_schema_notes(data.get(top_key, {}))
    except Exception as e:
        print(f"Error loading schema {schema_file}: {e}")
    return {}

def _strip_schema_notes(body):
    """Remove documentation-only keys (those starting with '_') from pillars/sub_pillars dicts.

    Templates use '_note' entries to document structure inline. These would break
    code that iterates pillars/sub_pillars expecting every value to be a dict.
    """
    if not isinstance(body, dict):
        return body
    for section in ('pillars', 'sub_pillars'):
        sect = body.get(section)
        if isinstance(sect, dict):
            body[section] = {k: v for k, v in sect.items()
                             if not (isinstance(k, str) and k.startswith('_'))}
    return body

def _schema_structure(schema_file=None):
    """Return 'nested', 'flat', or 'asmf' for the given schema."""
    if schema_file is None:
        schema_file = app_state.current_schema_file
    reg = SCHEMA_REGISTRY.get(schema_file)
    if reg:
        return reg['structure']
    # Guess from content
    schema = load_schema_data(schema_file)
    if 'dimensions' in schema and isinstance(schema['dimensions'], dict):
        return 'asmf'
    return 'flat' if 'sub_pillars' in schema else 'nested'

# Extract sub-pillar definitions (works for all schema versions)
def extract_sub_pillars(schema_file=None):
    """Extract all 20 sub-pillars with definitions from any schema version."""
    schema = load_schema_data(schema_file)
    structure = _schema_structure(schema_file)
    sub_pillars = []

    if structure == 'nested':
        # v3.2 style: pillars → sub_capabilities list
        if 'pillars' in schema:
            for pillar_key, pillar_data in schema['pillars'].items():
                pillar_code = pillar_data.get('code', '')
                pillar_name = pillar_data.get('focus', '')

                if 'sub_capabilities' in pillar_data:
                    for sub_cap in pillar_data['sub_capabilities']:
                        sub_pillars.append({
                            'id': sub_cap.get('id', ''),
                            'pillar_code': pillar_code,
                            'pillar_name': pillar_name,
                            'name': sub_cap.get('name', ''),
                            'definition': sub_cap.get('definition', ''),
                            'activities': sub_cap.get('granular_activities', [])
                        })
    elif structure == 'flat':
        # v4.0 / v5.0 style: pillars dict keyed by code + separate sub_pillars dict
        pillar_lookup = {}
        if 'pillars' in schema:
            for code, pdata in schema['pillars'].items():
                pillar_lookup[code] = pdata.get('name', pdata.get('focus', code))

        if 'sub_pillars' in schema:
            for sp_id, sp_data in schema['sub_pillars'].items():
                pillar_code = sp_id.split('-')[0] if '-' in sp_id else ''
                sub_pillars.append({
                    'id': sp_id,
                    'pillar_code': pillar_code,
                    'pillar_name': pillar_lookup.get(pillar_code, ''),
                    'name': sp_data.get('name', ''),
                    'definition': sp_data.get('expanded_definition', sp_data.get('definition', '')),
                    'activities': sp_data.get('what_to_verify_publicly',
                                              sp_data.get('gtm_evaluation_criteria',
                                                          sp_data.get('ai_evaluation_criteria',
                                                                      sp_data.get('maturity_criteria',
                                                                                  sp_data.get('granular_activities', [])))))
                })
    elif structure == 'asmf':
        if 'dimensions' in schema and isinstance(schema['dimensions'], dict):
            for dim_id, dim_data in schema['dimensions'].items():
                pillar_name = dim_data.get('name', '')
                for sd_id, sd_data in (dim_data.get('sub_dimensions') or {}).items():
                    sub_pillars.append({
                        'id': sd_id,
                        'pillar_code': dim_id,
                        'pillar_name': pillar_name,
                        'name': sd_data.get('name', ''),
                        'definition': sd_data.get('description', ''),
                        'activities': [sd_data.get('assessment_question')] if sd_data.get('assessment_question') else []
                    })

    return sub_pillars

# Load vendor data
def load_vendor_data(vendor_file=None):
    """Load vendor data from the specified JSON file (handles both wrapped and bare formats)"""
    if vendor_file is None:
        vendor_file = app_state.current_vendor_file
    
    vendors = []
    try:
        vendors.extend(dataset_repository.read_vendors(vendor_file))
    except (
        FileNotFoundError,
        json.JSONDecodeError,
        ValueError,
        UnicodeDecodeError,
    ) as e:
        print(f"Error loading {vendor_file}: {e}")
    
    return vendors

# Field metadata for descriptions and full names
FIELD_METADATA = {
    'vendor': {
        'name': 'Vendor Name',
        'description': 'The name of the vendor or service provider being evaluated'
    },
    'region': {
        'name': 'Region',
        'description': 'Geographic region where the vendor primarily operates (e.g., North America, Europe, APAC, Global)'
    },
    'specialization': {
        'name': 'Specialization',
        'description': 'The primary area of focus or expertise for the vendor'
    },
    'is_startup': {
        'name': 'Is Startup',
        'description': 'Whether the company is a startup (true/false)'
    },
    'is_ai_first': {
        'name': 'AI-First',
        'description': 'Whether the vendor is primarily AI-driven in their approach (true/false)'
    },
    'ir_focus_type': {
        'name': 'IR Focus Type',
        'description': 'Core Competency = IR is primary work product; Assistance Component = IR provided as feature/support'
    },
    'PLA': {
        'name': 'Planning (PLA)',
        'description': 'Organizational Readiness and Breach Preparation - capability score 1-5'
    },
    'INV': {
        'name': 'Investigation (INV)',
        'description': 'Evidence Identification, Collection, and Analytical Reconstruction - capability score 1-5'
    },
    'REM': {
        'name': 'Remediation (REM)',
        'description': 'Threat Containment and Business Restoration - capability score 1-5'
    },
    'PMG': {
        'name': 'Program Management (PMG)',
        'description': 'Incident Lifecycle Oversight and Communication - capability score 1-5'
    },
    'LAW': {
        'name': 'Legal (LAW)',
        'description': 'Legal Admissibility and Judicial Support - capability score 1-5'
    },
    'capability_analysis': {
        'name': 'Capability Analysis',
        'description': 'Detailed explanation of the vendor\'s key capabilities and differentiators'
    }
}

# Scoring legend
SCORE_LEGEND = {
    '1': 'Manual - Human-led with no technological automation',
    '2': 'Insufficient Evidence - Service provided but AI integration not verified',
    '3': 'AI-Augmented - Basic generative AI assistants for summarization/drafting',
    '4': 'Advanced AI - Specialized models perform deep correlation with human validation',
    '5': 'Fully Agentic - Autonomous systems that independently plan and execute'
}

@app.route('/')
def index():
    static_dir = os.path.join(os.path.dirname(__file__), 'static')
    def mtime_or_now(name: str) -> int:
        try:
            return int(os.path.getmtime(os.path.join(static_dir, name)))
        except Exception:
            return int(time.time())

    static_versions = {
        'app_js': mtime_or_now('app.js'),
        'style_css': mtime_or_now('style.css'),
    }

    return render_template(
        'index.html',
        field_metadata=FIELD_METADATA,
        score_legend=SCORE_LEGEND,
        static_versions=static_versions,
    )

@app.route('/api/vendors')
def get_vendors():
    """Get all vendors with optional filtering"""
    vendors = load_vendor_data()
    
    # Get filter parameters
    search_query = request.args.get('search', '').lower()
    filters = {}
    
    # Parse filter parameters
    for key in request.args:
        if key.startswith('filter_'):
            field_name = key.replace('filter_', '')
            value = request.args.get(key)
            if value:
                filters[field_name] = value.lower()
    
    # Apply filters
    filtered_vendors = []
    for vendor in vendors:
        match = True
        
        # Apply field-specific filters
        for field, filter_value in filters.items():
            if field in vendor:
                vendor_value = str(vendor[field]).lower()
                if filter_value not in vendor_value:
                    match = False
                    break
        
        # Apply search across all fields
        if search_query and match:
            found = False
            for key, value in vendor.items():
                if search_query in str(value).lower():
                    found = True
                    break
            match = found
        
        if match:
            filtered_vendors.append(vendor)
    
    return jsonify(filtered_vendors)

@app.route('/api/field-values/<field>')
def get_field_values(field):
    """Get unique values for a specific field"""
    vendors = load_vendor_data()
    values = set()
    
    for vendor in vendors:
        if field in vendor:
            val = vendor[field]
            if isinstance(val, bool):
                val = str(val)
            elif isinstance(val, dict):
                continue
            elif isinstance(val, list):
                continue
            values.add(str(val))
    
    return jsonify(sorted(list(values)))

@app.route('/api/filter-options')
def get_filter_options():
    """Return values for several vendor filter fields from one dataset read."""
    requested_fields = [
        field.strip()
        for field in request.args.get('fields', '').split(',')
        if field.strip()
    ]
    if not requested_fields:
        return jsonify({})

    vendors = load_vendor_data()
    options = {field: set() for field in requested_fields}
    for vendor in vendors:
        for field in requested_fields:
            if field not in vendor:
                continue
            value = vendor[field]
            if isinstance(value, (dict, list)):
                continue
            options[field].add(str(value))

    return jsonify({field: sorted(values) for field, values in options.items()})

@app.route('/api/metadata')
def get_metadata():
    """Get field metadata and scoring legend – schema-aware."""
    schema_file = app_state.current_schema_file
    schema = load_schema_data(schema_file)
    sub_pillars = extract_sub_pillars(schema_file)

    # ── Score legend from schema (fall back to hardcoded) ──
    meta = schema.get('metadata', {})
    scoring_logic = (
        meta.get('scoring_scale', {}).get('scoring_logic', {})
        or meta.get('capability_scoring_scale', {}).get('scoring_logic', {})
        or meta.get('scoring_logic', {})
    )
    if scoring_logic:
        score_legend = {str(k): v for k, v in scoring_logic.items()}
    elif 'dual_scoring_methodology' in meta:
        # PMR-style dual scoring: build legend from GTM + Proof scales
        dsm = meta['dual_scoring_methodology']
        gtm_scale = dsm.get('gtm_messaging_score', {}).get('scale', {})
        proof_scale = dsm.get('proof_of_execution_score', {}).get('scale', {})
        score_legend = {}
        for k, v in gtm_scale.items():
            proof_v = proof_scale.get(k, '')
            # Combine both scales: "GTM: X | Proof: Y"
            score_legend[str(k)] = f'GTM: {v} | Proof: {proof_v}'
    elif _schema_structure(schema_file) == 'asmf':
        score_legend = {}
    else:
        score_legend = dict(SCORE_LEGEND)

    # ── Pricing evaluation (MDR schema has a separate pricing scale) ──
    pricing_evaluation = schema.get('pricing_evaluation', None)
    pricing_score_legend = None
    if pricing_evaluation:
        pricing_scoring = meta.get('pricing_scoring_scale', {}).get('scoring_logic', {})
        if pricing_scoring:
            pricing_score_legend = {str(k): v for k, v in pricing_scoring.items()}

    # ── Pillar descriptions from schema ──
    # Start with base fields but strip any hardcoded pillar codes;
    # only the pillars defined in the active schema should appear.
    all_known_pillar_codes = {'PLA', 'INV', 'REM', 'PMG', 'LAW', 'GOV', 'RUN', 'INF',
                               'PPD', 'PCS', 'TDT', 'PCM', 'CTL'}
    field_metadata = {k: v for k, v in FIELD_METADATA.items() if k not in all_known_pillar_codes}
    pillars_in_schema = {}
    if _schema_structure(schema_file) == 'asmf' and 'dimensions' in schema:
        for code, dim in schema['dimensions'].items():
            pillars_in_schema[code] = {
                'name': dim.get('name', code),
                'focus': dim.get('description', ''),
                'ai_evidence_signals': dim.get('evidence_signals', []),
                'validated_pillar_score_rule': dim.get('validated_pillar_score_rule', ''),
            }
    else:
        pillars_in_schema = schema.get('pillars', {})

    for code, pdata in pillars_in_schema.items():
        field_metadata[code] = {
            'name': pdata.get('name', code),
            'description': pdata.get('focus', ''),
        }

    # ── Group sub-pillars by pillar ──
    pillars_grouped = {}
    for sub_pillar in sub_pillars:
        pillar_code = sub_pillar['pillar_code']
        if pillar_code not in pillars_grouped:
            pdata = pillars_in_schema.get(pillar_code, {})
            pillars_grouped[pillar_code] = {
                'code': pillar_code,
                'name': pdata.get('name', sub_pillar.get('pillar_name', '')),
                'description': pdata.get('focus', field_metadata.get(pillar_code, {}).get('description', '')),
                'ai_evidence_signals': pdata.get('ai_evidence_signals', pdata.get('evidence_signals', pdata.get('maturity_signals', []))),
                'validated_pillar_score_rule': pdata.get('validated_pillar_score_rule', ''),
                'sub_pillars': []
            }
        pillars_grouped[pillar_code]['sub_pillars'].append({
            'id': sub_pillar['id'],
            'name': sub_pillar['name'],
            'definition': sub_pillar.get('definition', ''),
            'activities': sub_pillar.get('activities', []),
        })

    # ── Schema intent ──
    intent = schema.get('intent', '')

    response = {
        'field_metadata': field_metadata,
        'score_legend': score_legend,
        'pillars_grouped': list(pillars_grouped.values()),
        'schema_file': schema_file,
        'schema_intent': intent,
    }
    if pricing_evaluation:
        response['pricing_evaluation'] = {
            'description': pricing_evaluation.get('description', ''),
            'framework_source': pricing_evaluation.get('framework_source', ''),
            'dimensions': pricing_evaluation.get('dimensions', {}),
        }
        if pricing_score_legend:
            response['pricing_score_legend'] = pricing_score_legend
    return jsonify(response)

@app.route('/api/sub-pillars')
def get_sub_pillars():
    """Get all 20 sub-pillars with definitions and filtering options"""
    sub_pillars = extract_sub_pillars()
    
    # Group by pillar for easier frontend consumption
    pillars_grouped = {}
    for sub_pillar in sub_pillars:
        pillar_code = sub_pillar['pillar_code']
        if pillar_code not in pillars_grouped:
            pillars_grouped[pillar_code] = {
                'code': pillar_code,
                'name': sub_pillar['pillar_name'],
                'sub_pillars': []
            }
        pillars_grouped[pillar_code]['sub_pillars'].append(sub_pillar)
    
    return jsonify({
        'pillars': list(pillars_grouped.values()),
        'all_sub_pillars': sub_pillars,
        'total_count': len(sub_pillars)
    })

@app.route('/api/vendors/by-sub-pillar')
def get_vendors_by_sub_pillar():
    """Filter vendors by sub-pillar score"""
    sub_pillar_id = request.args.get('sub_pillar', '').upper()
    min_score = request.args.get('min_score', '1')
    
    try:
        min_score = int(min_score)
    except ValueError:
        min_score = 1
    
    vendors = load_vendor_data()
    filtered = []
    
    for vendor in vendors:
        if 'granular_mapping' in vendor:
            granular = vendor['granular_mapping']
            # Extract pillar code from sub_pillar_id (e.g., "PLA" from "PLA-01")
            pillar_code = sub_pillar_id.split('-')[0] if '-' in sub_pillar_id else sub_pillar_id
            
            if pillar_code in granular and sub_pillar_id in granular[pillar_code]:
                score = granular[pillar_code][sub_pillar_id]
                if isinstance(score, (int, float)) and score >= min_score:
                    filtered.append({
                        'vendor': vendor,
                        'sub_pillar_score': score
                    })
    
    return jsonify(filtered)

@app.route('/api/update-definition', methods=['POST'])
def update_definition():
    """Update pillar, sub-pillar, score, or field definitions"""
    data = request.json
    edit_type = data.get('type')
    edit_id = data.get('id')
    name = data.get('name')
    description = data.get('description')
    
    try:
        if edit_type == 'pillar':
            # Update pillar in FIELD_METADATA
            FIELD_METADATA[edit_id]['name'] = name
            FIELD_METADATA[edit_id]['description'] = description
        elif edit_type == 'score':
            # Update score in SCORE_LEGEND
            SCORE_LEGEND[edit_id] = description
        elif edit_type == 'field':
            # Update field metadata
            FIELD_METADATA[edit_id]['name'] = name
            FIELD_METADATA[edit_id]['description'] = description
        elif edit_type == 'sub-pillar':
            # Update sub-pillar definition in the currently active schema
            schema_filename = app_state.current_schema_file
            full_schema = read_dataset(schema_filename)
            
            # Find the top-level key
            reg = SCHEMA_REGISTRY.get(schema_filename)
            top_key = reg['top_key'] if reg else list(full_schema.keys())[0]
            schema = full_schema.get(top_key, full_schema)
            structure = _schema_structure(schema_filename)

            if structure == 'nested':
                if 'pillars' in schema:
                    for pillar_data in schema['pillars'].values():
                        if 'sub_capabilities' in pillar_data:
                            for sub_cap in pillar_data['sub_capabilities']:
                                if sub_cap.get('id') == edit_id:
                                    sub_cap['name'] = name
                                    sub_cap['definition'] = description
            else:
                # flat structure (v4.0 / v5.0)
                if 'sub_pillars' in schema and edit_id in schema['sub_pillars']:
                    schema['sub_pillars'][edit_id]['name'] = name
                    if 'expanded_definition' in schema['sub_pillars'][edit_id]:
                        schema['sub_pillars'][edit_id]['expanded_definition'] = description
                    else:
                        schema['sub_pillars'][edit_id]['definition'] = description
            
            write_error = persist_dataset(schema_filename, full_schema)
            if write_error is not None:
                return write_error
        
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error updating definition: {e}")
        return jsonify({'success': False, 'error': str(e)}), 400

@app.route('/api/vendor-files', methods=['GET'])
def get_vendor_files():
    """Get list of available vendor JSON files, isolated by project.
    
    Uses project-based isolation so only vendor files belonging to the
    active schema's project are shown:
      - 'trism'      → AI TRiSM vendor files
      - 'preemptive'  → Preemptive Cybersecurity vendor files
      - everything else → DFIR vendor files
    Pass ?schema=<filename> to filter, or ?all=1 for unfiltered.
    """
    def _detect_project(name):
        """Return project tag from a filename (schema or vendor)."""
        lower = name.lower()
        if 'schema_template' in lower:
            return 'template'
        if ('agentic_soc' in lower or 'agentic_enterprise' in lower or
            'ai_platform_ecosystem' in lower or 'platform_ecosystem' in lower or
            'asmf' in lower):
            return 'asmf'
        if 'trism' in lower:
            return 'trism'
        if 'preemptive' in lower or 'precyber' in lower:
            return 'precyber'
        if 'secure_by_design' in lower or 'sbd_ai' in lower or 'sbdai' in lower:
            return 'sbdai'
        if 'cnapp_mq' in lower or 'cnapp mq' in lower:
            return 'cnapp_mq'
        if 'mq_gap' in lower or 'mq gap' in lower:
            return 'mq_gap'
        if 'cnapp' in lower:
            return 'cnapp'
        if 'mdr' in lower:
            return 'mdr'
        if 'offensive' in lower:
            return 'offsec'
        if 'product market readiness' in lower or 'pmr' in lower:
            return 'pmr'
        # DFIR schemas and vendor files (Vendor 3-x, 4-x, 5-x, 6-x, schema3-3, etc.)
        return 'dfir'

    try:
        schema_filter = request.args.get('schema', '')
        show_all = request.args.get('all', '0') == '1'
        # Determine active project from schema param or current state
        active_schema = schema_filter or app_state.current_schema_file
        active_project = _detect_project(active_schema)
        available_files = []
        app_dir = os.path.dirname(__file__)

        # Files to exclude: scoreless files, old superseded files, partial delta files
        _EXCLUDED_VENDOR_FILES = {
            'vendor 5-1 incomplete research summary.json',   # no scores
            'vendor 5-1 incomplete research delta.json',     # partial (43/138)
            'vendor3-3.json',                                # old, superseded by Vendor 3-7
            'vendor3-4.json',                                # old, superseded by Vendor 3-7
            'vendor3-5.json',                                # old, superseded by Vendor 3-7
        }
        
        # Find all vendor JSON files dynamically
        for filename in canonical_paths('vendor_score'):
            # Match vendor files: vendor*.json, Vendor*.json, or *researched*.json
            filename_lower = filename.lower()
            if (filename.endswith('.json') and 
                ('vendor' in filename_lower or 'researched' in filename_lower or 'validated' in filename_lower)):

                # Skip excluded files (scoreless, old, partial)
                if filename_lower in _EXCLUDED_VENDOR_FILES:
                    continue

                # ── Project isolation: only show files matching the active project ──
                if not show_all:
                    file_project = _detect_project(filename)
                    if active_project != file_project:
                        continue  # skip files from other projects

                try:
                    data = read_dataset(filename)
                    # Extract schema_ref and count
                    file_schema_ref = ''
                    count = 0
                    if isinstance(data, dict) and 'vendors' in data and isinstance(data['vendors'], list):
                        # Wrapped format
                        file_schema_ref = data.get('schema_ref', '')
                        count = len(data['vendors'])
                    elif isinstance(data, list):
                        count = len(data)
                        # Check first vendor for schema_ref
                        if data and isinstance(data[0], dict):
                            file_schema_ref = data[0].get('schema_ref', '')
                    elif isinstance(data, dict):
                        file_schema_ref = data.get('schema_ref', '')
                        for key, value in data.items():
                            if isinstance(value, list):
                                count = len(value)
                                break

                    # Create display name from filename
                    display_name = filename.replace('.json', '')
                    available_files.append({
                        'filename': filename,
                        'name': display_name,
                        'count': count,
                        'schema_ref': file_schema_ref
                    })
                except Exception as err:
                    print(f"Error reading {filename}: {err}")
                    pass
        
        # Sort by filename for consistent ordering
        available_files.sort(key=lambda x: x['filename'].lower())
        
        return jsonify({
            'files': available_files,
            'current': app_state.current_vendor_file,
            'current_schema': app_state.current_schema_file
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 400

@app.route('/api/switch-vendor-file', methods=['POST'])
def switch_vendor_file():
    """Switch to a different vendor file and auto-switch schema to match."""
    try:
        data = request.get_json(force=True, silent=True) or {}
        new_file = data.get('filename')
        schema_hint = data.get('schema')
        
        # Only source-controlled active datasets may be selected.
        if new_file not in canonical_paths('vendor_score'):
            return jsonify({'success': False, 'error': 'File not found'}), 404
        
        # Verify it's a valid JSON file
        try:
            read_dataset(new_file)
        except Exception:
            return jsonify({'success': False, 'error': 'Invalid JSON file'}), 400
        
        # Switch the file
        app_state.current_vendor_file = new_file

        # Keep explicitly-selected framework schema pinned for framework flows.
        # Frontend sends the active schema while switching vendor files.
        schema_switched = False
        if schema_hint in SCHEMA_REGISTRY and SCHEMA_REGISTRY.get(schema_hint, {}).get('structure') == 'asmf':
            if app_state.current_schema_file != schema_hint:
                app_state.current_schema_file = schema_hint
                schema_switched = True
        else:
            # Auto-switch schema to match the vendor file's project
            best_schema = _detect_best_schema(new_file)
            if best_schema and best_schema != app_state.current_schema_file:
                app_state.current_schema_file = best_schema
                schema_switched = True
        
        return jsonify({
            'success': True,
            'current': app_state.current_vendor_file,
            'current_schema': app_state.current_schema_file,
            'schema_switched': schema_switched,
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 400


# Preferred schema per project (latest/best version)
_PROJECT_SCHEMA_MAP = {
    'precyber_v3': 'Preemptive_Cybersecurity_Schema_v3.json',
    'precyber': 'Preemptive_Cybersecurity_Schema_v2.json',
    'trism': 'AI TriSM Schema 1_1.json',
    'dfir': 'schema3-3.json',
    'sbdai': 'Secure_by_Design_AI_Controls_Schema_v2.json',
    'mdr': 'MDR_Services_Schema.json',
    'cnapp': 'CNAPP_Schema.json',
    'cnapp_mq': 'CNAPP_MQ_Gap_Schema_App.json',
    'offsec': 'Offensive_Security_Schema.json',
    'pmr': 'Product Market Readiness Schema 1_0.json',
    'mq_gap': 'MDR_MQ_Gap_Schema_App.json',
}

def _detect_project_from_name(name):
    """Return project tag from a filename."""
    lower = name.lower()
    if 'schema_template' in lower:
        return 'template'
    if (
        'agentic_soc' in lower or
        'asmf' in lower or
        'agentic_enterprise' in lower or
        'ai_platform_ecosystem' in lower or
        'platform_ecosystem' in lower or
        'ai platform ecosystem' in lower
    ):
        return 'asmf'
    if 'trism' in lower:
        return 'trism'
    if ('preemptive' in lower or 'precyber' in lower) and ('6-0' in lower or 'v3' in lower):
        return 'precyber_v3'
    if 'preemptive' in lower or 'precyber' in lower:
        return 'precyber'
    if 'secure_by_design' in lower or 'sbd_ai' in lower or 'sbdai' in lower: return 'sbdai'
    if 'cnapp_mq' in lower or 'cnapp mq' in lower: return 'cnapp_mq'
    if 'mq_gap' in lower or 'mq gap' in lower: return 'mq_gap'
    if 'cnapp' in lower: return 'cnapp'
    if 'mdr' in lower: return 'mdr'
    if 'offensive' in lower: return 'offsec'
    if 'product market readiness' in lower or 'pmr' in lower: return 'pmr'
    return 'dfir'

def _detect_best_schema(vendor_file):
    """Given a vendor filename, return the best matching schema filename."""
    project = _detect_project_from_name(vendor_file)
    schema = _PROJECT_SCHEMA_MAP.get(project)
    if schema and schema in canonical_paths('schema'):
        return schema
    return None

# ── Schema API endpoints ─────────────────────────────────────────────

@app.route('/api/schema-files', methods=['GET'])
def get_schema_files():
    """List all available schema files with metadata."""
    schemas = []
    for fn in discover_schema_files():
        try:
            body = load_schema_data(fn)
            meta = body.get('metadata', {})
            display = SCHEMA_DISPLAY.get(fn, {'title': fn.replace('.json', ''), 'abbr': '', 'subtitle': ''})
            schemas.append({
                'filename': fn,
                'name': fn.replace('.json', ''),
                'intent': body.get('intent', ''),
                'scoring_logic': meta.get('scoring_scale', meta.get('scoring_logic', {})),
                'display': display,
                'kind': 'framework' if SCHEMA_REGISTRY.get(fn, {}).get('structure') == 'asmf' else 'schema',
                'capabilities': _framework_capabilities(fn) if SCHEMA_REGISTRY.get(fn, {}).get('structure') == 'asmf' else None,
            })
        except Exception as e:
            print(f"Error reading schema {fn}: {e}")
    return jsonify({
        'schemas': schemas,
        'current': app_state.current_schema_file
    })

@app.route('/api/switch-schema', methods=['POST'])
def switch_schema():
    """Switch the active schema and return matching vendor files."""
    data = request.get_json(force=True, silent=True) or {}
    new_schema = data.get('filename', '')
    if new_schema not in canonical_paths('schema', 'framework'):
        return jsonify({'success': False, 'error': 'Schema file not found'}), 404
    try:
        read_dataset(new_schema)
    except Exception:
        return jsonify({'success': False, 'error': 'Schema file not found'}), 404

    app_state.current_schema_file = new_schema
    if SCHEMA_REGISTRY.get(new_schema, {}).get('structure') == 'asmf':
        if new_schema == 'AI_platform_ecosystem_framework_v1.json':
            app_state.current_vendor_file = 'ai_platform_ecosystem_vendors_v1.json'
        else:
            app_state.current_vendor_file = ''
    return jsonify({
        'success': True,
        'current_schema': app_state.current_schema_file
    })

@app.route('/api/schema-detail', methods=['GET'])
def get_schema_detail():
    """Return the full parsed schema body for the currently-selected schema."""
    schema_file = request.args.get('schema', app_state.current_schema_file)
    body = load_schema_data(schema_file)
    sub_pillars = extract_sub_pillars(schema_file)
    structure = _schema_structure(schema_file)

    # Build scoring_logic in a uniform way
    meta = body.get('metadata', {})
    scoring = (
        meta.get('scoring_scale', {}).get('scoring_logic', {})
        or meta.get('capability_scoring_scale', {}).get('scoring_logic', {})
        or meta.get('scoring_logic', {})
    )

    # Pillars summary with full detail
    pillars_list = []
    if structure == 'asmf' and 'dimensions' in body:
        for code, dim in body['dimensions'].items():
            entry = {
                'code': dim.get('id', code),
                'name': dim.get('name', code),
                'focus': dim.get('description', ''),
                'ai_evidence_signals': dim.get('evidence_signals', []),
                'validated_pillar_score_rule': dim.get('validated_pillar_score_rule', ''),
            }
            for passthrough_key in ('aiuc1_requirements', 'aiuc1_coverage', 'aiuc1_categories',
                                    'ai_rmf_functions', 'nist_references', 'maturity_signals'):
                if passthrough_key in dim:
                    entry[passthrough_key] = dim[passthrough_key]
            pillars_list.append(entry)
    elif 'pillars' in body:
        for code, pdata in body['pillars'].items():
            entry = {
                'code': pdata.get('code', code),
                'name': pdata.get('name', pdata.get('focus', code)),
                'focus': pdata.get('focus', ''),
                'ai_evidence_signals': pdata.get('ai_evidence_signals', []),
                'validated_pillar_score_rule': pdata.get('validated_pillar_score_rule', ''),
            }
            # Pass through AIUC-1 and other enrichment fields
            for passthrough_key in ('aiuc1_requirements', 'aiuc1_coverage', 'aiuc1_categories',
                                    'ai_rmf_functions', 'nist_references', 'maturity_signals'):
                if passthrough_key in pdata:
                    entry[passthrough_key] = pdata[passthrough_key]
            pillars_list.append(entry)

    # Enrich sub-pillars with full schema detail
    raw_sub_pillars = body.get('sub_pillars', {})
    for sp in sub_pillars:
        sp_id = sp.get('id', '')
        if sp_id in raw_sub_pillars:
            raw = raw_sub_pillars[sp_id]
            sp['ai_evaluation_criteria'] = raw.get('ai_evaluation_criteria', [])
            sp['what_to_verify_publicly'] = raw.get('what_to_verify_publicly', [])
            sp['ai_specific_evidence'] = raw.get('ai_specific_evidence', [])
            # Pass through AIUC-1 and maturity enrichment fields
            for passthrough_key in ('aiuc1_mapping', 'ai_rmf_mapping', 'maturity_criteria'):
                if passthrough_key in raw:
                    sp[passthrough_key] = raw[passthrough_key]

    return dataset_json_response({
        'schema_file': schema_file,
        'intent': body.get('intent', ''),
        'scoring_logic': scoring,
        'pillars': pillars_list,
        'sub_pillars': sub_pillars,
        'sub_pillar_count': len(sub_pillars),
    }, schema_file)


@app.route('/api/export-schema-html', methods=['GET'])
def export_schema_html():
    """Render the currently-selected (or ?schema=) schema as a styled standalone HTML document.

    Returns an HTML attachment suitable for sharing or archiving.
    """
    from flask import Response
    from html import escape as _esc

    schema_file = request.args.get('schema', app_state.current_schema_file)
    body = load_schema_data(schema_file) or {}
    display = SCHEMA_DISPLAY.get(schema_file, {
        'title': schema_file.replace('.json', ''), 'abbr': '', 'subtitle': ''
    })
    meta = body.get('metadata', {}) or {}
    lineage = body.get('schema_lineage', {}) or {}
    intent = body.get('intent', '')

    # Scoring scale (handle multiple shapes)
    scoring_block = (
        meta.get('scoring_scale')
        or meta.get('capability_scoring_scale')
        or meta.get('scoring_logic')
        or {}
    )
    scoring_logic = scoring_block.get('scoring_logic', scoring_block) if isinstance(scoring_block, dict) else {}

    pillars = body.get('pillars', {}) or {}
    sub_pillars = body.get('sub_pillars', {}) or {}

    # Group sub-pillars by pillar code prefix (e.g., "VIA-01" → VIA)
    grouped_subs = {}
    for sp_id, sp_data in sub_pillars.items():
        code = sp_id.split('-', 1)[0] if '-' in sp_id else sp_id
        grouped_subs.setdefault(code, []).append((sp_id, sp_data or {}))

    def _list_html(items):
        if not items:
            return ''
        if isinstance(items, dict):
            items = [f"{k}: {v}" for k, v in items.items()]
        if not isinstance(items, list):
            items = [str(items)]
        return '<ul>' + ''.join(f'<li>{_esc(str(x))}</li>' for x in items) + '</ul>'

    def _kv_block(d):
        if not isinstance(d, dict):
            return ''
        rows = ''.join(
            f'<tr><th>{_esc(str(k))}</th><td>{_esc(str(v))}</td></tr>'
            for k, v in d.items() if not isinstance(v, (dict, list))
        )
        return f'<table class="kv">{rows}</table>' if rows else ''

    # Build pillars + sub-pillars HTML
    pillars_html_parts = []
    for code, pdata in pillars.items():
        pdata = pdata or {}
        name = pdata.get('name', code)
        focus = pdata.get('focus', '')
        evidence = pdata.get('evidence_signals') or pdata.get('ai_evidence_signals') or []
        # Pillar header card
        meta_rows = []
        for k in ('mq_criterion', 'mq_criterion_name', 'mq_weight', 'mq_axis'):
            if k in pdata:
                meta_rows.append(f'<tr><th>{_esc(k.replace("_", " ").title())}</th><td>{_esc(str(pdata[k]))}</td></tr>')
        meta_table = f'<table class="kv">{"".join(meta_rows)}</table>' if meta_rows else ''

        sp_cards = []
        for sp_id, sp in grouped_subs.get(code, []):
            sg = sp.get('scoring_guidance', {}) or {}
            sg_rows = ''.join(
                f'<tr><th>{_esc(str(k))}</th><td>{_esc(str(v))}</td></tr>'
                for k, v in sorted(sg.items(), key=lambda x: str(x[0]))
            )
            sp_cards.append(f"""
              <div class="sub-pillar">
                <div class="sp-id">{_esc(sp_id)}</div>
                <div class="sp-name">{_esc(sp.get('name', ''))}</div>
                <p class="sp-def">{_esc(sp.get('expanded_definition', ''))}</p>
                <div class="sp-section"><strong>What to verify publicly</strong>{_list_html(sp.get('what_to_verify_publicly', []))}</div>
                <div class="sp-section"><strong>Search terms</strong>{_list_html(sp.get('search_terms', []))}</div>
                {f'<div class="sp-section"><strong>Scoring guidance</strong><table class="scoring">{sg_rows}</table></div>' if sg_rows else ''}
              </div>
            """)

        pillars_html_parts.append(f"""
          <section class="pillar">
            <h2><span class="pillar-code">{_esc(code)}</span> {_esc(name)}</h2>
            {f'<p class="pillar-focus">{_esc(focus)}</p>' if focus else ''}
            {meta_table}
            {f'<div class="sp-section"><strong>Evidence signals</strong>{_list_html(evidence)}</div>' if evidence else ''}
            <div class="sub-pillars-grid">{''.join(sp_cards)}</div>
          </section>
        """)

    # MQ scoring mode (if present)
    mq_mode = body.get('mq_scoring_mode', {}) or {}
    mq_html = ''
    if mq_mode:
        axes_html = []
        for axis_key, axis in (mq_mode.get('axes', {}) or {}).items():
            criteria_rows = ''.join(
                f"""<tr>
                  <td>{_esc(c.get('criterion_id', ''))}</td>
                  <td>{_esc(c.get('criterion_name', ''))}</td>
                  <td>{_esc(c.get('weight', ''))}</td>
                  <td>{_esc(str(c.get('weight_value', '')))}</td>
                  <td>{_esc(c.get('formula', ''))}</td>
                  <td>{_esc(c.get('source', ''))}</td>
                </tr>"""
                for c in (axis.get('criteria', []) or [])
            )
            axes_html.append(f"""
              <h3>{_esc(axis_key.replace('_', ' ').title())}</h3>
              <table class="mq-table">
                <thead><tr><th>ID</th><th>Criterion</th><th>Weight</th><th>Value</th><th>Formula</th><th>Source</th></tr></thead>
                <tbody>{criteria_rows}</tbody>
              </table>
            """)
        mq_html = f"""
          <section class="mq-mode">
            <h2>Magic Quadrant Scoring Mode</h2>
            <p>{_esc(mq_mode.get('description', ''))}</p>
            {''.join(axes_html)}
          </section>
        """

    # Scoring logic table
    scoring_rows = ''.join(
        f'<tr><th>{_esc(str(k))}</th><td>{_esc(str(v))}</td></tr>'
        for k, v in sorted((scoring_logic or {}).items(), key=lambda x: str(x[0]))
    )

    title = display.get('title', schema_file)
    subtitle = display.get('subtitle', '')

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>{_esc(title)}</title>
<style>
  body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; max-width: 1100px; margin: 32px auto; padding: 0 24px; color: #1a1a1a; line-height: 1.55; }}
  header.doc-header {{ border-bottom: 3px solid #0066cc; padding-bottom: 16px; margin-bottom: 24px; }}
  header.doc-header .abbr {{ display: inline-block; background: #0066cc; color: #fff; padding: 4px 10px; border-radius: 4px; font-size: 12px; font-weight: 700; letter-spacing: 0.5px; }}
  header.doc-header h1 {{ margin: 8px 0 4px; font-size: 28px; }}
  header.doc-header p.subtitle {{ color: #555; margin: 0 0 8px; font-size: 15px; }}
  header.doc-header p.meta {{ color: #888; font-size: 12px; margin: 0; }}
  section {{ margin: 32px 0; }}
  h2 {{ font-size: 22px; border-bottom: 1px solid #ddd; padding-bottom: 6px; margin-top: 40px; }}
  h3 {{ font-size: 17px; margin-top: 24px; color: #333; }}
  .pillar-code {{ display: inline-block; background: #eef4ff; color: #0066cc; padding: 2px 8px; border-radius: 3px; font-family: ui-monospace, Menlo, Consolas, monospace; font-size: 14px; margin-right: 8px; }}
  .pillar-focus {{ color: #444; font-style: italic; margin: 6px 0 12px; }}
  table {{ border-collapse: collapse; width: 100%; margin: 8px 0 16px; font-size: 14px; }}
  table.kv th {{ text-align: left; background: #fafafa; padding: 6px 10px; width: 200px; border: 1px solid #eaeaea; vertical-align: top; }}
  table.kv td {{ padding: 6px 10px; border: 1px solid #eaeaea; }}
  table.scoring th {{ text-align: center; background: #f3f7ff; padding: 6px; border: 1px solid #d8e6f7; width: 50px; font-family: ui-monospace, Menlo, Consolas, monospace; }}
  table.scoring td {{ padding: 6px 10px; border: 1px solid #d8e6f7; font-size: 13px; }}
  table.mq-table th {{ background: #0066cc; color: #fff; padding: 8px 10px; text-align: left; font-size: 13px; }}
  table.mq-table td {{ padding: 8px 10px; border-bottom: 1px solid #eee; font-size: 13px; vertical-align: top; }}
  ul {{ margin: 4px 0 8px 22px; padding: 0; }}
  li {{ margin: 2px 0; font-size: 13px; }}
  .sub-pillars-grid {{ display: grid; grid-template-columns: 1fr; gap: 16px; margin-top: 16px; }}
  .sub-pillar {{ border: 1px solid #e6e6e6; border-radius: 6px; padding: 14px 18px; background: #fcfcfd; }}
  .sub-pillar .sp-id {{ font-family: ui-monospace, Menlo, Consolas, monospace; color: #0066cc; font-size: 12px; font-weight: 700; letter-spacing: 0.5px; }}
  .sub-pillar .sp-name {{ font-size: 16px; font-weight: 600; margin: 4px 0 8px; }}
  .sub-pillar .sp-def {{ font-size: 13px; color: #444; margin: 0 0 10px; }}
  .sp-section {{ margin: 8px 0; font-size: 13px; }}
  .sp-section strong {{ display: block; color: #333; margin-bottom: 4px; font-size: 12px; text-transform: uppercase; letter-spacing: 0.5px; }}
  footer {{ margin-top: 48px; padding-top: 16px; border-top: 1px solid #ddd; font-size: 12px; color: #888; text-align: center; }}
  @media print {{ body {{ max-width: 100%; }} }}
</style>
</head>
<body>
<header class="doc-header">
  <span class="abbr">{_esc(display.get('abbr', ''))}</span>
  <h1>{_esc(title)}</h1>
  <p class="subtitle">{_esc(subtitle)}</p>
  <p class="meta">Source: {_esc(schema_file)} &bull; Schema ID: {_esc(lineage.get('schema_id', ''))} &bull; Version: {_esc(str(lineage.get('version', '')))} &bull; Exported: {_esc(datetime.datetime.now().strftime('%Y-%m-%d %H:%M'))}</p>
</header>

<section>
  <h2>Intent</h2>
  <p>{_esc(intent)}</p>
</section>

<section>
  <h2>Metadata</h2>
  {_kv_block({k: v for k, v in meta.items() if not isinstance(v, (dict, list))})}
  {f'<h3>Scoring Scale</h3><table class="scoring">{scoring_rows}</table>' if scoring_rows else ''}
</section>

<section>
  <h2>Pillars &amp; Sub-Pillars ({len(pillars)} pillars, {len(sub_pillars)} sub-pillars)</h2>
  {''.join(pillars_html_parts)}
</section>

{mq_html}

<footer>
  Generated by the Gartner Vendor Analysis Platform &bull; {_esc(schema_file)}
</footer>
</body>
</html>
"""

    safe_name = schema_file.replace('.json', '').replace(' ', '_')
    return Response(
        html,
        mimetype='text/html; charset=utf-8',
        headers={'Content-Disposition': f'attachment; filename="{safe_name}.html"'}
    )


@app.route('/api/report/<report_id>', methods=['GET'])
def get_report(report_id):
    """Return a Market Insight report JSON by ID."""
    for source_path in canonical_paths('report_definition'):
        if source_path.startswith('Reports/') and report_id in source_path:
            try:
                data = read_dataset(source_path)
                return jsonify(data)
            except Exception as e:
                return jsonify({'error': str(e)}), 500
    return jsonify({'error': f'Report {report_id} not found'}), 404


@app.route('/api/reports', methods=['GET'])
def list_reports():
    """List all available report JSON files."""
    reports = []
    for source_path in sorted(canonical_paths('report_definition')):
        if not source_path.startswith('Reports/'):
            continue
        fname = source_path.split('/', 1)[1]
        try:
            data = read_dataset(source_path)
            reports.append({
                'file': fname,
                'id': data.get('id', fname),
                'title': data.get('title', fname),
                'content_type': data.get('content_type', 'Unknown'),
                'last_updated': data.get('last_updated', '')
            })
        except Exception:
            pass
    return jsonify(reports)


@app.route('/api/adoption-plan', methods=['GET'])
def get_adoption_plan():
    """Return the 90-day adoption plan from the active schema, enriched with sub-pillar names."""
    schema = load_schema_data()
    plan = schema.get('90_day_adoption_plan', {})
    if not plan:
        return jsonify({'error': 'No adoption plan in active schema'}), 404

    # Build sub-pillar lookup (name + pillar code)
    sp_lookup = {}   # sp_id → {name, pillar_code}
    pillar_lookup = {}  # pillar_code → pillar_name
    if 'sub_pillars' in schema:
        for sp_id, sp_data in schema['sub_pillars'].items():
            pillar_code = sp_id.rsplit('-', 1)[0] if '-' in sp_id else ''
            sp_lookup[sp_id] = {
                'name': sp_data.get('name', sp_id),
                'pillar_code': pillar_code,
            }
    if 'pillars' in schema:
        for p_code, p_data in schema['pillars'].items():
            pillar_lookup[p_code] = p_data.get('name', p_code)

    # Enrich phases with sub-pillar names + pillar grouping
    phases = plan.get('phases', {})
    enriched_phases = []
    for phase_key in sorted(phases.keys()):
        phase = phases[phase_key]
        target_sps = phase.get('target_sub_pillars', [])

        enriched_sps = []
        pillars_in_phase = {}  # pillar_code → [{code, name}]
        for sp in target_sps:
            info = sp_lookup.get(sp, {'name': sp, 'pillar_code': ''})
            pc = info['pillar_code']
            entry = {'code': sp, 'name': info['name'], 'pillar_code': pc, 'pillar_name': pillar_lookup.get(pc, pc)}
            enriched_sps.append(entry)
            pillars_in_phase.setdefault(pc, []).append({'code': sp, 'name': info['name'], 'pillar_name': pillar_lookup.get(pc, pc)})

        enriched_phases.append({
            'key': phase_key,
            'name': phase.get('name', ''),
            'activities': phase.get('activities', []),
            'ai_rmf_functions': phase.get('ai_rmf_functions', []),
            'target_sub_pillars': enriched_sps,
            'pillars_in_phase': pillars_in_phase,
            'week_start': _parse_week_start(phase_key),
            'week_end': _parse_week_end(phase_key),
        })

    return dataset_json_response({
        'description': plan.get('description', ''),
        'self_assessment_schedule': plan.get('self_assessment_schedule', ''),
        'milestones': plan.get('milestones', []),
        'phases': enriched_phases,
        'total_weeks': 12,
    }, app_state.current_schema_file)


def _parse_week_start(key):
    """Extract start week from phase key like 'weeks_1_3'."""
    import re
    m = re.search(r'(\d+)_(\d+)$', key)
    return int(m.group(1)) if m else 1

def _parse_week_end(key):
    """Extract end week from phase key like 'weeks_1_3'."""
    import re
    m = re.search(r'(\d+)_(\d+)$', key)
    return int(m.group(2)) if m else 12


@app.route('/api/update-adoption-plan', methods=['POST'])
def update_adoption_plan():
    """Update a phase in the adoption plan and persist to schema JSON."""
    data = request.json
    phase_key = data.get('key')
    if not phase_key:
        return jsonify({'error': 'Missing phase key'}), 400

    schema_filename = app_state.current_schema_file
    try:
        full_schema = read_dataset(schema_filename)

        reg = SCHEMA_REGISTRY.get(schema_filename)
        top_key = reg['top_key'] if reg else list(full_schema.keys())[0]
        schema = full_schema.get(top_key, full_schema)

        plan = schema.get('90_day_adoption_plan', {})
        phases = plan.get('phases', {})

        if phase_key not in phases:
            return jsonify({'error': f'Phase {phase_key} not found'}), 404

        phase = phases[phase_key]
        # Update fields if provided
        if 'name' in data:
            phase['name'] = data['name']
        if 'activities' in data:
            phase['activities'] = data['activities']
        if 'ai_rmf_functions' in data:
            phase['ai_rmf_functions'] = data['ai_rmf_functions']
        if 'target_sub_pillars' in data:
            phase['target_sub_pillars'] = data['target_sub_pillars']
        if 'week_start' in data and 'week_end' in data:
            # Rename the phase key if week range changed
            new_key = f"weeks_{data['week_start']}_{data['week_end']}"
            if new_key != phase_key:
                phases[new_key] = phases.pop(phase_key)
                phase_key = new_key

        write_error = persist_dataset(schema_filename, full_schema)
        if write_error is not None:
            return write_error

        return jsonify({'success': True, 'phase_key': phase_key})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-pricing', methods=['GET'])
def get_mdr_pricing():
    """Return MDR pricing analysis data for the pricing report."""
    import statistics
    pricing_sources = canonical_paths('vendor_pricing')
    if not pricing_sources:
        return jsonify({'error': 'MDR pricing data file not found'}), 404
    pricing_file = pricing_sources[0]
    is_enriched = 'Enriched' in pricing_file

    try:
        data = read_dataset(pricing_file)

        vendors = data.get('vendors', [])
        dims = data.get('dimensions', [])
        dim_labels = data.get('dimension_labels', {})

        # Per-dimension statistics
        dim_stats = {}
        for dim in dims:
            scores = [v.get('pricing_dimension_scores', {}).get(dim) for v in vendors]
            scores = [s for s in scores if s is not None]
            dim_stats[dim] = {
                'label': dim_labels.get(dim, dim),
                'mean': round(statistics.mean(scores), 2) if scores else 0,
                'median': round(statistics.median(scores), 1) if scores else 0,
                'min': min(scores) if scores else 0,
                'max': max(scores) if scores else 0,
                'count': len(scores),
            }

        # Per-vendor summary (sorted by pricing_overall_score desc)
        vendor_summaries = []
        for v in vendors:
            vs = {
                'vendor': v.get('vendor', ''),
                'website': v.get('website', ''),
                'headquarters': v.get('headquarters', ''),
                'region': v.get('region', ''),
                'mdr_service_type': v.get('mdr_service_type', ''),
                'delivery_model': v.get('delivery_model', ''),
                'target_market': v.get('target_market', ''),
                'pricing_model_type': v.get('pricing_model_type', ''),
                'description': v.get('description', ''),
                'product_names': v.get('product_names', []),
                'pricing_dimension_scores': v.get('pricing_dimension_scores', {}),
                'pricing_overall_score': v.get('pricing_overall_score', 0),
                'outcome_maturity_rating': v.get('outcome_maturity_rating', 0),
                'pricing_analysis': v.get('pricing_analysis', ''),
                'pricing_model_details': v.get('pricing_model_details', {}),
                'pricing_dimension_scores_v2': v.get('pricing_dimension_scores_v2', {}),
                'pricing_overall_score_v2': v.get('pricing_overall_score_v2', 0),
                'pricing_dimension_rationale_v2_text': v.get('pricing_dimension_rationale_v2_text', {}),
                'outcome_maturity_rationale_v2': v.get('outcome_maturity_rationale_v2', ''),
                'outcome_signals_v2': v.get('outcome_signals_v2', []),
                'pricing_research_confidence': v.get('pricing_research_confidence', ''),
                'pricing_evidence': v.get('pricing_evidence', {}),
            }
            # AI enrichment fields (present in 2-1 file)
            if is_enriched:
                vs['ai_pricing_influence'] = v.get('ai_pricing_influence', 0)
                vs['ai_pricing_influence_label'] = v.get('ai_pricing_influence_label', '')
                vs['ai_pricing_narrative'] = v.get('ai_pricing_narrative', '')
                vs['ai_capability_scores'] = v.get('ai_capability_scores', {})
                vs['pricing_strengths'] = v.get('pricing_strengths', [])
                vs['pricing_weaknesses'] = v.get('pricing_weaknesses', [])
                vs['pricing_roadmap'] = v.get('pricing_roadmap', [])
                vs['pricing_recommendations'] = v.get('pricing_recommendations', [])
            vendor_summaries.append(vs)
        vendor_summaries.sort(key=lambda x: x.get('pricing_overall_score', 0), reverse=True)

        # Cohort breakdowns
        model_types = {}
        regions = {}
        svc_types = {}
        delivery_models = {}
        target_markets = {}
        for v in vendors:
            mt = v.get('pricing_model_type', 'Unknown')
            model_types[mt] = model_types.get(mt, 0) + 1
            r = v.get('region', 'Unknown')
            regions[r] = regions.get(r, 0) + 1
            st = v.get('mdr_service_type', 'Unknown')
            svc_types[st] = svc_types.get(st, 0) + 1
            dm = v.get('delivery_model', 'Unknown')
            delivery_models[dm] = delivery_models.get(dm, 0) + 1
            tm = v.get('target_market', 'Unknown')
            target_markets[tm] = target_markets.get(tm, 0) + 1

        # Mean pricing scores by model type
        model_type_scores = {}
        for v in vendors:
            mt = v.get('pricing_model_type', 'Unknown')
            if mt not in model_type_scores:
                model_type_scores[mt] = []
            model_type_scores[mt].append(v.get('pricing_overall_score', 0))
        model_type_avgs = {k: round(statistics.mean(v), 2) for k, v in model_type_scores.items()}

        # Mean pricing scores by service type
        svc_type_scores = {}
        for v in vendors:
            st = v.get('mdr_service_type', 'Unknown')
            if st not in svc_type_scores:
                svc_type_scores[st] = []
            svc_type_scores[st].append(v.get('pricing_overall_score', 0))
        svc_type_avgs = {k: round(statistics.mean(v), 2) for k, v in svc_type_scores.items()}

        # Mean pricing scores by delivery model
        dm_scores = {}
        for v in vendors:
            dm = v.get('delivery_model', 'Unknown')
            if dm not in dm_scores:
                dm_scores[dm] = []
            dm_scores[dm].append(v.get('pricing_overall_score', 0))
        dm_avgs = {k: round(statistics.mean(v), 2) for k, v in dm_scores.items()}

        # Overall stats
        all_scores = [v.get('pricing_overall_score', 0) for v in vendors]
        all_outcome = [v.get('outcome_maturity_rating', 0) for v in vendors]

        return jsonify({
            'vendor_count': len(vendors),
            'dimensions': dims,
            'dimension_labels': dim_labels,
            'dimension_stats': dim_stats,
            'summary': data.get('summary', {}),
            'is_ai_enriched': is_enriched,
            'ai_influence_stats': data.get('ai_influence_stats', {}),
            'pricing_benchmarks': data.get('pricing_benchmarks', {}),
            'overall_stats': {
                'mean': round(statistics.mean(all_scores), 2) if all_scores else 0,
                'median': round(statistics.median(all_scores), 2) if all_scores else 0,
                'min': round(min(all_scores), 1) if all_scores else 0,
                'max': round(max(all_scores), 1) if all_scores else 0,
            },
            'outcome_stats': {
                'mean': round(statistics.mean(all_outcome), 2) if all_outcome else 0,
                'median': round(statistics.median(all_outcome), 1) if all_outcome else 0,
            },
            'cohorts': {
                'pricing_model_type': model_types,
                'region': regions,
                'mdr_service_type': svc_types,
                'delivery_model': delivery_models,
                'target_market': target_markets,
            },
            'cohort_scores': {
                'by_model_type': model_type_avgs,
                'by_service_type': svc_type_avgs,
                'by_delivery_model': dm_avgs,
            },
            'vendors': vendor_summaries,
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-market-insight/perspectives', methods=['GET'])
def get_mdr_market_insight_perspectives():
    """Return the list of available report perspectives."""
    json_file = os.path.join(os.path.dirname(__file__), 'mdr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Market Insight reports file not found'}), 404
    try:
        data = read_dataset('mdr_market_insight_reports.json')
        perspectives = [{'id': r['id'], 'label': r['label']} for r in data.get('reports', [])]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-mq-scores', methods=['GET'])
def get_mdr_mq_scores():
    """Return the MDR Magic Quadrant scoring data."""
    json_file = os.path.join(os.path.dirname(__file__), 'MDR Services Vendor MQ Scores.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'MQ Scores file not found'}), 404
    try:
        data = read_dataset('MDR Services Vendor MQ Scores.json')
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cnapp-mq-scores', methods=['GET'])
def get_cnapp_mq_scores():
    """Return the CNAPP Magic Quadrant scoring data.

    Query params:
        mode: 'v1' (heuristic, default) or 'v2' (evidence-enriched).
    """
    mode = (request.args.get('mode') or 'v1').lower()
    fname = 'CNAPP Vendor MQ Scores v2.json' if mode == 'v2' else 'CNAPP Vendor MQ Scores.json'
    json_file = os.path.join(os.path.dirname(__file__), fname)
    if not os.path.exists(json_file):
        return jsonify({'error': f'CNAPP MQ Scores file not found ({fname})'}), 404
    try:
        data = read_dataset(fname)
        data.setdefault('score_mode', mode)
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cnapp-mq-market-insight/perspectives', methods=['GET'])
def get_cnapp_mq_market_insight_perspectives():
    """Return the list of available CNAPP MQ Market Insight perspectives."""
    json_file = os.path.join(os.path.dirname(__file__), 'cnapp_mq_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'CNAPP MQ Market Insight reports file not found'}), 404
    try:
        data = read_dataset(
            'cnapp_mq_market_insight_reports.json'
        )
        perspectives = [{'id': r['id'], 'label': r['label']} for r in data.get('reports', [])]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cnapp-mq-market-insight', methods=['GET'])
def get_cnapp_mq_market_insight():
    """Return a CNAPP MQ Market Insight report for a given perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'cnapp_mq_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'CNAPP MQ Market Insight reports file not found'}), 404
    try:
        data = read_dataset(
            'cnapp_mq_market_insight_reports.json'
        )
        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No report perspectives available'}), 404
        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404
        return jsonify({
            'id': report['id'],
            'label': report['label'],
            'title': report['title'],
            'summary': report['summary'],
            'spa': report['spa'],
            'findings': report['findings'],
            'recommendations': report['recommendations'],
            'analysis_sections': report['analysis_sections'],
            'background': report['background'],
            'impact': report['impact'],
            'conclusion': report['conclusion'],
            'glossary': report['glossary'],
            'evidence': report['evidence'],
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cnapp-mq-market-insight', methods=['POST'])
def save_cnapp_mq_market_insight():
    """Save edits to a CNAPP MQ Market Insight report perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'cnapp_mq_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'CNAPP MQ Market Insight reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400
        data = read_dataset(
            'cnapp_mq_market_insight_reports.json'
        )
        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404
        editable_fields = ['title', 'summary', 'spa', 'background', 'impact', 'conclusion']
        for field in editable_fields:
            if field in payload:
                report[field] = payload[field]
        if 'findings' in payload and isinstance(payload['findings'], list):
            report['findings'] = payload['findings']
        if 'recommendations' in payload and isinstance(payload['recommendations'], list):
            report['recommendations'] = payload['recommendations']
        if 'analysis_sections' in payload and isinstance(payload['analysis_sections'], list):
            report['analysis_sections'] = payload['analysis_sections']
        write_error = persist_dataset(
            'cnapp_mq_market_insight_reports.json',
            data,
        )
        if write_error is not None:
            return write_error
        return jsonify({'success': True, 'message': f'Report "{perspective_id}" saved.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-market-insight', methods=['GET'])
def get_mdr_market_insight():
    """Return the MDR Pricing Market Insight for a given perspective as structured JSON."""
    json_file = os.path.join(os.path.dirname(__file__), 'mdr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Market Insight reports file not found'}), 404
    try:
        data = read_dataset('mdr_market_insight_reports.json')

        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No report perspectives available'}), 404

        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        return jsonify({
            'id': report['id'],
            'label': report['label'],
            'title': report['title'],
            'summary': report['summary'],
            'spa': report['spa'],
            'findings': report['findings'],
            'recommendations': report['recommendations'],
            'analysis_sections': report['analysis_sections'],
            'background': report['background'],
            'impact': report['impact'],
            'conclusion': report['conclusion'],
            'glossary': report['glossary'],
            'evidence': report['evidence'],
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-market-insight', methods=['POST'])
def save_mdr_market_insight():
    """Save edits to an MDR Market Insight report perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'mdr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'MDR Market Insight reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400

        data = read_dataset('mdr_market_insight_reports.json')

        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        editable_fields = ['title', 'summary', 'spa', 'background', 'impact', 'conclusion']
        for field in editable_fields:
            if field in payload:
                report[field] = payload[field]

        if 'findings' in payload and isinstance(payload['findings'], list):
            report['findings'] = payload['findings']

        if 'recommendations' in payload and isinstance(payload['recommendations'], list):
            report['recommendations'] = payload['recommendations']

        if 'analysis_sections' in payload and isinstance(payload['analysis_sections'], list):
            report['analysis_sections'] = payload['analysis_sections']

        write_error = persist_dataset('mdr_market_insight_reports.json', data)
        if write_error is not None:
            return write_error

        return jsonify({'success': True, 'message': f'Report "{perspective_id}" saved.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ─── PreCyber statistics helpers ───────────────────────────────────────────────

def _get_precyber_vendor_file():
    """Return the active PreCyber vendor source selected by the manifest."""
    for source in _CANONICAL_CATALOG:
        if (
            source.get('status') == 'active'
            and source.get('kind') == 'vendor_score'
            and source.get('market') == 'precyber'
        ):
            return source['path']
    return None


def _compute_precyber_stats(vendor_file):
    """Compute all PreCyber statistics from a vendor JSON file and return a dict."""
    vendors = dataset_repository.read_vendors(vendor_file)

    pillars = ['EXM', 'AMT', 'ADR', 'PPM', 'SVC']
    pillar_labels = {
        'EXM': 'Exposure Management',
        'AMT': 'Adversary Management & Threat Intelligence',
        'ADR': 'Adversary Disruption',
        'PPM': 'Posture & Policy Management',
        'SVC': 'Services Maturity & Delivery',
    }
    threshold = 2.0
    n = len(vendors)

    def get_pillar_avg(vendor, pillar_code):
        sub = vendor.get('sub_pillar_scores_current', {})
        keys = [k for k in sub if k.startswith(pillar_code + '-') and sub[k] >= 1]
        if keys:
            return sum(sub[k] for k in keys) / len(keys)
        return vendor.get('pillar_scores', {}).get(pillar_code, 0)

    # Pillar penetration
    pillar_penetration = {}
    pillar_avgs = {}
    for p in pillars:
        scores = [get_pillar_avg(v, p) for v in vendors]
        above = sum(1 for s in scores if s >= threshold)
        pillar_penetration[p] = {
            'label': pillar_labels[p],
            'pct': round(above / n * 100) if n else 0,
            'count': above,
            'avg': round(sum(scores) / n, 2) if n else 0,
        }
        pillar_avgs[p] = pillar_penetration[p]['avg']

    # Coverage distribution (0–5 pillars meeting threshold)
    coverage_dist = {i: 0 for i in range(6)}
    for v in vendors:
        cnt = sum(1 for p in pillars if get_pillar_avg(v, p) >= threshold)
        coverage_dist[cnt] += 1

    full_spectrum = coverage_dist.get(5, 0)
    majority_spectrum = coverage_dist.get(4, 0)
    narrow = sum(coverage_dist.get(i, 0) for i in range(4))
    avg_coverage = round(sum(
        sum(1 for p in pillars if get_pillar_avg(v, p) >= threshold)
        for v in vendors
    ) / n, 1) if n else 0

    # Blind spot (>=1 pillar below threshold) and no-AMT vendors
    blind_spot = sum(1 for v in vendors if any(get_pillar_avg(v, p) < threshold for p in pillars))
    no_amt = sum(1 for v in vendors if get_pillar_avg(v, 'AMT') < threshold)

    # Delivery model stats — group vendors by delivery_model key
    dm_vendors = {}
    for v in vendors:
        dm_key = v.get('delivery_model', 'unknown')
        dm_vendors.setdefault(dm_key, []).append(v)

    dm_labels = {
        'direct_service': 'Direct Service Providers',
        'platform_plus_partner': 'Platform + Partner',
        'platform_only': 'Platform Only',
    }

    delivery_models = {}
    for dm_key, dv in dm_vendors.items():
        c = len(dv)
        pillar_totals = {p: 0.0 for p in pillars}
        coverage_sum = 0
        for v in dv:
            for p in pillars:
                pillar_totals[p] += get_pillar_avg(v, p)
            coverage_sum += sum(1 for p in pillars if get_pillar_avg(v, p) >= threshold)
        pillar_avgs_dm = {p: round(pillar_totals[p] / c, 2) if c else 0 for p in pillars}
        pillar_below_pct = {
            p: round(sum(1 for v in dv if get_pillar_avg(v, p) < threshold) * 100 / c) if c else 0
            for p in pillars
        }
        delivery_models[dm_key] = {
            'label': dm_labels.get(dm_key, dm_key),
            'count': c,
            'avg_coverage': round(coverage_sum / c, 1) if c else 0,
            'pillar_avgs': pillar_avgs_dm,
            'pillar_below_pct': pillar_below_pct,
            'svc_avg': pillar_avgs_dm['SVC'],
            'overall_avg': round(sum(pillar_avgs_dm.values()) / len(pillars), 2) if pillars else 0,
        }

    # Top balanced vendors
    vendor_balance = []
    for v in vendors:
        scores = {p: get_pillar_avg(v, p) for p in pillars}
        vendor_balance.append({
            'vendor': v['vendor'],
            'min_score': round(min(scores.values()), 2),
            'delivery_model': v.get('delivery_model', ''),
            'pillar_scores': {p: round(scores[p], 2) for p in pillars},
            'coverage': sum(1 for s in scores.values() if s >= threshold),
        })
    vendor_balance.sort(key=lambda x: x['min_score'], reverse=True)
    top_balanced = vendor_balance[:6]

    _sub_pillar_labels = {
        'AMT-01': 'Polymorphic & Morphing Defense',
        'AMT-02': 'Runtime Application Protection',
        'AMT-03': 'Dynamic Network & Infrastructure Defense',
        'AMT-04': 'Identity & Credential Rotation',
        'AMT-05': 'AMTD Services Maturity',
        'SVC-01': 'Implementation & Onboarding',
        'SVC-02': 'Consultative & Advisory Services',
        'SVC-03': 'Managed Operations & Continuous Delivery',
        'SVC-04': 'AI-Driven & Autonomous Delivery',
        'SVC-05': 'SVC-05',
    }

    amt_sub_pillars = {}
    for key in ['AMT-01', 'AMT-02', 'AMT-03', 'AMT-04', 'AMT-05']:
        scores = [v.get('sub_pillar_scores_current', {}).get(key, 0) for v in vendors]
        above = sum(1 for s in scores if s >= threshold)
        label = vendors[0].get('sub_pillar_schema_labels', {}).get(key) or _sub_pillar_labels.get(key, key)
        amt_sub_pillars[key] = {
            'label': label,
            'avg': round(sum(scores) / n, 2) if n else 0,
            'pct_above': round(above / n * 100) if n else 0,
        }

    svc_sub_pillars = {}
    for key in ['SVC-01', 'SVC-02', 'SVC-03', 'SVC-04']:
        scores = [v.get('sub_pillar_scores_current', {}).get(key, 0) for v in vendors]
        above = sum(1 for s in scores if s >= threshold)
        label = vendors[0].get('sub_pillar_schema_labels', {}).get(key) or _sub_pillar_labels.get(key, key)
        svc_sub_pillars[key] = {
            'label': label,
            'avg': round(sum(scores) / n, 2) if n else 0,
            'pct_above': round(above / n * 100) if n else 0,
        }

    return {
        'vendor_count': n,
        'pillars': pillars,
        'pillar_labels': pillar_labels,
        'pillar_penetration': pillar_penetration,
        'pillar_avgs': pillar_avgs,
        'coverage_distribution': coverage_dist,
        'full_spectrum_count': full_spectrum,
        'full_spectrum_pct': round(full_spectrum / n * 100) if n else 0,
        'majority_spectrum_count': majority_spectrum,
        'majority_spectrum_pct': round(majority_spectrum / n * 100) if n else 0,
        'narrow_count': narrow,
        'narrow_pct': round(narrow / n * 100) if n else 0,
        'avg_coverage': avg_coverage,
        'delivery_models': delivery_models,
        'top_balanced': top_balanced,
        'amt_sub_pillars': amt_sub_pillars,
        'svc_sub_pillars': svc_sub_pillars,
        'blind_spot_count': blind_spot,
        'blind_spot_pct': round(blind_spot / n * 100) if n else 0,
        'no_amt_count': no_amt,
        'no_amt_pct': round(no_amt / n * 100) if n else 0,
    }


@app.route('/api/precyber-stats', methods=['GET'])
def get_precyber_stats():
    """Return computed statistics from PreCyber vendor data for live visualizations."""
    vendor_file = _get_precyber_vendor_file()
    if not vendor_file:
        return jsonify({'error': 'PreCyber vendor data not found'}), 404
    try:
        stats = _compute_precyber_stats(vendor_file)
        return jsonify(stats)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ═══════════════════════════════════════════════════════════════════════
#  Product Market Readiness — Credibility Gap Analytics
# ═══════════════════════════════════════════════════════════════════════

@app.route('/api/pmr-stats', methods=['GET'])
def get_pmr_stats():
    """Return computed statistics from PMR vendor data for credibility gap visualizations."""
    import statistics as stats_mod
    vendor_file = os.path.join(os.path.dirname(__file__),
                               'Product Market Readiness Vendor 1-0 Seed.json')
    if not os.path.exists(vendor_file):
        return jsonify({'error': 'PMR vendor data not found'}), 404
    try:
        data = read_dataset(
            'Product Market Readiness Vendor 1-0 Seed.json'
        )

        vendors = data.get('vendors', [])
        n = len(vendors)
        pillars = ['PPD', 'PCS', 'TDT', 'PCM', 'CTL']
        pillar_labels = {
            'PPD': 'Product Positioning & Differentiation',
            'PCS': 'Proof Points & Case Studies',
            'TDT': 'Technical Depth & Transparency',
            'PCM': 'Pricing & Commercial Model Clarity',
            'CTL': 'Content & Thought Leadership',
        }

        # Per-pillar GTM / Proof / Gap averages
        pillar_stats = {}
        for p in pillars:
            gtm_scores = [v.get('pillar_gtm_scores', {}).get(p, 0) for v in vendors]
            proof_scores = [v.get('pillar_proof_scores', {}).get(p, 0) for v in vendors]
            gaps = [v.get('pillar_gaps', {}).get(p, 0) for v in vendors]
            pillar_stats[p] = {
                'label': pillar_labels[p],
                'gtm_avg': round(stats_mod.mean(gtm_scores), 2) if gtm_scores else 0,
                'proof_avg': round(stats_mod.mean(proof_scores), 2) if proof_scores else 0,
                'gap_avg': round(stats_mod.mean(gaps), 2) if gaps else 0,
                'gtm_median': round(stats_mod.median(gtm_scores), 2) if gtm_scores else 0,
                'proof_median': round(stats_mod.median(proof_scores), 2) if proof_scores else 0,
            }

        # Overall averages
        all_gtm = [v.get('overall_gtm_score', 0) for v in vendors]
        all_proof = [v.get('overall_proof_score', 0) for v in vendors]
        all_gaps = [v.get('overall_credibility_gap', 0) for v in vendors]

        # Coverage grade distribution
        grade_dist = {}
        for v in vendors:
            g = v.get('coverage_grade', 'F')
            grade_dist[g] = grade_dist.get(g, 0) + 1

        # Gap severity distribution across all vendors
        gap_severity_totals = {'aligned': 0, 'minor': 0, 'moderate': 0, 'significant': 0, 'critical': 0}
        for v in vendors:
            gsd = v.get('gap_severity_distribution', {})
            for k in gap_severity_totals:
                gap_severity_totals[k] += gsd.get(k, 0)

        # Vendor type breakdown
        type_counts = {}
        type_gtm = {}
        type_proof = {}
        for v in vendors:
            vt = v.get('vendor_type', 'Unknown')
            type_counts[vt] = type_counts.get(vt, 0) + 1
            if vt not in type_gtm:
                type_gtm[vt] = []
                type_proof[vt] = []
            type_gtm[vt].append(v.get('overall_gtm_score', 0))
            type_proof[vt].append(v.get('overall_proof_score', 0))
        vendor_type_stats = {}
        for vt in type_counts:
            vendor_type_stats[vt] = {
                'count': type_counts[vt],
                'gtm_avg': round(stats_mod.mean(type_gtm[vt]), 2) if type_gtm[vt] else 0,
                'proof_avg': round(stats_mod.mean(type_proof[vt]), 2) if type_proof[vt] else 0,
                'gap_avg': round(stats_mod.mean(type_gtm[vt]), 2) - round(stats_mod.mean(type_proof[vt]), 2) if type_gtm[vt] else 0,
            }

        # Top over-claimers (highest positive gap) and under-marketers (most negative gap)
        vendor_gaps = []
        for v in vendors:
            vendor_gaps.append({
                'vendor': v.get('vendor', ''),
                'overall_gtm': v.get('overall_gtm_score', 0),
                'overall_proof': v.get('overall_proof_score', 0),
                'overall_gap': v.get('overall_credibility_gap', 0),
                'coverage_grade': v.get('coverage_grade', 'F'),
                'vendor_type': v.get('vendor_type', ''),
                'pillar_gtm': v.get('pillar_gtm_scores', {}),
                'pillar_proof': v.get('pillar_proof_scores', {}),
                'pillar_gaps': v.get('pillar_gaps', {}),
            })
        over_claimers = sorted(vendor_gaps, key=lambda x: x['overall_gap'], reverse=True)[:10]
        under_marketers = sorted(vendor_gaps, key=lambda x: x['overall_gap'])[:10]
        best_aligned = sorted(vendor_gaps, key=lambda x: abs(x['overall_gap']))[:10]

        # Cross-schema reference stats
        source_schema_counts = {}
        for v in vendors:
            for s in v.get('source_schemas', []):
                source_schema_counts[s] = source_schema_counts.get(s, 0) + 1

        return jsonify({
            'vendor_count': n,
            'pillars': pillars,
            'pillar_labels': pillar_labels,
            'pillar_stats': pillar_stats,
            'overall_stats': {
                'gtm_mean': round(stats_mod.mean(all_gtm), 2) if all_gtm else 0,
                'proof_mean': round(stats_mod.mean(all_proof), 2) if all_proof else 0,
                'gap_mean': round(stats_mod.mean(all_gaps), 2) if all_gaps else 0,
                'gtm_median': round(stats_mod.median(all_gtm), 2) if all_gtm else 0,
                'proof_median': round(stats_mod.median(all_proof), 2) if all_proof else 0,
            },
            'grade_distribution': grade_dist,
            'gap_severity_totals': gap_severity_totals,
            'vendor_type_stats': vendor_type_stats,
            'over_claimers': over_claimers,
            'under_marketers': under_marketers,
            'best_aligned': best_aligned,
            'source_schema_coverage': source_schema_counts,
            'vendors': vendor_gaps,
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/pmr-market-insight/perspectives', methods=['GET'])
def get_pmr_market_insight_perspectives():
    """Return the list of available PMR report perspectives."""
    json_file = os.path.join(os.path.dirname(__file__), 'pmr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PMR Market Insight reports file not found'}), 404
    try:
        data = read_dataset('pmr_market_insight_reports.json')
        perspectives = [{'id': r['id'], 'label': r['label']} for r in data.get('reports', [])]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/pmr-market-insight', methods=['GET'])
def get_pmr_market_insight():
    """Return PMR Market Insight for a given perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'pmr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PMR Market Insight reports file not found'}), 404
    try:
        data = read_dataset('pmr_market_insight_reports.json')
        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No PMR report perspectives available'}), 404
        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404
        return jsonify({
            'id': report['id'],
            'label': report['label'],
            'title': report['title'],
            'summary': report['summary'],
            'spa': report['spa'],
            'findings': report['findings'],
            'recommendations': report['recommendations'],
            'analysis_sections': report['analysis_sections'],
            'background': report['background'],
            'impact': report['impact'],
            'conclusion': report['conclusion'],
            'glossary': report.get('glossary', {}),
            'evidence': report.get('evidence', {}),
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/pmr-market-insight', methods=['POST'])
def save_pmr_market_insight():
    """Save edits to a PMR Market Insight report perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'pmr_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PMR Market Insight reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400

        data = read_dataset('pmr_market_insight_reports.json')

        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        editable = ['title', 'summary', 'spa', 'background', 'impact', 'conclusion',
                     'findings', 'recommendations', 'analysis_sections']
        for field in editable:
            if field in payload:
                report[field] = payload[field]

        write_error = persist_dataset('pmr_market_insight_reports.json', data)
        if write_error is not None:
            return write_error

        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/precyber-market-insight/perspectives', methods=['GET'])
def get_precyber_market_insight_perspectives():
    """Return the list of available PreCyber report perspectives."""
    json_file = os.path.join(os.path.dirname(__file__), 'precyber_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PreCyber Market Insight reports file not found'}), 404
    try:
        data = read_dataset(
            'precyber_market_insight_reports.json'
        )
        perspectives = [{'id': r['id'], 'label': r['label']} for r in data.get('reports', [])]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/precyber-market-insight', methods=['GET'])
def get_precyber_market_insight():
    """Return the PreCyber Market Insight for a given perspective as structured JSON."""
    json_file = os.path.join(os.path.dirname(__file__), 'precyber_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PreCyber Market Insight reports file not found'}), 404
    try:
        data = read_dataset(
            'precyber_market_insight_reports.json'
        )

        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No report perspectives available'}), 404

        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        return jsonify({
            'id': report['id'],
            'label': report['label'],
            'title': report['title'],
            'summary': report['summary'],
            'spa': report['spa'],
            'findings': report['findings'],
            'recommendations': report['recommendations'],
            'analysis_sections': report['analysis_sections'],
            'background': report['background'],
            'impact': report['impact'],
            'conclusion': report['conclusion'],
            'glossary': report['glossary'],
            'evidence': report['evidence'],
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/precyber-market-insight', methods=['POST'])
def save_precyber_market_insight():
    """Save edits to a PreCyber Market Insight report perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'precyber_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'PreCyber Market Insight reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400

        data = read_dataset(
            'precyber_market_insight_reports.json'
        )

        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        editable_fields = ['title', 'summary', 'spa', 'background', 'impact', 'conclusion']
        for field in editable_fields:
            if field in payload:
                report[field] = payload[field]

        if 'findings' in payload and isinstance(payload['findings'], list):
            report['findings'] = payload['findings']

        if 'recommendations' in payload and isinstance(payload['recommendations'], list):
            report['recommendations'] = payload['recommendations']

        if 'analysis_sections' in payload and isinstance(payload['analysis_sections'], list):
            report['analysis_sections'] = payload['analysis_sections']

        write_error = persist_dataset(
            'precyber_market_insight_reports.json',
            data,
        )
        if write_error is not None:
            return write_error

        return jsonify({'success': True, 'message': f'Report "{perspective_id}" saved.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── DFIR Market Insight endpoints ──────────────────────────────────────

@app.route('/api/dfir-market-insight/perspectives', methods=['GET'])
def get_dfir_market_insight_perspectives():
    """Return the list of available DFIR report perspectives."""
    json_file = os.path.join(os.path.dirname(__file__), 'dfir_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'DFIR Market Insight reports file not found'}), 404
    try:
        data = read_dataset('dfir_market_insight_reports.json')
        perspectives = [{'id': r['id'], 'label': r['label']} for r in data.get('reports', [])]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/dfir-market-insight', methods=['GET'])
def get_dfir_market_insight():
    """Return the DFIR Market Insight for a given perspective as structured JSON."""
    json_file = os.path.join(os.path.dirname(__file__), 'dfir_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'DFIR Market Insight reports file not found'}), 404
    try:
        data = read_dataset('dfir_market_insight_reports.json')

        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No report perspectives available'}), 404

        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        return jsonify({
            'id': report['id'],
            'label': report['label'],
            'title': report['title'],
            'summary': report['summary'],
            'spa': report['spa'],
            'findings': report['findings'],
            'recommendations': report['recommendations'],
            'analysis_sections': report['analysis_sections'],
            'background': report['background'],
            'impact': report['impact'],
            'conclusion': report['conclusion'],
            'glossary': report['glossary'],
            'evidence': report['evidence'],
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/dfir-market-insight', methods=['POST'])
def save_dfir_market_insight():
    """Save edits to a DFIR Market Insight report perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'dfir_market_insight_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'DFIR Market Insight reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400

        data = read_dataset('dfir_market_insight_reports.json')

        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        editable_fields = ['title', 'summary', 'spa', 'background', 'impact', 'conclusion']
        for field in editable_fields:
            if field in payload:
                report[field] = payload[field]

        if 'findings' in payload and isinstance(payload['findings'], list):
            report['findings'] = payload['findings']

        if 'recommendations' in payload and isinstance(payload['recommendations'], list):
            report['recommendations'] = payload['recommendations']

        if 'analysis_sections' in payload and isinstance(payload['analysis_sections'], list):
            report['analysis_sections'] = payload['analysis_sections']

        write_error = persist_dataset('dfir_market_insight_reports.json', data)
        if write_error is not None:
            return write_error

        return jsonify({'success': True, 'message': f'Report "{perspective_id}" saved.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Analyst Take (cross-schema) ──────────────────────────────────────────

@app.route('/api/analyst-take/perspectives', methods=['GET'])
def get_analyst_take_perspectives():
    """Return list of available Analyst Take perspectives, filtered by schema."""
    json_file = os.path.join(os.path.dirname(__file__), 'analyst_take_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Analyst Take reports file not found'}), 404
    try:
        data = read_dataset('analyst_take_reports.json')
        schema = request.args.get('schema', '')
        reports = data.get('reports', [])
        # Filter: always include template (schema_ref is null) + reports matching the active schema
        if schema:
            reports = [r for r in reports if not r.get('schema_ref') or r.get('schema_ref') == schema]
        else:
            # No schema specified – show only template
            reports = [r for r in reports if not r.get('schema_ref')]
        perspectives = [{'id': r['id'], 'label': r['label']} for r in reports]
        return jsonify({'perspectives': perspectives})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyst-take', methods=['GET'])
def get_analyst_take():
    """Return an Analyst Take report for a given perspective."""
    json_file = os.path.join(os.path.dirname(__file__), 'analyst_take_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Analyst Take reports file not found'}), 404
    try:
        data = read_dataset('analyst_take_reports.json')

        reports = data.get('reports', [])
        if not reports:
            return jsonify({'error': 'No Analyst Take perspectives available'}), 404

        perspective_id = request.args.get('perspective', reports[0]['id'])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        return jsonify(report)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyst-take', methods=['POST'])
def save_analyst_take():
    """Save edits to an Analyst Take report."""
    json_file = os.path.join(os.path.dirname(__file__), 'analyst_take_reports.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Analyst Take reports file not found'}), 404
    try:
        payload = request.json
        perspective_id = payload.get('id')
        if not perspective_id:
            return jsonify({'error': 'Missing perspective id'}), 400

        data = read_dataset('analyst_take_reports.json')

        reports = data.get('reports', [])
        report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            return jsonify({'error': f'Perspective "{perspective_id}" not found'}), 404

        for field in ['title', 'subtitle', 'notes']:
            if field in payload:
                report[field] = payload[field]

        if 'body_sections' in payload and isinstance(payload['body_sections'], list):
            report['body_sections'] = payload['body_sections']

        if 'recommended_reading' in payload and isinstance(payload['recommended_reading'], list):
            report['recommended_reading'] = payload['recommended_reading']

        write_error = persist_dataset('analyst_take_reports.json', data)
        if write_error is not None:
            return write_error

        return jsonify({'success': True, 'message': f'Analyst Take "{perspective_id}" saved.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/dfir-infographic-pptx', methods=['GET'])
def dfir_infographic_pptx():
    """Generate an editable PowerPoint slide of the DFIR Agentic Shift infographic."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        def rgb(hex_str):
            h = hex_str.lstrip('#')
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def add_box(left, top, width, height, fill_hex=None, border_hex=None):
            shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
            shape.fill.background() if not fill_hex else None
            if fill_hex:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(fill_hex)
            if border_hex:
                shape.line.color.rgb = rgb(border_hex)
                shape.line.width = Pt(1.5)
            else:
                shape.line.fill.background()
            return shape

        def add_textbox(left, top, width, height, text, size=11, bold=False, color_hex='333333', align='left'):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = rgb(color_hex)
            if align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            return txBox

        def add_multiline_textbox(left, top, width, height, lines, size=11, color_hex='555555', align='left', spacing=None):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(size)
                p.font.bold = bld
                p.font.color.rgb = rgb(col) if col else rgb(color_hex)
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                if spacing:
                    p.space_after = Pt(spacing)
            return txBox

        # Conversion helpers (slide is 13.333" x 7.5", SVG was 1200x900)
        def sx(v): return int(v / 1200 * Inches(13.333))
        def sy(v): return int(v / 900 * Inches(7.5))
        def sw(v): return int(v / 1200 * Inches(13.333))
        def sh(v): return int(v / 900 * Inches(7.5))

        # ── Title banner ──
        add_box(sx(0), sy(0), sw(1200), sh(70), '1a3a5c', '1a3a5c')
        add_textbox(sx(50), sy(5), sw(1100), sh(32),
                    'The Agentic Shift: Revolutionizing Digital Forensics & Incident Response (DFIR)',
                    22, True, 'ffffff', 'center')
        add_textbox(sx(100), sy(38), sw(1000), sh(28),
                    'From manual, human-dependent investigation \u2192 high-speed, transparent, legally defensible AI models',
                    11, False, 'ccd8e8', 'center')

        # ── Left section header ──
        add_box(sx(20), sy(85), sw(520), sh(32), 'e0ddd5', 'b0a898')
        add_textbox(sx(30), sy(88), sw(500), sh(28),
                    'TRADITIONAL MANUAL DFIR: FROM MANUAL GRINDING', 11, True, '5a503c', 'center')

        # ── Left Box 1: Hours to Days ──
        add_box(sx(30), sy(135), sw(240), sh(130), 'fff8f0', 'ca5010')
        add_multiline_textbox(sx(40), sy(145), sw(220), sh(120), [
            ('Hours to Days', 16, True, 'ca5010'),
            ('Analysis Speed', 13, True, '5a503c'),
            ('', 6, False, None),
            ('Manual data grinding through', 11, False, '777777'),
            ('massive evidence sets', 11, False, '777777'),
        ], align='center')

        # ── Left Box 2: Methodology ──
        add_box(sx(290), sy(135), sw(240), sh(130), 'f5f0f0', 'a80000')
        add_multiline_textbox(sx(300), sy(145), sw(220), sh(120), [
            ('Human-Driven / Opaque', 16, True, 'a80000'),
            ('Methodology', 13, True, '5a503c'),
            ('', 6, False, None),
            ('"Black box" concern blocks', 11, False, '777777'),
            ('AI adoption in investigation', 11, False, '777777'),
        ], align='center')

        # ── Left Box 3: 70% stat ──
        add_box(sx(30), sy(280), sw(240), sh(120), 'fff5f5', 'a80000')
        add_multiline_textbox(sx(40), sy(288), sw(220), sh(108), [
            ('70%', 28, True, 'a80000'),
            ('of Manual Models Replaced by 2030', 12, True, '5a503c'),
            ('', 4, False, None),
            ('Traditional forensic investigations are', 10, False, '777777'),
            ('losing ground to agentic solutions', 10, False, '777777'),
        ], align='center')

        # ── Left Box 4: Core Focus ──
        add_box(sx(290), sy(280), sw(240), sh(120), 'f8f5ff', '8764b8')
        add_multiline_textbox(sx(300), sy(288), sw(220), sh(108), [
            ('Core Focus:', 14, True, '5a503c'),
            ('Reporting & Admin', 14, True, 'ca5010'),
            ('', 4, False, None),
            ('AI applied to reports, planning,', 10, False, '777777'),
            ('compliance \u2014 not critical path', 10, False, '777777'),
        ], align='center')

        # ── Center: DFKG circle area ──
        add_box(sx(520), sy(420), sw(160), sh(160), 'e8f4ff', '0078d4')
        add_multiline_textbox(sx(525), sy(445), sw(150), sh(130), [
            ('Digital Forensic', 13, True, '1a3a5c'),
            ('Knowledge Graph', 13, True, '1a3a5c'),
            ('(DFKGs)', 12, False, '0078d4'),
            ('', 4, False, None),
            ('Visualizes linkages + timelines', 10, False, '555555'),
            ('at machine speed', 10, False, '555555'),
        ], align='center')

        # ── Center: Tandem Model ──
        add_box(sx(505), sy(600), sw(190), sh(100), 'f0fff0', '107c10')
        add_multiline_textbox(sx(510), sy(608), sw(180), sh(88), [
            ('Tandem Human-AI Model', 13, True, '107c10'),
            ('', 4, False, None),
            ('AI speed + human validation', 11, False, '555555'),
            ('overcomes black box trust', 11, False, '555555'),
        ], align='center')

        # ── Right section header ──
        add_box(sx(660), sy(85), sw(520), sh(32), '00b4d8', '0095b3')
        add_textbox(sx(670), sy(88), sw(500), sh(28),
                    'AGENTIC AI DFIR: TO MACHINE-SPEED ANALYSIS', 11, True, 'ffffff', 'center')

        # ── Right Box 1: Under 12 Minutes ──
        add_box(sx(670), sy(135), sw(240), sh(130), 'f0faff', '0078d4')
        add_multiline_textbox(sx(680), sy(145), sw(220), sh(120), [
            ('Under 12 Minutes', 16, True, '0078d4'),
            ('Investigation Speed', 13, True, '1a3a5c'),
            ('', 6, False, None),
            ('AI agents complete complex', 11, False, '555555'),
            ('investigations replacing hours', 11, False, '555555'),
            ('of manual data ingestion', 11, False, '555555'),
        ], align='center')

        # ── Right Box 2: 85% Accuracy ──
        add_box(sx(930), sy(135), sw(240), sh(130), 'f0fff4', '107c10')
        add_multiline_textbox(sx(940), sy(145), sw(220), sh(120), [
            ('85%', 28, True, '107c10'),
            ('Source Attribution', 13, True, '1a3a5c'),
            ('', 4, False, None),
            ('ForensicLLM pinpoints specific', 10, False, '555555'),
            ('file paths for every claim', 10, False, '555555'),
        ], align='center')

        # ── Right Box 3: Legal Admissibility ──
        add_box(sx(670), sy(280), sw(240), sh(120), 'f5f0ff', '8764b8')
        add_multiline_textbox(sx(680), sy(288), sw(220), sh(108), [
            ('Legal Admissibility', 14, True, '8764b8'),
            ('& Rule 901', 14, True, '8764b8'),
            ('', 4, False, None),
            ('Documented agent interactions meet', 10, False, '555555'),
            ('Daubert + Federal Rule standards', 10, False, '555555'),
        ], align='center')

        # ── Right Box 4: Explainable ──
        add_box(sx(930), sy(280), sw(240), sh(120), 'fffff0', '0078d4')
        add_multiline_textbox(sx(940), sy(288), sw(220), sh(108), [
            ('Explainable /', 14, True, '0078d4'),
            ('Methodology Engine', 14, True, '0078d4'),
            ('', 4, False, None),
            ('SHAP + LIME frameworks deliver', 10, False, '555555'),
            ('transparent, auditable reasoning', 10, False, '555555'),
        ], align='center')

        # ── Vendor stat ovals ──
        # 62% Traditional
        add_box(sx(65), sy(425), sw(170), sh(75), 'fff5f0', 'ca5010')
        add_multiline_textbox(sx(70), sy(430), sw(160), sh(68), [
            ('62%', 20, True, 'ca5010'),
            ('Traditional \u2022 86 vendors', 10, False, '5a503c'),
        ], align='center')

        # 23% AI-First Startups
        add_box(sx(325), sy(440), sw(150), sh(75), 'f0f8ff', '0078d4')
        add_multiline_textbox(sx(330), sy(445), sw(140), sh(68), [
            ('23%', 20, True, '0078d4'),
            ('AI-First Startups \u2022 32', 10, False, '1a3a5c'),
        ], align='center')

        # 11% AI-First Other
        add_box(sx(80), sy(530), sw(140), sh(65), 'f8f0ff', '8764b8')
        add_multiline_textbox(sx(85), sy(535), sw(130), sh(58), [
            ('11%', 20, True, '8764b8'),
            ('AI-First Other \u2022 15', 10, False, '5a3c7c'),
        ], align='center')

        # 138 Vendors
        add_box(sx(960), sy(435), sw(180), sh(80), 'f0fff4', '107c10')
        add_multiline_textbox(sx(965), sy(440), sw(170), sh(72), [
            ('138', 28, True, '107c10'),
            ('Vendors Analyzed', 12, True, '1a3a5c'),
        ], align='center')

        # 80%+ ForensicLLM
        add_box(sx(765), sy(495), sw(170), sh(70), 'f0faff', '0078d4')
        add_multiline_textbox(sx(770), sy(500), sw(160), sh(62), [
            ('80%+', 20, True, '0078d4'),
            ('ForensicLLM Accuracy', 10, False, '1a3a5c'),
        ], align='center')

        # ── Bottom comparison bars ──
        add_box(sx(20), sy(730), sw(1160), sh(150), 'f8f8f5', 'e0ddd5')

        # Row labels
        add_textbox(sx(50), sy(738), sw(200), sh(20), 'TRADITIONAL', 11, True, 'ca5010')
        add_textbox(sx(550), sy(738), sw(100), sh(20), 'vs.', 12, True, '888888', 'center')
        add_textbox(sx(950), sy(738), sw(200), sh(20), 'AI-POWERED', 11, True, '0078d4', 'right')

        bars = [
            (775, 'Analysis Speed: Hours to Days', 'Analysis Speed: Under 12 Minutes'),
            (810, 'Methodology: Human-Driven / Opaque', 'Methodology: Explainable / Engine'),
            (845, 'Core Focus: Reporting & Admin', 'Core Focus: Triage, Containment & Forensics'),
        ]
        for y, left_txt, right_txt in bars:
            add_box(sx(50), sy(y), sw(460), sh(28), 'fce8d8', 'ca5010')
            add_textbox(sx(55), sy(y + 3), sw(450), sh(24), left_txt, 11, False, 'ca5010', 'center')
            add_textbox(sx(520), sy(y + 1), sw(160), sh(24), '\u2192', 16, True, '107c10', 'center')
            add_box(sx(690), sy(y), sw(460), sh(28), 'd8eeff', '0078d4')
            add_textbox(sx(695), sy(y + 3), sw(450), sh(24), right_txt, 11, False, '0078d4', 'center')

        # Footer
        add_textbox(sx(800), sy(875), sw(380), sh(20),
                    '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # Write to bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from flask import send_file
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='DFIR_Agentic_Shift_Infographic.pptx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/dfir-all-graphics-pptx', methods=['GET'])
def dfir_all_graphics_pptx():
    """Generate a multi-slide editable PowerPoint deck with all 13 DFIR graphics (6 Story Arc + 7 Data Deep Dives)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        def box(slide, l, t, w, h, fill=None, border=None):
            s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
            if fill:
                s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
            else:
                s.fill.background()
            if border:
                s.line.color.rgb = rgb(border); s.line.width = Pt(1.5)
            else:
                s.line.fill.background()
            return s

        def tb(slide, l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = text
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = rgb(col)
            if al == 'center': p.alignment = PP_ALIGN.CENTER
            elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        def ml(slide, l, t, w, h, lines, al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(11)
                p.font.bold = bld
                if col: p.font.color.rgb = rgb(col)
                if al == 'center': p.alignment = PP_ALIGN.CENTER
                elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        SW = Inches(13.333); SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 750 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 750 * SH)

        # ════════════════════════════════════════════════════════════
        #  STORY ARC SLIDES (1–6): Complete Research Narrative
        # ════════════════════════════════════════════════════════════

        # ────────────────────────────────────────────────────────────
        # STORY 1: Market Landscape Pie Chart
        # ────────────────────────────────────────────────────────────
        sa1 = prs.slides.add_slide(BLANK)
        box(sa1, sx(0), sy(0), sw(1200), sh(60), '0078d4', '0078d4')
        tb(sa1, sx(50), sy(8), sw(1100), sh(28), 'Story 1: The DFIR Market Landscape', 22, True, 'ffffff', 'center')
        tb(sa1, sx(100), sy(36), sw(1000), sh(22), '138 vendors analyzed across three distinct investment profiles', 11, False, 'ccd8e8', 'center')

        # Pie segments as proportional boxes with labels
        pie_data = [
            ('Traditional', 86, '62%', 'ca5010', 'fce8d8'),
            ('AI-First Startups', 32, '23%', '0078d4', 'd8eeff'),
            ('AI-First Non-Startups', 15, '11%', '8764b8', 'f0e8ff'),
            ('Other', 5, '4%', '888888', 'f0f0f0'),
        ]
        # Large center number
        tb(sa1, sx(400), sy(100), sw(400), sh(100), '138', 60, True, '0078d4', 'center')
        tb(sa1, sx(400), sy(190), sw(400), sh(30), 'Vendors Analyzed', 16, True, '333333', 'center')

        # Pie segments as horizontal proportional bar
        bar_x = 100
        bar_y = 260
        bar_w = 1000
        bar_h = 70
        colors = ['ca5010', '0078d4', '8764b8', '888888']
        pcts = [62, 23, 11, 4]
        cx = bar_x
        for i, (label, count, pct_str, col, bg) in enumerate(pie_data):
            seg_w = int(bar_w * pcts[i] / 100)
            box(sa1, sx(cx), sy(bar_y), sw(seg_w), sh(bar_h), col, col)
            if seg_w > 60:
                tb(sa1, sx(cx + 5), sy(bar_y + 10), sw(seg_w - 10), sh(24), pct_str, 22, True, 'ffffff', 'center')
                tb(sa1, sx(cx + 5), sy(bar_y + 38), sw(seg_w - 10), sh(20), label, 11, False, 'ffffff', 'center')
            cx += seg_w

        # Category detail cards
        card_y = 370
        card_h = 130
        for i, (label, count, pct_str, col, bg) in enumerate(pie_data):
            cx_card = 30 + i * 290
            box(sa1, sx(cx_card), sy(card_y), sw(270), sh(card_h), bg, col)
            tb(sa1, sx(cx_card + 15), sy(card_y + 10), sw(240), sh(30), label, 14, True, col)
            tb(sa1, sx(cx_card + 15), sy(card_y + 42), sw(240), sh(30), f'{count} vendors ({pct_str})', 20, True, col, 'center')
            insights = [
                'Avg 3.92 across 5 pillars\nLegacy-integrated architecture',
                'Avg 4.11 across 5 pillars\n+0.60 on investigative path',
                'Avg 3.96 across 5 pillars\nMixed architecture profile',
                'Niche specialists\nLimited pillar coverage',
            ]
            tb(sa1, sx(cx_card + 15), sy(card_y + 78), sw(240), sh(44), insights[i], 10, False, '555555', 'center')

        # Key insight callout
        box(sa1, sx(50), sy(530), sw(1100), sh(80), 'e8f4fd', '0078d4')
        tb(sa1, sx(70), sy(538), sw(1060), sh(18), 'Key Insight', 13, True, '0078d4')
        tb(sa1, sx(70), sy(560), sw(1060), sh(44), 'Traditional vendors (62%) average 3.92 across five capability pillars. AI-first startups (23%) average 4.11, outperforming significantly on the investigative critical path. Neither category scores above 4.0 in every pillar.', 11, False, '333333')

        tb(sa1, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # STORY 2: Trust Barrier Mindmap
        # ────────────────────────────────────────────────────────────
        sa2 = prs.slides.add_slide(BLANK)
        box(sa2, sx(0), sy(0), sw(1200), sh(60), 'ca5010', 'ca5010')
        tb(sa2, sx(50), sy(8), sw(1100), sh(28), 'Story 2: The Trust Barrier \u2014 Why AI Adoption Stalls', 22, True, 'ffffff', 'center')
        tb(sa2, sx(100), sy(36), sw(1000), sh(22), 'Trust \u2014 not capability \u2014 is the primary barrier to AI adoption in DFIR', 11, False, 'fce8d8', 'center')

        # Central node
        box(sa2, sx(480), sy(85), sw(240), sh(55), 'a80000', 'a80000')
        tb(sa2, sx(480), sy(93), sw(240), sh(40), 'TRUST DEFICIT', 18, True, 'ffffff', 'center')

        # Branch 1: Black Box Problem
        box(sa2, sx(30), sy(170), sw(350), sh(36), '3a1e0e', 'ca5010')
        tb(sa2, sx(35), sy(174), sw(340), sh(28), 'The Black Box Problem', 14, True, 'f0a070', 'center')
        bb_items = [
            'No visibility into AI methodology',
            'Practitioners distrust opaque models',
            'Cultural tradition favors manual work',
            'AI focused on admin not investigation',
        ]
        for j, item in enumerate(bb_items):
            yy = 216 + j * 32
            box(sa2, sx(40), sy(yy), sw(330), sh(28), 'fce8d8', 'e0c8b0')
            tb(sa2, sx(50), sy(yy + 4), sw(310), sh(20), item, 11, False, '5a3020')

        # Branch 2: Legal & Evidentiary Barriers
        box(sa2, sx(420), sy(170), sw(360), sh(36), '3a1e0e', 'ca5010')
        tb(sa2, sx(425), sy(174), sw(350), sh(28), 'Legal & Evidentiary Barriers', 14, True, 'f0a070', 'center')
        le_items = [
            'Daubert Standard for admissibility',
            'Federal Rule 901 requirements',
            'Chain of Custody must be documented',
            'Court demands repeatable process',
        ]
        for j, item in enumerate(le_items):
            yy = 216 + j * 32
            box(sa2, sx(430), sy(yy), sw(340), sh(28), 'fce8d8', 'e0c8b0')
            tb(sa2, sx(440), sy(yy + 4), sw(320), sh(20), item, 11, False, '5a3020')

        # Branch 3: Solution Pathways (right column)
        box(sa2, sx(820), sy(170), sw(360), sh(36), '0e2a0e', '107c10')
        tb(sa2, sx(825), sy(174), sw(350), sh(28), 'Solution Pathways', 14, True, '70f080', 'center')
        solutions = [
            ('SHAP & LIME Frameworks', 'Explain every AI decision'),
            ('Digital Forensic Knowledge Graphs', 'Visualize evidence linkages'),
            ('Tandem Operating Model', 'AI speed + human validation'),
            ('Immutable Audit Trails', 'UIDs for every artifact & action'),
        ]
        for j, (title, desc) in enumerate(solutions):
            yy = 216 + j * 64
            box(sa2, sx(830), sy(yy), sw(340), sh(56), 'e8f8e8', '107c10')
            tb(sa2, sx(840), sy(yy + 4), sw(320), sh(20), title, 12, True, '107c10')
            tb(sa2, sx(840), sy(yy + 26), sw(320), sh(24), desc, 10, False, '555555')

        # Bottom line callout
        box(sa2, sx(50), sy(500), sw(1100), sh(70), 'fef3e8', 'ca5010')
        tb(sa2, sx(70), sy(508), sw(1060), sh(18), 'Bottom Line', 13, True, 'ca5010')
        tb(sa2, sx(70), sy(530), sw(1060), sh(34), 'The largest obstacle is not AI capability but trust, driven by "black box" concerns. SHAP/LIME frameworks already satisfy chain of custody, Daubert, and Federal Rule 901 requirements \u2014 the barrier is adoption, not technology.', 11, False, '333333')

        tb(sa2, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # STORY 3: Detection → Methodology Engine Flowchart
        # ────────────────────────────────────────────────────────────
        sa3 = prs.slides.add_slide(BLANK)
        box(sa3, sx(0), sy(0), sw(1200), sh(60), '107c10', '107c10')
        tb(sa3, sx(50), sy(8), sw(1100), sh(28), 'Story 3: From Detection Tool to Methodology Engine', 22, True, 'ffffff', 'center')
        tb(sa3, sx(100), sy(36), sw(1000), sh(22), 'AI must evolve from automating admin tasks to powering the investigative critical path', 11, False, 'd4edce', 'center')

        # LEFT column: Where AI Is Today (orange)
        box(sa3, sx(40), sy(90), sw(440), sh(36), '3a1e0e', 'ca5010')
        tb(sa3, sx(45), sy(94), sw(430), sh(28), 'WHERE AI IS TODAY', 15, True, 'f0a070', 'center')
        today_items = ['Report Writing', 'SOAR Automation', 'Program Management', 'Compliance Docs', 'Planning Support']
        for j, item in enumerate(today_items):
            yy = 140 + j * 44
            box(sa3, sx(60), sy(yy), sw(400), sh(38), '3a1e0e', 'ca5010')
            tb(sa3, sx(70), sy(yy + 8), sw(380), sh(22), item, 13, False, 'f0a070', 'center')

        # RIGHT column: Where AI Should Focus (green)
        box(sa3, sx(720), sy(90), sw(440), sh(36), '0e2a0e', '107c10')
        tb(sa3, sx(725), sy(94), sw(430), sh(28), 'WHERE AI SHOULD FOCUS', 15, True, '70f080', 'center')
        focus_items = ['Deep Forensic Analysis', 'Triage & Scoping', 'Timeline Reconstruction', 'Malware Reverse Engineering', 'Containment & Isolation']
        for j, item in enumerate(focus_items):
            yy = 140 + j * 44
            box(sa3, sx(740), sy(yy), sw(400), sh(38), '0e2a0e', '107c10')
            tb(sa3, sx(750), sy(yy + 8), sw(380), sh(22), item, 13, False, '70f080', 'center')

        # Center arrows
        tb(sa3, sx(480), sy(180), sw(240), sh(24), '\u2192  Marginal Impact', 12, True, 'a80000', 'center')
        box(sa3, sx(500), sy(210), sw(200), sh(36), '3a0e0e', 'a80000')
        tb(sa3, sx(505), sy(216), sw(190), sh(24), 'Process Efficiency Only', 11, True, 'ff8888', 'center')

        tb(sa3, sx(480), sy(300), sw(240), sh(24), '\u2192  Critical Path Impact', 12, True, '107c10', 'center')
        box(sa3, sx(480), sy(330), sw(240), sh(36), '0e3a0e', '107c10')
        tb(sa3, sx(485), sy(336), sw(230), sh(24), 'Faster Recovery &\nRisk Mitigation', 11, True, '88ff88', 'center')

        # The Shift callout
        box(sa3, sx(50), sy(420), sw(1100), sh(70), 'e8f8e8', '107c10')
        tb(sa3, sx(70), sy(428), sw(1060), sh(18), 'The Shift', 13, True, '107c10')
        tb(sa3, sx(70), sy(450), sw(1060), sh(34), 'ForensicLLM achieves 80%+ accuracy in source attribution. DFKGs shift focus from "what happened" to "how evidence was collected." AI-first startups scoring +0.60 above traditional vendors in triage, containment, and malware analysis prove the model works.', 11, False, '333333')

        tb(sa3, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # STORY 4: Tandem Operating Model Sequence
        # ────────────────────────────────────────────────────────────
        sa4 = prs.slides.add_slide(BLANK)
        box(sa4, sx(0), sy(0), sw(1200), sh(60), '8764b8', '8764b8')
        tb(sa4, sx(50), sy(8), sw(1100), sh(28), 'Story 4: The Tandem Operating Model', 22, True, 'ffffff', 'center')
        tb(sa4, sx(100), sy(36), sw(1000), sh(22), 'AI speed and human judgment working in concert, from raw evidence to court-admissible deliverables', 11, False, 'e8d8f8', 'center')

        # Sequence participants (horizontal lane headers)
        participants = [
            ('Evidence', '0078d4', 'd8eeff'),
            ('AI Agent', '107c10', 'e8f8e8'),
            ('DFKG', 'ca5010', 'fce8d8'),
            ('Analyst', '8764b8', 'f0e8ff'),
            ('Output', '1a3a5c', 'e0e8f0'),
        ]
        lane_w = 200
        lane_gap = 20
        lane_start = 50
        for i, (name, col, bg) in enumerate(participants):
            lx = lane_start + i * (lane_w + lane_gap)
            box(sa4, sx(lx), sy(80), sw(lane_w), sh(40), col, col)
            tb(sa4, sx(lx), sy(86), sw(lane_w), sh(28), name, 14, True, 'ffffff', 'center')

        # Sequence steps (horizontal flow with from→to indicators)
        seq_steps = [
            (0, 1, 'Raw telemetry & artifacts'),
            (1, 1, 'Data ingestion & correlation'),
            (1, 2, 'Build knowledge graph'),
            (2, 1, 'Linkages & timeline'),
            (1, 3, 'Findings + methodology doc'),
            (None, None, 'SHAP/LIME explains each decision'),
            (3, 1, 'Methodology validated \u2713'),
            (1, 1, 'Deep forensic analysis'),
            (1, 3, 'Results + evidence lineage'),
            (3, 3, 'Strategic interpretation'),
            (3, 4, 'Court-admissible report'),
            (None, None, 'Every agent action logged with UIDs'),
        ]
        step_y = 140
        for j, (frm, to, label) in enumerate(seq_steps):
            yy = step_y + j * 40
            if frm is None:
                # Note spanning full width
                box(sa4, sx(lane_start), sy(yy), sw(1100), sh(34), 'f8f0ff', '8764b8')
                tb(sa4, sx(lane_start + 10), sy(yy + 6), sw(1080), sh(22), '\u2605 ' + label, 11, True, '8764b8', 'center')
            else:
                frm_x = lane_start + frm * (lane_w + lane_gap) + lane_w // 2
                to_x = lane_start + to * (lane_w + lane_gap) + lane_w // 2
                left_x = min(frm_x, to_x) - 20
                span_w = abs(to_x - frm_x) + 40 if frm != to else lane_w
                if frm == to:
                    left_x = lane_start + frm * (lane_w + lane_gap)
                    span_w = lane_w
                arrow_col = participants[to][1] if frm != to else participants[frm][1]
                box(sa4, sx(left_x), sy(yy), sw(span_w), sh(34), participants[to][2] if frm != to else participants[frm][2], arrow_col)
                direction = '\u2192' if to >= frm else '\u2190'
                display_label = (direction + ' ' + label) if frm != to else label
                tb(sa4, sx(left_x + 5), sy(yy + 6), sw(span_w - 10), sh(22), display_label, 10, False, '333333', 'center')

        # Operating principle callout
        box(sa4, sx(50), sy(630), sw(1100), sh(70), 'f0e8ff', '8764b8')
        tb(sa4, sx(70), sy(638), sw(1060), sh(18), 'Operating Principle', 13, True, '8764b8')
        tb(sa4, sx(70), sy(660), sw(1060), sh(34), 'AI agents handle data ingestion, correlation, and timeline reconstruction. Human experts validate methodology and provide strategic interpretation. Every agent action \u2014 model versions, prompts, tool invocations \u2014 is logged in an immutable audit trail.', 11, False, '333333')

        tb(sa4, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # STORY 5: Evidence Chain of Custody State Diagram
        # ────────────────────────────────────────────────────────────
        sa5 = prs.slides.add_slide(BLANK)
        box(sa5, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(sa5, sx(50), sy(8), sw(1100), sh(28), 'Story 5: Evidence Chain of Custody Lifecycle', 22, True, 'ffffff', 'center')
        tb(sa5, sx(100), sy(36), sw(1000), sh(22), 'How AI maintains defensible chain of custody from raw evidence to court-ready deliverables', 11, False, 'ccd8e8', 'center')

        # State flow: left to right with 3 composite states
        states = [
            ('AI Ingestion', '0078d4', 'd8eeff', ['Data Capture', 'Hash Verification', 'UID Assignment']),
            ('AI Analysis', '107c10', 'e8f8e8', ['Correlation', 'Timeline Reconstruction', 'Artifact Linkage']),
            ('Human Validation', '8764b8', 'f0e8ff', ['Methodology Review', 'Reasoning Verification', 'Evidence Confirmation']),
        ]
        # Start node
        box(sa5, sx(40), sy(110), sw(140), sh(50), '555555', '333333')
        tb(sa5, sx(40), sy(116), sw(140), sh(38), 'Incident\nDetected', 12, True, 'ffffff', 'center')

        # Raw Evidence
        tb(sa5, sx(185), sy(122), sw(40), sh(24), '\u2192', 18, True, '555555', 'center')
        box(sa5, sx(225), sy(110), sw(140), sh(50), 'fff3e0', 'ca5010')
        tb(sa5, sx(225), sy(116), sw(140), sh(38), 'Raw\nEvidence', 12, True, 'ca5010', 'center')

        # Three composite state boxes
        state_x = 30
        state_y = 200
        sbox_w = 350
        sbox_h = 220
        for si, (sname, scol, sbg, sub_items) in enumerate(states):
            sx_pos = state_x + si * (sbox_w + 55)
            # Outer state box
            box(sa5, sx(sx_pos), sy(state_y), sw(sbox_w), sh(sbox_h), sbg, scol)
            # State label
            tb(sa5, sx(sx_pos + 10), sy(state_y + 8), sw(sbox_w - 20), sh(24), sname, 15, True, scol, 'center')
            # Sub-states
            for sj, sub in enumerate(sub_items):
                sub_y = state_y + 44 + sj * 52
                box(sa5, sx(sx_pos + 20), sy(sub_y), sw(sbox_w - 40), sh(42), 'ffffff', scol)
                tb(sa5, sx(sx_pos + 30), sy(sub_y + 10), sw(sbox_w - 60), sh(22), sub, 12, False, '333333', 'center')
                if sj < len(sub_items) - 1:
                    tb(sa5, sx(sx_pos + 150), sy(sub_y + 38), sw(50), sh(18), '\u2193', 14, True, scol, 'center')

            # Transition arrows between states
            if si < 2:
                arrow_x = sx_pos + sbox_w + 5
                transitions = ['AI Agent Acquires', 'Documented Methods', 'SHAP/LIME Explanations']
                tb(sa5, sx(arrow_x), sy(state_y + sbox_h // 2 - 20), sw(45), sh(24), '\u2192', 18, True, scol, 'center')
                tb(sa5, sx(arrow_x - 10), sy(state_y + sbox_h // 2 + 4), sw(65), sh(30), transitions[si], 8, False, '777777', 'center')

        # Court Ready + end
        box(sa5, sx(440), sy(450), sw(200), sh(50), '107c10', '107c10')
        tb(sa5, sx(440), sy(456), sw(200), sh(38), 'Court Ready', 14, True, 'ffffff', 'center')
        tb(sa5, sx(440), sy(500), sw(200), sh(24), 'Immutable Audit Trail \u2192', 9, False, '777777', 'center')

        # Daubert/Rule 901 badge
        box(sa5, sx(690), sy(450), sw(260), sh(50), 'e8f8e8', '107c10')
        tb(sa5, sx(700), sy(456), sw(240), sh(38), 'Daubert & Rule 901\nCompliant', 13, True, '107c10', 'center')

        tb(sa5, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # STORY 6: Path to 2030 Timeline
        # ────────────────────────────────────────────────────────────
        sa6 = prs.slides.add_slide(BLANK)
        box(sa6, sx(0), sy(0), sw(1200), sh(60), 'a80000', 'a80000')
        tb(sa6, sx(50), sy(8), sw(1100), sh(28), 'Story 6: Path to 2030 \u2014 The DFIR AI Convergence', 22, True, 'ffffff', 'center')
        tb(sa6, sx(100), sy(36), sw(1000), sh(22), '"By 2030, traditional models of manual, human-dependent forensic investigation will largely be irrelevant."', 11, False, 'ffc8c8', 'center')

        timeline_data = [
            ('2024', 'a80000', 'ffe0e0', [
                'AI applied to reports and admin only',
                'Traditional vendors dominate at 62%',
                'Trust deficit blocks investigation AI',
            ]),
            ('2025', 'ca5010', 'fce8d8', [
                'ForensicLLM achieves 80%+ accuracy',
                'DFKGs enable evidence visualization',
                'AI-first startups gain investigation edge',
            ]),
            ('2026', '0078d4', 'd8eeff', [
                'Tandem operating models emerge',
                'SHAP & LIME frameworks standardize',
                '138 vendors assessed across 5 pillars',
            ]),
            ('2027', '8764b8', 'f0e8ff', [
                'Govt oversight of AI forensic tools',
                'Insurance validation of AI-powered LLMs',
                'Traditional vendors redirect investment',
            ]),
            ('2028', '107c10', 'e8f8e8', [
                'Vendors below 4.0 in investigation at risk',
                'AI-native architectures become standard',
                'Chain of custody automation matures',
            ]),
            ('2030', '1a3a5c', 'e0e8f0', [
                'Manual investigation models irrelevant',
                'Full convergence: AI speed + human trust',
                'Court admissibility of AI evidence standard',
            ]),
        ]
        tl_y = 90
        col_w = 185
        col_gap = 10
        for ti, (year, col, bg, items) in enumerate(timeline_data):
            tx = 20 + ti * (col_w + col_gap)
            # Year header
            box(sa6, sx(tx), sy(tl_y), sw(col_w), sh(44), col, col)
            tb(sa6, sx(tx), sy(tl_y + 6), sw(col_w), sh(28), year, 22, True, 'ffffff', 'center')
            # Connector arrow
            if ti < len(timeline_data) - 1:
                tb(sa6, sx(tx + col_w), sy(tl_y + 8), sw(col_gap + 4), sh(28), '\u2192', 16, True, '888888', 'center')
            # Items
            for tj, item in enumerate(items):
                iy = tl_y + 54 + tj * 50
                box(sa6, sx(tx), sy(iy), sw(col_w), sh(44), bg, col)
                tb(sa6, sx(tx + 8), sy(iy + 6), sw(col_w - 16), sh(32), item, 10, False, '333333', 'center')

        # Critical threshold callout
        box(sa6, sx(50), sy(470), sw(1100), sh(80), 'ffe0e0', 'a80000')
        tb(sa6, sx(70), sy(478), sw(1060), sh(18), 'Critical Threshold \u2014 2028', 13, True, 'a80000')
        tb(sa6, sx(70), sy(500), sw(1060), sh(44), 'Vendors scoring below 4.0 in investigative and remediation pillars by 2028 risk being unable to compete for enterprise DFIR engagements by 2030. Today, 100% of traditional vendors are legacy-integrated; only 50% of AI-first startups have moved to AI-native architecture.', 11, False, '333333')

        tb(sa6, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ════════════════════════════════════════════════════════════
        #  DATA DEEP DIVE SLIDES (7–13): Evidence & Analysis
        # ════════════════════════════════════════════════════════════

        # ────────────────────────────────────────────────────────────
        # SLIDE 7: Five-Pillar Capability Heatmap
        # ────────────────────────────────────────────────────────────
        s1 = prs.slides.add_slide(BLANK)
        box(s1, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s1, sx(50), sy(8), sw(1100), sh(28), '7. Five-Pillar Capability Heatmap', 22, True, 'ffffff', 'center')
        tb(s1, sx(100), sy(36), sw(1000), sh(22), '138 DFIR vendors scored across five capability pillars by vendor category', 11, False, 'ccd8e8', 'center')

        heatmap = [
            ('Investigation', 4.02, 4.55, 4.52),
            ('Planning & Preparation', 3.94, 4.26, 4.12),
            ('Containment & Remediation', 3.72, 4.00, 3.93),
            ('Program Management', 3.97, 3.91, 3.82),
            ('Legal & Compliance', 3.95, 3.77, 3.88),
        ]
        avg = ('Overall Average', 3.92, 4.11, 4.04)

        # Column headers
        colW = 220; startX = 380; rowH = 55
        headers = [('Traditional\n86 vendors (62%)', '#ca5010'), ('AI-First Startups\n32 vendors (23%)', '#0078d4'), ('AI-First Non-Startups\n15 vendors (11%)', '#8764b8')]
        for ci, (ht, hc) in enumerate(headers):
            box(s1, sx(startX + ci * colW), sy(75), sw(colW - 10), sh(50), hc, hc)
            tb(s1, sx(startX + ci * colW + 5), sy(78), sw(colW - 20), sh(46), ht, 11, True, 'ffffff', 'center')

        def hm_color(v):
            if v >= 4.3: return '107c10'
            if v >= 4.0: return '0078d4'
            if v >= 3.85: return 'ca5010'
            return 'a80000'
        def hm_bg(v):
            if v >= 4.3: return 'e6f4e6'
            if v >= 4.0: return 'e0f0ff'
            if v >= 3.85: return 'fff0e0'
            return 'ffe0e0'

        for ri, (name, t, ai, non) in enumerate(heatmap):
            y = 135 + ri * rowH
            box(s1, sx(40), sy(y), sw(330), sh(rowH - 5), 'f8f8f5', 'e0ddd5')
            tb(s1, sx(50), sy(y + 8), sw(310), sh(36), name, 13, True, '333333')
            for ci, v in enumerate([t, ai, non]):
                box(s1, sx(startX + ci * colW), sy(y), sw(colW - 10), sh(rowH - 5), hm_bg(v), hm_color(v))
                tb(s1, sx(startX + ci * colW + 5), sy(y + 8), sw(colW - 20), sh(36), f'{v:.2f}', 20, True, hm_color(v), 'center')

        # Average row
        y = 135 + 5 * rowH + 10
        box(s1, sx(40), sy(y), sw(330), sh(rowH - 5), 'e0ddd5', 'b0a898')
        tb(s1, sx(50), sy(y + 8), sw(310), sh(36), avg[0], 13, True, '333333')
        for ci, v in enumerate([avg[1], avg[2], avg[3]]):
            box(s1, sx(startX + ci * colW), sy(y), sw(colW - 10), sh(rowH - 5), hm_bg(v), hm_color(v))
            tb(s1, sx(startX + ci * colW + 5), sy(y + 8), sw(colW - 20), sh(36), f'{v:.2f}', 20, True, hm_color(v), 'center')

        # Legend
        legend_y = y + rowH + 20
        legends = [('\u25a0 \u2265 4.30 Strong', '107c10'), ('\u25a0 4.00\u20134.29 Competitive', '0078d4'), ('\u25a0 3.85\u20133.99 Below Target', 'ca5010'), ('\u25a0 < 3.85 Critical Gap', 'a80000')]
        for li, (lt, lc) in enumerate(legends):
            tb(s1, sx(40 + li * 270), sy(legend_y), sw(260), sh(22), lt, 11, True, lc)

        tb(s1, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # SLIDE 8: 2030 Convergence Roadmap
        # ────────────────────────────────────────────────────────────
        s2 = prs.slides.add_slide(BLANK)
        box(s2, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s2, sx(50), sy(8), sw(1100), sh(28), '8. Path to 2030: Vendor Category Convergence Roadmap', 22, True, 'ffffff', 'center')
        tb(s2, sx(100), sy(36), sw(1000), sh(22), 'Current maturity as % of 5.0 target, with capability gap tags per category', 11, False, 'ccd8e8', 'center')

        categories = [
            ('Traditional', '86 vendors', 3.92, 78.4, '#ca5010',
             [('\u26a0 Containment 3.72', True), ('\u26a0 Investigation needs AI lift', True), ('100% legacy-integrated', False), ('\u2713 Program Mgmt 3.97', False), ('\u2713 Legal 3.95', False)]),
            ('AI-First Startups', '32 vendors', 4.11, 82.2, '#0078d4',
             [('\u26a0 Legal/Compliance 3.77', True), ('\u26a0 Chain of Custody \u22120.37', True), ('50% AI-native architecture', False), ('\u2713 Investigation 4.55', False), ('\u2713 Planning 4.26', False)]),
            ('AI-First Non-Startups', '15 vendors', 4.04, 80.8, '#8764b8',
             [('\u26a0 Program Mgmt 3.82', True), ('\u26a0 Containment 3.93', True), ('\u2713 Investigation 4.52', False)]),
        ]
        for ci, (name, sub, score, pct, color, gaps) in enumerate(categories):
            baseY = 80 + ci * 200
            # Category label
            tb(s2, sx(40), sy(baseY), sw(200), sh(24), name, 16, True, color.lstrip('#'))
            tb(s2, sx(40), sy(baseY + 26), sw(200), sh(18), f'{sub} \u2022 {score:.2f} avg', 10, False, '666666')
            # Progress bar background
            box(s2, sx(260), sy(baseY + 5), sw(880), sh(38), 'e8e8e5', 'cccccc')
            # Progress bar fill
            fillW = int(880 * pct / 100)
            box(s2, sx(260), sy(baseY + 5), sw(fillW), sh(38), color.lstrip('#'), color.lstrip('#'))
            tb(s2, sx(260 + fillW // 2 - 60), sy(baseY + 10), sw(120), sh(28), f'{score:.2f} / 5.00', 14, True, 'ffffff', 'center')
            # Target marker
            targetX = int(880 * 80 / 100)
            box(s2, sx(260 + targetX - 1), sy(baseY + 2), sw(3), sh(44), 'a80000', 'a80000')
            tb(s2, sx(260 + targetX - 35), sy(baseY + 48), sw(70), sh(16), '4.00 Target', 8, True, 'a80000', 'center')
            # Gap tags
            for gi, (gtxt, critical) in enumerate(gaps):
                gx = 260 + (gi % 3) * 300
                gy = baseY + 72 + (gi // 3) * 30
                gc = 'a80000' if critical else '107c10'
                gb = 'fff0f0' if critical else 'f0fff0'
                box(s2, sx(gx), sy(gy), sw(280), sh(24), gb, gc)
                tb(s2, sx(gx + 5), sy(gy + 2), sw(270), sh(20), gtxt, 10, False, gc, 'center')

        # SPA callout
        box(s2, sx(40), sy(660), sw(1120), sh(70), 'fff5f5', 'a80000')
        ml(s2, sx(50), sy(663), sw(1100), sh(64), [
            ('SPA \u2014 2030 Strategic Planning Assumption', 12, True, 'a80000'),
            ('"By 2030, traditional models of manual forensic investigation will largely be irrelevant." Vendors below 4.0 by 2028 risk being unable to compete.', 11, False, '555555'),
        ])
        tb(s2, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # SLIDE 9: Investment Mismatch
        # ────────────────────────────────────────────────────────────
        s3 = prs.slides.add_slide(BLANK)
        box(s3, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s3, sx(50), sy(8), sw(1100), sh(28), '9. Investment Mismatch: Where the Industry Is vs. Where It Should Be', 22, True, 'ffffff', 'center')
        tb(s3, sx(100), sy(36), sw(1000), sh(22), 'Traditional vs AI-First investment focus compared to recommended allocation', 11, False, 'ccd8e8', 'center')

        mismatch = [
            ('Investigation', 40, 85, 90),
            ('Containment', 30, 70, 85),
            ('Planning', 50, 75, 70),
            ('Program Mgmt', 80, 45, 40),
            ('Legal/Compliance', 75, 40, 35),
        ]
        # Legend
        tb(s3, sx(60), sy(72), sw(160), sh(18), '\u25a0 Traditional Focus', 11, True, 'ca5010')
        tb(s3, sx(240), sy(72), sw(140), sh(18), '\u25a0 AI-First Focus', 11, True, '0078d4')
        tb(s3, sx(400), sy(72), sw(180), sh(18), '\u25a0 Recommended Focus', 11, True, '107c10')

        for mi, (name, trad, ai, ideal) in enumerate(mismatch):
            baseY = 100 + mi * 115
            tb(s3, sx(40), sy(baseY), sw(250), sh(22), name, 14, True, '333333')
            barW = 700
            labels_w = 80
            # Traditional bar
            tb(s3, sx(40), sy(baseY + 28), sw(labels_w), sh(18), 'Traditional', 10, False, '777777')
            box(s3, sx(40 + labels_w), sy(baseY + 26), sw(barW), sh(22), 'f0f0ee', 'dddddd')
            box(s3, sx(40 + labels_w), sy(baseY + 26), sw(int(barW * trad / 100)), sh(22), 'ca5010', 'ca5010')
            tb(s3, sx(40 + labels_w + barW + 10), sy(baseY + 28), sw(50), sh(18), f'{trad}%', 12, True, 'ca5010')
            # AI-First bar
            tb(s3, sx(40), sy(baseY + 54), sw(labels_w), sh(18), 'AI-First', 10, False, '777777')
            box(s3, sx(40 + labels_w), sy(baseY + 52), sw(barW), sh(22), 'f0f0ee', 'dddddd')
            box(s3, sx(40 + labels_w), sy(baseY + 52), sw(int(barW * ai / 100)), sh(22), '0078d4', '0078d4')
            tb(s3, sx(40 + labels_w + barW + 10), sy(baseY + 54), sw(50), sh(18), f'{ai}%', 12, True, '0078d4')
            # Recommended bar
            tb(s3, sx(40), sy(baseY + 80), sw(labels_w), sh(18), 'Recommended', 10, False, '777777')
            box(s3, sx(40 + labels_w), sy(baseY + 78), sw(barW), sh(22), 'f0f0ee', 'dddddd')
            box(s3, sx(40 + labels_w), sy(baseY + 78), sw(int(barW * ideal / 100)), sh(22), '107c10', '107c10')
            tb(s3, sx(40 + labels_w + barW + 10), sy(baseY + 80), sw(50), sh(18), f'{ideal}%', 12, True, '107c10')

        # Key insight
        box(s3, sx(40), sy(680), sw(1120), sh(50), 'f0faff', '0078d4')
        ml(s3, sx(50), sy(683), sw(1100), sh(44), [
            ('Key Insight', 11, True, '0078d4'),
            ('Traditional vendors over-invest in governance (+40 above recommended in Pgm Mgmt) while under-investing in investigation (\u221250 below recommended). AI-first are closer but under-invest in legal/compliance.', 10, False, '555555'),
        ])
        tb(s3, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # SLIDE 10: Executive Summary Infographic
        # ────────────────────────────────────────────────────────────
        s4 = prs.slides.add_slide(BLANK)
        box(s4, sx(0), sy(0), sw(1200), sh(55), '0078d4', '0078d4')
        tb(s4, sx(50), sy(5), sw(1100), sh(28), '10. Agentic AI Is the New Digital Forensics Workhorse', 20, True, 'ffffff', 'center')
        tb(s4, sx(100), sy(32), sw(1000), sh(20), 'Market Insight Executive Summary \u2014 138 vendors across 5 capability pillars', 11, False, 'b0d8ff', 'center')

        # Stats bar
        stats = [('138', 'Vendors'), ('62%', 'Traditional'), ('23%', 'AI-First Startups'), ('11%', 'Non-Startups'), ('80%+', 'ForensicLLM')]
        stat_colors = ['0078d4', 'ca5010', '107c10', '8764b8', '0078d4']
        for si, ((sv, sl), sc) in enumerate(zip(stats, stat_colors)):
            bx = 40 + si * 232
            box(s4, sx(bx), sy(65), sw(215), sh(60), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(bx + 5), sy(68), sw(205), sh(28), sv, 22, True, sc, 'center')
            tb(s4, sx(bx + 5), sy(98), sw(205), sh(20), sl, 10, False, '666666', 'center')

        # Key Findings
        tb(s4, sx(40), sy(140), sw(300), sh(22), 'KEY FINDINGS', 13, True, '0078d4')
        findings = [
            ('\U0001f517 Trust Through Transparency', 'SHAP + LIME explainability frameworks satisfy chain of custody, Daubert, and Federal Rule 901 admissibility.'),
            ('\U0001f9e0 AI-Driven Investigation Is Viable', 'ForensicLLM achieves 80%+ source attribution. DFKGs enable real-time AI reasoning verification.'),
            ('\u26a1 Offensive\u2013Defensive Gap Widening', 'Threat actors weaponize zero-days in hours. Defensive AI lags due to regulatory and cultural barriers.'),
        ]
        for fi, (ft, fb) in enumerate(findings):
            fy = 168 + fi * 85
            box(s4, sx(40), sy(fy), sw(1120), sh(75), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(55), sy(fy + 5), sw(1090), sh(22), ft, 13, True, '333333')
            tb(s4, sx(55), sy(fy + 30), sw(1090), sh(40), fb, 11, False, '555555')

        # Recommendations flow
        tb(s4, sx(40), sy(430), sw(400), sh(22), 'RECOMMENDATIONS FOR PRODUCT LEADERS', 13, True, 'ca5010')
        recs = [
            ('01', 'Explainable AI', 'Prioritize SHAP/LIME workflows', '0078d4'),
            ('02', 'Tandem Model', 'AI ingestion + human validation', '107c10'),
            ('03', 'Evidence Lineage', 'UIDs + location records in AI reports', 'ca5010'),
            ('04', 'Immutable Audit', 'Log every agent action & prompt', '8764b8'),
        ]
        for ri, (rn, rt, rd, rc) in enumerate(recs):
            rx = 40 + ri * 285
            box(s4, sx(rx), sy(460), sw(260), sh(90), 'f8f8f5', rc)
            tb(s4, sx(rx + 10), sy(463), sw(240), sh(18), rn, 10, True, rc)
            tb(s4, sx(rx + 10), sy(482), sw(240), sh(20), rt, 14, True, rc, 'center')
            tb(s4, sx(rx + 10), sy(505), sw(240), sh(40), rd, 10, False, '555555', 'center')
            if ri < 3:
                tb(s4, sx(rx + 265), sy(490), sw(20), sh(22), '\u2192', 16, True, '888888', 'center')

        # SPA callout
        box(s4, sx(40), sy(570), sw(1120), sh(55), 'fff0f0', 'a80000')
        ml(s4, sx(55), sy(575), sw(1090), sh(45), [
            ('\u26a0 "By 2030, the traditional models of manual, human-dependent forensic investigation will largely be irrelevant."', 12, True, 'a80000'),
            ('\u2014 Strategic Planning Assumption', 10, False, '888888'),
        ], 'center')
        tb(s4, sx(800), sy(720), sw(380), sh(20), '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ────────────────────────────────────────────────────────────
        # SLIDE 11: AI Arms Race
        # ────────────────────────────────────────────────────────────
        s5 = prs.slides.add_slide(BLANK)
        box(s5, sx(0), sy(0), sw(1200), sh(55), 'a80000', 'a80000')
        tb(s5, sx(50), sy(5), sw(1100), sh(28), '11. The AI Arms Race: Offensive Speed vs. Defensive Lag', 20, True, 'ffffff', 'center')
        tb(s5, sx(100), sy(32), sw(1000), sh(20), 'Threat actors adopt AI with no constraints \u2014 defensive DFIR is held back by regulation and tradition', 11, False, 'ffcccc', 'center')

        # Offensive side
        box(s5, sx(30), sy(70), sw(550), sh(260), 'fff5f5', 'a80000')
        tb(s5, sx(40), sy(75), sw(530), sh(24), '\u2694 Offensive (Threat Actors)', 15, True, 'a80000', 'center')
        off_items = [
            '\U0001f534 Zero-day weaponization: months \u2192 hours',
            '\U0001f534 Nation-state AI agents casting wider nets',
            '\U0001f534 No regulatory constraints on AI usage',
            '\U0001f534 Automated reconnaissance at machine speed',
            '\U0001f534 Polymorphic malware generated by LLMs',
        ]
        for oi, ot in enumerate(off_items):
            tb(s5, sx(55), sy(106 + oi * 40), sw(510), sh(35), ot, 12, False, '660000')

        # VS divider
        tb(s5, sx(580), sy(170), sw(40), sh(30), 'VS', 16, True, '888888', 'center')

        # Defensive side
        box(s5, sx(620), sy(70), sw(550), sh(260), 'f0f8ff', '0078d4')
        tb(s5, sx(630), sy(75), sw(530), sh(24), '\U0001f6e1 Defensive (DFIR Vendors)', 15, True, '0078d4', 'center')
        def_items = [
            '\U0001f535 AI adoption: lagging significantly',
            '\U0001f535 Regulatory & legislative constraints',
            '\U0001f535 "Black box" trust deficit with practitioners',
            '\U0001f535 100% of traditional vendors still legacy-integrated',
            '\U0001f535 AI focused on reports & admin, not investigation',
        ]
        for di, dt in enumerate(def_items):
            tb(s5, sx(635), sy(106 + di * 40), sw(520), sh(35), dt, 12, False, '003366')

        # Current vs recommended focus
        tb(s5, sx(40), sy(345), sw(550), sh(22), 'Where AI Is Applied Today vs. Where It Should Be', 13, True, 'a80000')

        tb(s5, sx(40), sy(375), sw(250), sh(18), '\U0001f7e0 Current Focus (Traditional)', 11, True, 'ca5010')
        curr_tags = ['Report Writing', 'Program Mgmt', 'SOAR Automation', 'Compliance Docs', 'Planning']
        for ci, ct in enumerate(curr_tags):
            box(s5, sx(40 + ci * 110), sy(400), sw(105), sh(28), 'fff0e0', 'ca5010')
            tb(s5, sx(42 + ci * 110), sy(403), sw(101), sh(24), ct, 9, False, 'ca5010', 'center')

        tb(s5, sx(40), sy(440), sw(280), sh(18), '\U0001f7e2 Where AI Should Focus (Critical Path)', 11, True, '107c10')
        rec_tags = ['Deep Forensic Analysis', 'Triage & Scoping', 'Containment', 'Timeline Recon', 'Malware Rev Eng']
        for ri, rt in enumerate(rec_tags):
            box(s5, sx(40 + ri * 130), sy(465), sw(125), sh(28), 'f0fff0', '107c10')
            tb(s5, sx(42 + ri * 130), sy(468), sw(121), sh(24), rt, 9, False, '107c10', 'center')

        # Bridging trust divide
        tb(s5, sx(40), sy(510), sw(400), sh(22), 'Bridging the Trust Divide', 13, True, '0078d4')
        trust_items = [
            ('\U0001f4cb Evidentiary Standards', 'Daubert and Federal Rule 901 \u2014 AI must produce repeatable, documented processes.'),
            ('\U0001f578 Digital Forensic Knowledge Graphs', 'DFKGs visualize timelines and evidence linkages at machine speed.'),
            ('\U0001f3db Government & Insurance Validation', 'Government oversight of AI forensic tools and cyber insurance validation ahead.'),
        ]
        for ti, (tt, td) in enumerate(trust_items):
            ty = 540 + ti * 55
            box(s5, sx(40), sy(ty), sw(1120), sh(48), 'f8f8f5', '0078d4')
            tb(s5, sx(55), sy(ty + 3), sw(400), sh(20), tt, 12, True, '0078d4')
            tb(s5, sx(55), sy(ty + 22), sw(1090), sh(22), td, 10, False, '555555')

        # Bottom callout
        box(s5, sx(40), sy(710), sw(1120), sh(30), 'fff0f0', 'a80000')
        tb(s5, sx(50), sy(713), sw(1100), sh(24), 'The question is not whether AI will transform DFIR \u2014 but which vendors will lead and which will be left behind.', 12, True, 'a80000', 'center')

        # ────────────────────────────────────────────────────────────
        # SLIDE 12: Tandem Model
        # ────────────────────────────────────────────────────────────
        s6 = prs.slides.add_slide(BLANK)
        box(s6, sx(0), sy(0), sw(1200), sh(55), '107c10', '107c10')
        tb(s6, sx(50), sy(5), sw(1100), sh(28), '12. The Tandem Model: AI + Human Investigation', 20, True, 'ffffff', 'center')
        tb(s6, sx(100), sy(32), sw(1000), sh(20), 'Neither full automation nor manual-only delivers optimal outcomes', 11, False, 'c0efc0', 'center')

        # How it works flow
        tb(s6, sx(40), sy(65), sw(400), sh(22), 'How the Tandem Model Works', 13, True, '107c10')
        flow = [
            ('\U0001f916 AI Agent Layer', 'Data ingestion, correlation,\ntimeline reconstruction', '0078d4', 'f0f8ff'),
            ('\U0001f91d Validation Layer', 'Human experts verify methodology,\nconfirm evidence chain', '107c10', 'f0fff0'),
            ('\U0001f4ca Actionable Output', 'Court-admissible reports with\nevidence lineage and audit trails', '8764b8', 'f5f0ff'),
        ]
        for fi, (ft, fd, fc, fb) in enumerate(flow):
            fx = 40 + fi * 380
            box(s6, sx(fx), sy(92), sw(350), sh(100), fb, fc)
            tb(s6, sx(fx + 10), sy(97), sw(330), sh(22), ft, 14, True, fc, 'center')
            tb(s6, sx(fx + 10), sy(125), sw(330), sh(60), fd, 11, False, '555555', 'center')
            if fi < 2:
                tb(s6, sx(fx + 355), sy(125), sw(25), sh(30), '\u2192' if fi == 1 else '\u27f7', 18, True, '888888', 'center')

        # Category donut representations (as score boxes)
        tb(s6, sx(40), sy(210), sw(400), sh(22), 'Who Owns What: Category Strengths', 13, True, '0078d4')
        donuts = [
            ('Traditional', '3.92', '\u2713 Governance, CoC', 'ca5010', 'fff8f0'),
            ('AI-First Startups', '4.11', '\u2713 Investigation, Speed', '0078d4', 'f0f8ff'),
            ('AI-First Non-Startups', '4.04', '\u2713 Investigation, Balance', '8764b8', 'f5f0ff'),
        ]
        for di, (dn, dv, ds, dc, db) in enumerate(donuts):
            dx = 40 + di * 380
            box(s6, sx(dx), sy(240), sw(350), sh(90), db, dc)
            tb(s6, sx(dx + 10), sy(248), sw(330), sh(28), dv, 28, True, dc, 'center')
            tb(s6, sx(dx + 10), sy(280), sw(330), sh(18), dn, 12, True, '333333', 'center')
            tb(s6, sx(dx + 10), sy(300), sw(330), sh(22), ds, 10, False, '107c10', 'center')

        # Critical Sub-Capability Gaps
        tb(s6, sx(40), sy(345), sw(500), sh(22), 'Critical Sub-Capability Gaps (AI-First vs. Traditional)', 13, True, 'ca5010')

        tb(s6, sx(40), sy(375), sw(200), sh(18), '\U0001f535 AI-First Leads', 11, True, '0078d4')
        ai_leads = [('Containment & Isolation', '+0.62'), ('Visibility Gap Analysis', '+0.61'), ('Triage & Scoping', '+0.60'), ('Malware Reverse Eng.', '+0.59')]
        for ai, (an, av) in enumerate(ai_leads):
            box(s6, sx(40), sy(398 + ai * 35), sw(540), sh(30), 'f0f8ff', '0078d4')
            tb(s6, sx(50), sy(401 + ai * 35), sw(400), sh(24), an, 11, False, '333333')
            tb(s6, sx(450), sy(401 + ai * 35), sw(120), sh(24), av, 14, True, '0078d4', 'right')

        tb(s6, sx(620), sy(375), sw(200), sh(18), '\U0001f7e0 Traditional Leads', 11, True, 'ca5010')
        trad_leads = [('Chain of Custody', '+0.37'), ('Post-Incident Learning', '+0.26'), ('Expert Witness Support', '+0.13')]
        for ti, (tn, tv) in enumerate(trad_leads):
            box(s6, sx(620), sy(398 + ti * 35), sw(540), sh(30), 'fff8f0', 'ca5010')
            tb(s6, sx(630), sy(401 + ti * 35), sw(400), sh(24), tn, 11, False, '333333')
            tb(s6, sx(1030), sy(401 + ti * 35), sw(120), sh(24), tv, 14, True, 'ca5010', 'right')

        box(s6, sx(620), sy(398 + 3 * 35), sw(540), sh(30), 'f8f8f5', 'cccccc')
        tb(s6, sx(630), sy(401 + 3 * 35), sw(520), sh(24), 'Neither category scores 4.0+ in every pillar', 11, False, '888888', 'center')

        # Path to Convergence
        tb(s6, sx(40), sy(550), sw(400), sh(22), 'The Path to Convergence', 13, True, '107c10')
        convergence = [
            ('\U0001f527 Traditional Must Do', 'Apply AI to investigation and remediation \u2014 move beyond SOAR and report automation toward the critical path.', 'ca5010'),
            ('\U0001f3d7 AI-First Must Do', 'Build governance, compliance, and legal defensibility \u2014 chain of custody, expert witness, post-incident docs.', '0078d4'),
            ('\U0001f3af The Destination', 'AI shifts investigators from data grinding to validation and interpretation of actionable evidence.', '107c10'),
        ]
        for ci, (ct, cd, cc) in enumerate(convergence):
            cy = 578 + ci * 45
            box(s6, sx(40), sy(cy), sw(1120), sh(40), 'f8f8f5', cc)
            tb(s6, sx(55), sy(cy + 2), sw(300), sh(18), ct, 11, True, cc)
            tb(s6, sx(360), sy(cy + 2), sw(790), sh(35), cd, 10, False, '555555')

        # Bottom callout
        box(s6, sx(40), sy(715), sw(1120), sh(28), 'f0fff0', '107c10')
        tb(s6, sx(50), sy(717), sw(1100), sh(24), 'Product leaders not fully invested in agentic AI for DFIR will find their solutions irrelevant by 2030.', 11, True, '107c10', 'center')

        # ────────────────────────────────────────────────────────────
        # SLIDE 13: Hand-Drawn Agentic Shift (reuse existing layout)
        # ────────────────────────────────────────────────────────────
        s7 = prs.slides.add_slide(BLANK)

        def sx7(v): return int(v / 1200 * SW)
        def sy7(v): return int(v / 900 * SH)
        def sw7(v): return int(v / 1200 * SW)
        def sh7(v): return int(v / 900 * SH)

        def box7(l, t, w, h, fill=None, border=None):
            return box(s7, sx7(l), sy7(t), sw7(w), sh7(h), fill, border)
        def tb7(l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            return tb(s7, sx7(l), sy7(t), sw7(w), sh7(h), text, sz, bold, col, al)
        def ml7(l, t, w, h, lines, al='left'):
            return ml(s7, sx7(l), sy7(t), sw7(w), sh7(h), lines, al)

        # Title banner
        box7(0, 0, 1200, 70, '1a3a5c', '1a3a5c')
        tb7(50, 5, 1100, 32, 'The Agentic Shift: Revolutionizing Digital Forensics & Incident Response (DFIR)', 22, True, 'ffffff', 'center')
        tb7(100, 38, 1000, 28, 'From manual, human-dependent investigation \u2192 high-speed, transparent, legally defensible AI models', 11, False, 'ccd8e8', 'center')

        # Left header
        box7(20, 85, 520, 32, 'e0ddd5', 'b0a898')
        tb7(30, 88, 500, 28, 'TRADITIONAL MANUAL DFIR: FROM MANUAL GRINDING', 11, True, '5a503c', 'center')

        # Left boxes
        box7(30, 135, 240, 130, 'fff8f0', 'ca5010')
        ml7(40, 145, 220, 120, [('Hours to Days', 16, True, 'ca5010'), ('Analysis Speed', 13, True, '5a503c'), ('', 6, False, None), ('Manual data grinding through', 11, False, '777777'), ('massive evidence sets', 11, False, '777777')], 'center')

        box7(290, 135, 240, 130, 'f5f0f0', 'a80000')
        ml7(300, 145, 220, 120, [('Human-Driven / Opaque', 16, True, 'a80000'), ('Methodology', 13, True, '5a503c'), ('', 6, False, None), ('"Black box" concern blocks', 11, False, '777777'), ('AI adoption in investigation', 11, False, '777777')], 'center')

        box7(30, 280, 240, 120, 'fff5f5', 'a80000')
        ml7(40, 288, 220, 108, [('70%', 28, True, 'a80000'), ('of Manual Models Replaced by 2030', 12, True, '5a503c'), ('', 4, False, None), ('Traditional forensic investigations', 10, False, '777777'), ('losing ground to agentic solutions', 10, False, '777777')], 'center')

        box7(290, 280, 240, 120, 'f8f5ff', '8764b8')
        ml7(300, 288, 220, 108, [('Core Focus:', 14, True, '5a503c'), ('Reporting & Admin', 14, True, 'ca5010'), ('', 4, False, None), ('AI applied to reports, planning,', 10, False, '777777'), ('compliance \u2014 not critical path', 10, False, '777777')], 'center')

        # Center: DFKG + Tandem
        box7(520, 420, 160, 160, 'e8f4ff', '0078d4')
        ml7(525, 445, 150, 130, [('Digital Forensic', 13, True, '1a3a5c'), ('Knowledge Graph', 13, True, '1a3a5c'), ('(DFKGs)', 12, False, '0078d4'), ('', 4, False, None), ('Visualizes linkages + timelines', 10, False, '555555'), ('at machine speed', 10, False, '555555')], 'center')

        box7(505, 600, 190, 100, 'f0fff0', '107c10')
        ml7(510, 608, 180, 88, [('Tandem Human-AI Model', 13, True, '107c10'), ('', 4, False, None), ('AI speed + human validation', 11, False, '555555'), ('overcomes black box trust', 11, False, '555555')], 'center')

        # Right header
        box7(660, 85, 520, 32, '00b4d8', '0095b3')
        tb7(670, 88, 500, 28, 'AGENTIC AI DFIR: TO MACHINE-SPEED ANALYSIS', 11, True, 'ffffff', 'center')

        # Right boxes
        box7(670, 135, 240, 130, 'f0faff', '0078d4')
        ml7(680, 145, 220, 120, [('Under 12 Minutes', 16, True, '0078d4'), ('Investigation Speed', 13, True, '1a3a5c'), ('', 6, False, None), ('AI agents complete complex', 11, False, '555555'), ('investigations replacing hours', 11, False, '555555')], 'center')

        box7(930, 135, 240, 130, 'f0fff4', '107c10')
        ml7(940, 145, 220, 120, [('85%', 28, True, '107c10'), ('Source Attribution', 13, True, '1a3a5c'), ('', 4, False, None), ('ForensicLLM pinpoints specific', 10, False, '555555'), ('file paths for every claim', 10, False, '555555')], 'center')

        box7(670, 280, 240, 120, 'f5f0ff', '8764b8')
        ml7(680, 288, 220, 108, [('Legal Admissibility', 14, True, '8764b8'), ('& Rule 901', 14, True, '8764b8'), ('', 4, False, None), ('Documented agent interactions meet', 10, False, '555555'), ('Daubert + Federal Rule standards', 10, False, '555555')], 'center')

        box7(930, 280, 240, 120, 'fffff0', '0078d4')
        ml7(940, 288, 220, 108, [('Explainable /', 14, True, '0078d4'), ('Methodology Engine', 14, True, '0078d4'), ('', 4, False, None), ('SHAP + LIME frameworks deliver', 10, False, '555555'), ('transparent, auditable reasoning', 10, False, '555555')], 'center')

        # Vendor stats
        box7(65, 425, 170, 75, 'fff5f0', 'ca5010')
        ml7(70, 430, 160, 68, [('62%', 20, True, 'ca5010'), ('Traditional \u2022 86 vendors', 10, False, '5a503c')], 'center')

        box7(325, 440, 150, 75, 'f0f8ff', '0078d4')
        ml7(330, 445, 140, 68, [('23%', 20, True, '0078d4'), ('AI-First Startups \u2022 32', 10, False, '1a3a5c')], 'center')

        box7(80, 530, 140, 65, 'f8f0ff', '8764b8')
        ml7(85, 535, 130, 58, [('11%', 20, True, '8764b8'), ('AI-First Other \u2022 15', 10, False, '5a3c7c')], 'center')

        box7(960, 435, 180, 80, 'f0fff4', '107c10')
        ml7(965, 440, 170, 72, [('138', 28, True, '107c10'), ('Vendors Analyzed', 12, True, '1a3a5c')], 'center')

        box7(765, 495, 170, 70, 'f0faff', '0078d4')
        ml7(770, 500, 160, 62, [('80%+', 20, True, '0078d4'), ('ForensicLLM Accuracy', 10, False, '1a3a5c')], 'center')

        # Bottom comparison bars
        box7(20, 730, 1160, 150, 'f8f8f5', 'e0ddd5')
        tb7(50, 738, 200, 20, 'TRADITIONAL', 11, True, 'ca5010')
        tb7(550, 738, 100, 20, 'vs.', 12, True, '888888', 'center')
        tb7(950, 738, 200, 20, 'AI-POWERED', 11, True, '0078d4', 'right')

        s7_bars = [
            (775, 'Analysis Speed: Hours to Days', 'Analysis Speed: Under 12 Minutes'),
            (810, 'Methodology: Human-Driven / Opaque', 'Methodology: Explainable / Engine'),
            (845, 'Core Focus: Reporting & Admin', 'Core Focus: Triage, Containment & Forensics'),
        ]
        for y, left_txt, right_txt in s7_bars:
            box7(50, y, 460, 28, 'fce8d8', 'ca5010')
            tb7(55, y + 3, 450, 24, left_txt, 11, False, 'ca5010', 'center')
            tb7(520, y + 1, 160, 24, '\u2192', 16, True, '107c10', 'center')
            box7(690, y, 460, 28, 'd8eeff', '0078d4')
            tb7(695, y + 3, 450, 24, right_txt, 11, False, '0078d4', 'center')

        tb7(800, 875, 380, 20, '\u00a9 Gartner Research \u2022 DFIR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        # ── Save and return ──
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from flask import send_file
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='DFIR_All_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/mdr-infographic-pptx', methods=['GET'])
def mdr_infographic_pptx():
    """Generate an editable PowerPoint slide of the MDR Pricing Evolution infographic."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        def rgb(hex_str):
            h = hex_str.lstrip('#')
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def add_box(left, top, width, height, fill_hex=None, border_hex=None):
            shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
            if fill_hex:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(fill_hex)
            else:
                shape.fill.background()
            if border_hex:
                shape.line.color.rgb = rgb(border_hex)
                shape.line.width = Pt(1.5)
            else:
                shape.line.fill.background()
            return shape

        def add_textbox(left, top, width, height, text, size=11, bold=False, color_hex='333333', align='left'):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = rgb(color_hex)
            if align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            return txBox

        def add_ml(left, top, width, height, lines, size=11, color_hex='555555', align='left'):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(size)
                p.font.bold = bld
                p.font.color.rgb = rgb(col) if col else rgb(color_hex)
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
            return txBox

        SW = Inches(13.333)
        SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 920 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 920 * SH)

        # Title banner
        add_box(sx(0), sy(0), sw(1200), sh(70), '1a3a5c', '1a3a5c')
        add_textbox(sx(50), sy(5), sw(1100), sh(32),
                    "AI Is Reshaping MDR Pricing \u2014 But Most Vendors Aren't Ready",
                    22, True, 'ffffff', 'center')
        add_textbox(sx(100), sy(38), sw(1000), sh(28),
                    'From subscription-only models \u2192 outcome-based, AI-driven commercial innovation',
                    11, False, 'ccd8e8', 'center')

        # Left header
        add_box(sx(20), sy(85), sw(520), sh(32), 'e0ddd5', 'b0a898')
        add_textbox(sx(30), sy(88), sw(500), sh(28), 'SUBSCRIPTION-ONLY: THE OLD MODEL', 11, True, '5a503c', 'center')

        # Left Box 1
        add_box(sx(30), sy(135), sw(240), sh(130), 'fff8f0', 'ca5010')
        add_ml(sx(40), sy(140), sw(220), sh(120), [
            ('67%', 28, True, 'ca5010'),
            ('Subscription-Only', 15, True, 'ca5010'),
            ('64 of 95 Vendors', 12, True, '5a503c'),
            ('Flat-rate pricing with no', 11, False, '777777'),
            ('outcome alignment', 11, False, '777777'),
        ], align='center')

        # Left Box 2
        add_box(sx(290), sy(135), sw(240), sh(130), 'f5f0f0', 'a80000')
        add_ml(sx(300), sy(140), sw(220), sh(120), [
            ('1.63 / 5.0', 24, True, 'a80000'),
            ('Outcome Pricing Score', 13, True, '5a503c'),
            ('86% of vendors score 2.0 or', 11, False, '777777'),
            ('below \u2014 median is just 1.0', 11, False, '777777'),
        ], align='center')

        # Left Box 3
        add_box(sx(30), sy(280), sw(240), sh(120), 'fff5f5', 'a80000')
        add_ml(sx(40), sy(288), sw(220), sh(108), [
            ('AI Dividend Captured', 14, True, 'a80000'),
            ('', 4, False, None),
            ('Vendors keep AI efficiency', 10, False, '777777'),
            ('gains as margin, not', 10, False, '777777'),
            ('sharing with buyers', 10, False, '777777'),
        ], align='center')

        # Left Box 4
        add_box(sx(290), sy(280), sw(240), sh(120), 'f8f5ff', '8764b8')
        add_ml(sx(300), sy(288), sw(220), sh(108), [
            ('2.35 / 5.0', 22, True, '8764b8'),
            ('Market Average', 14, True, '5a503c'),
            ('Overall pricing maturity:', 11, False, '777777'),
            ('"Developing" stage', 11, False, '777777'),
        ], align='center')

        # Center: AI Pricing Catalyst
        add_box(sx(520), sy(420), sw(160), sh(160), 'e8f4ff', '0078d4')
        add_ml(sx(525), sy(440), sw(150), sh(140), [
            ('AI Pricing', 14, True, '1a3a5c'),
            ('Catalyst', 14, True, '1a3a5c'),
            ('', 4, False, None),
            ('1.02-point gap between', 10, False, '555555'),
            ('AI-mature & minimal', 10, False, '555555'),
            ('vendors on all dimensions', 10, False, '555555'),
        ], align='center')

        # Center: Composable
        add_box(sx(505), sy(600), sw(190), sh(100), 'f0fff0', '107c10')
        add_ml(sx(510), sy(608), sw(180), sh(88), [
            ('Composable Models', 13, True, '107c10'),
            ('2.80 vs 2.16', 16, True, '107c10'),
            ('Modular pricing beats', 10, False, '555555'),
            ('subscription-only by 30%', 10, False, '555555'),
        ], align='center')

        # Right header
        add_box(sx(660), sy(85), sw(520), sh(32), '00b4d8', '0095b3')
        add_textbox(sx(670), sy(88), sw(500), sh(28), 'AI-DRIVEN: OUTCOME-BASED FUTURE', 11, True, 'ffffff', 'center')

        # Right Box 1
        add_box(sx(670), sy(135), sw(240), sh(130), 'f0faff', '0078d4')
        add_ml(sx(680), sy(140), sw(220), sh(120), [
            ('40%', 28, True, '0078d4'),
            ('Outcome-Linked', 15, True, '0078d4'),
            ('by 2028', 13, True, '1a3a5c'),
            ('MDR contracts with at least', 11, False, '555555'),
            ('one outcome component', 11, False, '555555'),
        ], align='center')

        # Right Box 2
        add_box(sx(930), sy(135), sw(240), sh(130), 'f0fff4', '107c10')
        add_ml(sx(940), sy(140), sw(220), sh(120), [
            ('2.86', 28, True, '107c10'),
            ('AI-Mature Average', 13, True, '1a3a5c'),
            ('vs 1.84 for AI-Minimal:', 11, False, '555555'),
            ('55% higher pricing', 11, False, '555555'),
            ('sophistication', 11, False, '555555'),
        ], align='center')

        # Right Box 3
        add_box(sx(670), sy(280), sw(240), sh(120), 'f5f0ff', '8764b8')
        add_ml(sx(680), sy(288), sw(220), sh(108), [
            ('Breach Warranties', 14, True, '8764b8'),
            ('', 4, False, None),
            ('Performance credits, risk-', 10, False, '555555'),
            ('score pricing, shared', 10, False, '555555'),
            ('savings models', 10, False, '555555'),
        ], align='center')

        # Right Box 4
        add_box(sx(930), sy(280), sw(240), sh(120), 'fffff0', '0078d4')
        add_ml(sx(940), sy(288), sw(220), sh(108), [
            ('AI Efficiency Sharing', 14, True, '0078d4'),
            ('', 4, False, None),
            ('Cost-per-incident trending,', 10, False, '555555'),
            ('AI utilization dashboards,', 10, False, '555555'),
            ('graduated discounts', 10, False, '555555'),
        ], align='center')

        # Stat callouts
        add_box(sx(40), sy(425), sw(180), sh(75), 'fff5f0', 'ca5010')
        add_ml(sx(45), sy(430), sw(170), sh(68), [
            ('86%', 22, True, 'ca5010'),
            ('Score \u2264 2.0 on PRC-SUC', 10, False, '5a503c'),
        ], align='center')

        add_box(sx(300), sy(440), sw(160), sh(70), 'f0f8ff', '0078d4')
        add_ml(sx(305), sy(445), sw(150), sh(62), [
            ('19', 22, True, '0078d4'),
            ('AI-Significant Vendors', 10, False, '1a3a5c'),
        ], align='center')

        add_box(sx(60), sy(530), sw(160), sh(65), 'f8f0ff', '8764b8')
        add_ml(sx(65), sy(535), sw(150), sh(58), [
            ('0', 22, True, '8764b8'),
            ('Transformative-Tier', 10, False, '5a3c7c'),
        ], align='center')

        add_box(sx(960), sy(430), sw(190), sh(80), 'f0fff4', '107c10')
        add_ml(sx(965), sy(435), sw(180), sh(72), [
            ('95', 28, True, '107c10'),
            ('Vendors Assessed', 12, False, '1a3a5c'),
        ], align='center')

        add_box(sx(760), sy(490), sw(175), sh(70), 'f0faff', '0078d4')
        add_ml(sx(765), sy(495), sw(165), sh(62), [
            ('3.52', 20, True, '0078d4'),
            ('Best: Sub Transparency', 10, False, '1a3a5c'),
        ], align='center')

        # Bottom comparison bars
        add_box(sx(20), sy(730), sw(1160), sh(170), 'f8f8f5', 'e0ddd5')
        add_textbox(sx(50), sy(738), sw(200), sh(20), 'SUBSCRIPTION-ONLY', 11, True, 'ca5010')
        add_textbox(sx(550), sy(738), sw(100), sh(20), 'vs.', 12, True, '888888', 'center')
        add_textbox(sx(950), sy(738), sw(200), sh(20), 'OUTCOME-BASED', 11, True, '0078d4', 'right')

        bars = [
            (775, 'Revenue: Flat-rate, predictable, opaque', 'Revenue: Performance-linked, transparent'),
            (810, 'AI Efficiency: Captured as vendor margin', 'AI Efficiency: Shared with buyers via dashboards'),
            (845, 'Risk: Buyer bears all risk', 'Risk: Vendor co-owns via breach warranties'),
            (880, 'Score: 2.16 avg (subscription-only)', 'Score: 2.80 avg (composable models)'),
        ]
        for y, left_txt, right_txt in bars:
            add_box(sx(50), sy(y), sw(460), sh(28), 'fce8d8', 'ca5010')
            add_textbox(sx(55), sy(y + 3), sw(450), sh(24), left_txt, 10, False, 'ca5010', 'center')
            add_textbox(sx(520), sy(y + 1), sw(160), sh(24), '\u2192', 16, True, '107c10', 'center')
            add_box(sx(690), sy(y), sw(460), sh(28), 'd8eeff', '0078d4')
            add_textbox(sx(695), sy(y + 3), sw(450), sh(24), right_txt, 10, False, '0078d4', 'center')

        add_textbox(sx(800), sy(900), sw(380), sh(20),
                    '\u00a9 Gartner Research \u2022 MDR Market Insight 2026', 9, False, 'aaaaaa', 'right')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='MDR_Pricing_Evolution_Infographic.pptx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/precyber-infographic-pptx', methods=['GET'])
def precyber_infographic_pptx():
    """Generate an editable PowerPoint slide of the PreCyber Market Fragmentation infographic."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io

        vendor_file = _get_precyber_vendor_file()
        if not vendor_file:
            return jsonify({'error': 'PreCyber vendor data not found'}), 404
        st = _compute_precyber_stats(vendor_file)
        n = st['vendor_count']
        pp = st['pillar_penetration']
        dm = st['delivery_models']
        fs_count = st['full_spectrum_count']
        fs_pct = st['full_spectrum_pct']
        maj_count = st['majority_spectrum_count']
        maj_pct = st['majority_spectrum_pct']
        narrow_count = st['narrow_count']
        narrow_pct = st['narrow_pct']
        blind_pct = st['blind_spot_pct']
        no_amt_pct = st['no_amt_pct']
        dm_ds = dm.get('direct_service', {})
        dm_pp = dm.get('platform_plus_partner', {})
        dm_po = dm.get('platform_only', {})
        ds_count = dm_ds.get('count', 0)
        ppp_count = dm_pp.get('count', 0)
        po_count = dm_po.get('count', 0)
        ds_pct = round(ds_count * 100 / n) if n else 0
        ppp_pct = round(ppp_count * 100 / n) if n else 0
        po_pct = round(po_count * 100 / n) if n else 0
        ds_pa = dm_ds.get('pillar_avgs', {})
        ppp_pa = dm_pp.get('pillar_avgs', {})
        po_pa = dm_po.get('pillar_avgs', {})

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        def rgb(hex_str):
            h = hex_str.lstrip('#')
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def add_box(left, top, width, height, fill_hex=None, border_hex=None):
            shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
            if fill_hex:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(fill_hex)
            else:
                shape.fill.background()
            if border_hex:
                shape.line.color.rgb = rgb(border_hex)
                shape.line.width = Pt(1.5)
            else:
                shape.line.fill.background()
            return shape

        def add_textbox(left, top, width, height, text, size=11, bold=False, color_hex='333333', align='left'):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = rgb(color_hex)
            if align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            return txBox

        def add_ml(left, top, width, height, lines, size=11, color_hex='555555', align='left'):
            txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(size)
                p.font.bold = bld
                p.font.color.rgb = rgb(col) if col else rgb(color_hex)
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
            return txBox

        SW = Inches(13.333)
        SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 920 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 920 * SH)

        # Title banner
        add_box(sx(0), sy(0), sw(1200), sh(70), '1a3a5c', '1a3a5c')
        add_textbox(sx(50), sy(5), sw(1100), sh(32),
                    'The Preemptive Cybersecurity Market Is Dangerously Fragmented',
                    22, True, 'ffffff', 'center')
        add_textbox(sx(100), sy(38), sw(1000), sh(28),
                    f'Only {fs_pct}% of {n} vendors achieve full-spectrum coverage across all five pillars',
                    11, False, 'ccd8e8', 'center')

        # Pillar header
        add_box(sx(20), sy(85), sw(1160), sh(32), 'e0ddd5', 'b0a898')
        add_textbox(sx(30), sy(88), sw(1140), sh(28),
                    'THE FIVE PILLARS OF PREEMPTIVE CYBERSECURITY', 11, True, '5a503c', 'center')

        # Five pillar boxes
        pillars = [
            ('EXM', 'Exposure Mgmt', f"{pp['EXM']['pct']}%", '107c10', 'f0fff4'),
            ('PPM', 'Posture & Policy', f"{pp['PPM']['pct']}%", '0078d4', 'f0faff'),
            ('ADR', 'Detection & Response', f"{pp['ADR']['pct']}%", '8764b8', 'f5f0ff'),
            ('SVC', 'Services & Capability', f"{pp['SVC']['pct']}%", 'ca5010', 'fff8f0'),
            ('AMT', 'Adversary Mgmt', f"{pp['AMT']['pct']}%", 'a80000', 'fff5f5'),
        ]
        for i, (label, name, pct, color, bg) in enumerate(pillars):
            x = 30 + i * 232
            add_box(sx(x), sy(130), sw(215), sh(100), bg, color)
            add_ml(sx(x + 5), sy(133), sw(205), sh(92), [
                (pct, 24, True, color),
                (label, 14, True, color),
                (name, 11, False, '555555'),
                ('Vendor Coverage', 9, False, '888888'),
            ], align='center')

        # Delivery models header
        add_box(sx(20), sy(250), sw(1160), sh(28), '00b4d8', '0095b3')
        add_textbox(sx(30), sy(253), sw(1140), sh(24),
                    'THREE DELIVERY MODELS \u2014 THREE STRATEGIC TRADE-OFFS', 11, True, 'ffffff', 'center')

        # Direct Service
        add_box(sx(30), sy(290), sw(360), sh(160), 'f0fff4', '107c10')
        add_ml(sx(35), sy(295), sw(350), sh(150), [
            ('Direct Service Providers', 14, True, '107c10'),
            (f'{ds_count} vendors ({ds_pct}%)', 12, True, '1a3a5c'),
            (f"SVC: {ds_pa.get('SVC', 0):.2f} (highest) \u2022 ADR: {ds_pa.get('ADR', 0):.2f}", 10, False, '555555'),
            ('Own SOC + analyst teams', 10, False, '555555'),
            ('Single accountability point', 10, False, '555555'),
            ('\u26a0 Limited platform depth', 10, False, '555555'),
        ], align='center')

        # Platform + Partner
        add_box(sx(420), sy(290), sw(360), sh(160), 'f0faff', '0078d4')
        add_ml(sx(425), sy(295), sw(350), sh(150), [
            ('Platform + Partner', 14, True, '0078d4'),
            (f'{ppp_count} vendors ({ppp_pct}%)', 12, True, '1a3a5c'),
            (f"AMT: {ppp_pa.get('AMT', 0):.2f} (highest) \u2022 PPM: {ppp_pa.get('PPM', 0):.2f}", 10, False, '555555'),
            ('Tech platform + MSSP delivery', 10, False, '555555'),
            ('Broadest pillar coverage', 10, False, '555555'),
            (f"\u26a0 Accountability gaps (SVC: {ppp_pa.get('SVC', 0):.2f})", 10, False, '555555'),
        ], align='center')

        # Platform-Only
        add_box(sx(810), sy(290), sw(360), sh(160), 'fff8f0', 'ca5010')
        add_ml(sx(815), sy(295), sw(350), sh(150), [
            ('Platform-Only', 14, True, 'ca5010'),
            (f'{po_count} vendors ({po_pct}%)', 12, True, '1a3a5c'),
            (f"SVC: {po_pa.get('SVC', 0):.2f} \u2022 AMT: {po_pa.get('AMT', 0):.2f} (lowest)", 10, False, '555555'),
            ('No service delivery mechanism', 10, False, '555555'),
            ('88% score < 2.0 on services', 10, False, '555555'),
            ('\u26a0 Structural service deficit', 10, False, '555555'),
        ], align='center')

        # Full-Spectrum center box
        add_box(sx(370), sy(470), sw(460), sh(90), 'e8f4ff', '0078d4')
        add_ml(sx(375), sy(475), sw(450), sh(82), [
            ('Full-Spectrum Vendors', 16, True, '1a3a5c'),
            (f'{fs_count} of {n} ({fs_pct}%) score 2.0+ across all 5 pillars', 12, True, '0078d4'),
            ('Top 3: Mandiant \u2022 SentinelOne \u2022 Fortinet', 10, False, '555555'),
            ('No platform-only vendor achieves full coverage', 10, False, '555555'),
        ], align='center')

        # Stat callouts
        add_box(sx(30), sy(490), sw(160), sh(70), 'fff5f0', 'ca5010')
        add_ml(sx(35), sy(495), sw(150), sh(62), [
            (f'{blind_pct}%', 22, True, 'ca5010'),
            ('Have \u22651 Blind Spot', 10, False, '5a503c'),
        ], align='center')

        add_box(sx(210), sy(510), sw(140), sh(55), 'f8f0ff', '8764b8')
        add_ml(sx(215), sy(515), sw(130), sh(48), [
            (f'{no_amt_pct}%', 20, True, '8764b8'),
            ('No AMT Capability', 9, False, '5a3c7c'),
        ], align='center')

        add_box(sx(850), sy(490), sw(155), sh(70), 'f0fff4', '107c10')
        add_ml(sx(855), sy(495), sw(145), sh(62), [
            (str(n), 24, True, '107c10'),
            ('Vendors Assessed', 10, False, '1a3a5c'),
        ], align='center')

        add_box(sx(1030), sy(510), sw(140), sh(55), 'f0faff', '0078d4')
        add_ml(sx(1035), sy(515), sw(130), sh(48), [
            (f'{narrow_pct}%', 20, True, '0078d4'),
            ('Narrow (\u22643 pillars)', 9, False, '1a3a5c'),
        ], align='center')

        # Market segments
        add_box(sx(20), sy(590), sw(1160), sh(250), 'f8f8f5', 'e0ddd5')
        add_textbox(sx(30), sy(598), sw(1140), sh(24),
                    'MARKET SEGMENT BREAKDOWN', 12, True, '5a503c', 'center')

        add_box(sx(50), sy(630), sw(308), sh(50), '107c10', '107c10')
        add_ml(sx(55), sy(633), sw(298), sh(44), [
            (f'{fs_count} \u2013 Full Spectrum ({fs_pct}%)', 12, True, 'ffffff'),
            ('All 5 pillars \u2265 2.0 \u2022 Best positioned', 10, False, 'c0efc0'),
        ], align='center')

        add_box(sx(50), sy(690), sw(422), sh(50), '0078d4', '0078d4')
        add_ml(sx(55), sy(693), sw(412), sh(44), [
            (f'{maj_count} \u2013 Majority Spectrum ({maj_pct}%)', 12, True, 'ffffff'),
            ('4 pillars \u2022 One investment from full coverage', 10, False, 'b0d8ff'),
        ], align='center')

        add_box(sx(50), sy(750), sw(400), sh(50), 'ca5010', 'ca5010')
        add_ml(sx(55), sy(753), sw(390), sh(44), [
            (f'{narrow_count} \u2013 Narrow Spectrum ({narrow_pct}%)', 12, True, 'ffffff'),
            ('\u22643 pillars \u2022 Niche specialists or incomplete', 10, False, 'fce0c8'),
        ], align='center')

        add_textbox(sx(50), sy(815), sw(1100), sh(20),
                    '\u2192 Service delivery (SVC) is the missing link: no platform-only vendor achieves full-spectrum \u2190',
                    11, True, '1a3a5c', 'center')

        add_textbox(sx(800), sy(835), sw(380), sh(20),
                    '\u00a9 Gartner Research \u2022 PreCyber Market Insight 2026', 9, False, 'aaaaaa', 'right')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='PreCyber_Market_Fragmentation_Infographic.pptx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mdr-all-graphics-pptx', methods=['GET'])
def mdr_all_graphics_pptx():
    """Generate a multi-slide editable PowerPoint deck with ALL 9 MDR graphics."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from flask import send_file
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        def box(slide, l, t, w, h, fill=None, border=None):
            s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
            if fill:
                s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
            else:
                s.fill.background()
            if border:
                s.line.color.rgb = rgb(border); s.line.width = Pt(1.5)
            else:
                s.line.fill.background()
            return s

        def tb(slide, l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = text
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = rgb(col)
            if al == 'center': p.alignment = PP_ALIGN.CENTER
            elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        def ml(slide, l, t, w, h, lines, al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(11)
                p.font.bold = bld
                if col: p.font.color.rgb = rgb(col)
                if al == 'center': p.alignment = PP_ALIGN.CENTER
                elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        SW = Inches(13.333); SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 750 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 750 * SH)
        footer = '\u00a9 Gartner Research \u2022 MDR Market Insight 2026'

        # ── SLIDE 1: Pricing Dimension Heatmap ──
        s1 = prs.slides.add_slide(BLANK)
        box(s1, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s1, sx(50), sy(8), sw(1100), sh(28), '1. Pricing Dimension Heatmap', 22, True, 'ffffff', 'center')
        tb(s1, sx(100), sy(36), sw(1000), sh(22), '95 MDR vendors scored across 6 pricing dimensions by commercial model type', 11, False, 'ccd8e8', 'center')

        dims = [
            ('Subscription Transparency', 'PRC-SUB', 3.80, 3.42, 3.52),
            ('Usage-Based Alignment', 'PRC-USG', 2.60, 1.90, 2.14),
            ('Fixed Delivery Pricing', 'PRC-FIX', 2.95, 2.10, 2.38),
            ('Success & Outcome Fees', 'PRC-SUC', 2.41, 1.34, 1.63),
            ('Composability & Maturity', 'PRC-COM', 2.88, 2.05, 2.30),
            ('Pricing-to-Outcomes', 'PRC-OUT', 2.18, 1.50, 1.77),
        ]
        headers = [('Composable\n27 vendors (28%)', '107c10'), ('Subscription-Only\n64 vendors (67%)', 'ca5010'), ('Market Average\n95 vendors', '0078d4')]
        colW = 220; startX = 380; rowH = 55
        for ci, (ht, hc) in enumerate(headers):
            box(s1, sx(startX + ci * colW), sy(75), sw(colW - 10), sh(50), hc, hc)
            tb(s1, sx(startX + ci * colW + 5), sy(78), sw(colW - 20), sh(46), ht, 11, True, 'ffffff', 'center')

        def hm_color(v):
            if v >= 3.5: return '107c10'
            if v >= 2.5: return '0078d4'
            if v >= 2.0: return 'ca5010'
            return 'a80000'
        def hm_bg(v):
            if v >= 3.5: return 'e6f4e6'
            if v >= 2.5: return 'e0f0ff'
            if v >= 2.0: return 'fff0e0'
            return 'ffe0e0'

        for ri, (name, code, comp, sub, mkt) in enumerate(dims):
            y = 135 + ri * rowH
            box(s1, sx(40), sy(y), sw(330), sh(rowH - 5), 'f8f8f5', 'e0ddd5')
            ml(s1, sx(50), sy(y + 5), sw(310), sh(40), [(name, 12, True, '333333'), (code, 9, False, '888888')])
            for ci, v in enumerate([comp, sub, mkt]):
                box(s1, sx(startX + ci * colW), sy(y), sw(colW - 10), sh(rowH - 5), hm_bg(v), hm_color(v))
                tb(s1, sx(startX + ci * colW + 5), sy(y + 10), sw(colW - 20), sh(32), f'{v:.2f}', 20, True, hm_color(v), 'center')

        y = 135 + 6 * rowH + 10
        box(s1, sx(40), sy(y), sw(330), sh(rowH - 5), 'e0ddd5', 'b0a898')
        tb(s1, sx(50), sy(y + 10), sw(310), sh(32), 'Overall Average', 13, True, '333333')
        for ci, v in enumerate([2.80, 2.16, 2.35]):
            box(s1, sx(startX + ci * colW), sy(y), sw(colW - 10), sh(rowH - 5), hm_bg(v), hm_color(v))
            tb(s1, sx(startX + ci * colW + 5), sy(y + 10), sw(colW - 20), sh(32), f'{v:.2f}', 20, True, hm_color(v), 'center')

        legend_y = y + rowH + 15
        for li, (lt, lc) in enumerate([('\u25a0 \u2265 3.50 Strong', '107c10'), ('\u25a0 2.50\u20133.49 Competitive', '0078d4'), ('\u25a0 2.00\u20132.49 Below Target', 'ca5010'), ('\u25a0 < 2.00 Critical', 'a80000')]):
            tb(s1, sx(40 + li * 270), sy(legend_y), sw(260), sh(22), lt, 11, True, lc)
        tb(s1, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 2: 2028 Outcome Pricing Roadmap ──
        s2 = prs.slides.add_slide(BLANK)
        box(s2, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s2, sx(50), sy(8), sw(1100), sh(28), '2. Path to 2028: Outcome-Based Pricing Roadmap', 22, True, 'ffffff', 'center')
        tb(s2, sx(100), sy(36), sw(1000), sh(22), 'Current maturity of each pricing dimension relative to the 3.0/5.0 commercial viability target', 11, False, 'ccd8e8', 'center')

        roadmap = [
            ('Subscription Transparency', 3.52, 3.5, '107c10', '\u2713 Only mature dimension \u2014 already at target'),
            ('Usage-Based Alignment', 2.14, 3.0, 'ca5010', '\u26a0 Buyers see no AI cost reduction benefits'),
            ('Composability & Maturity', 2.30, 3.0, '0078d4', '\u26a0 67% still subscription-only monoliths'),
            ('Fixed Delivery Pricing', 2.38, 3.0, '8764b8', '\u26a0 Limited retainer / project-based options'),
            ('Pricing-to-Outcomes', 1.77, 3.0, 'a80000', '\u26a0 No measurement infrastructure'),
            ('Success & Outcome Fees', 1.63, 3.0, 'a80000', '\u26a0 86% score \u2264 2.0 \u2014 median is 1.0'),
        ]
        for ri, (name, score, target, color, note) in enumerate(roadmap):
            baseY = 80 + ri * 100
            tb(s2, sx(40), sy(baseY), sw(400), sh(22), name, 14, True, color)
            tb(s2, sx(40), sy(baseY + 22), sw(400), sh(16), f'{score:.2f} / 5.00', 10, False, '666666')
            barW = 880
            box(s2, sx(260), sy(baseY + 5), sw(barW), sh(32), 'e8e8e5', 'cccccc')
            fillW = int(barW * score / 5.0)
            box(s2, sx(260), sy(baseY + 5), sw(fillW), sh(32), color, color)
            tb(s2, sx(260 + fillW // 2 - 40), sy(baseY + 9), sw(80), sh(24), f'{score:.2f}', 13, True, 'ffffff', 'center')
            targetX = int(barW * target / 5.0)
            box(s2, sx(260 + targetX - 1), sy(baseY + 2), sw(3), sh(38), 'a80000', 'a80000')
            tb(s2, sx(260 + targetX - 30), sy(baseY + 42), sw(60), sh(14), f'{target:.1f} Target', 8, True, 'a80000', 'center')
            tb(s2, sx(260), sy(baseY + 58), sw(barW), sh(16), note, 10, False, '666666')

        box(s2, sx(40), sy(680), sw(1120), sh(55), 'fff5f5', 'a80000')
        ml(s2, sx(55), sy(683), sw(1090), sh(48), [
            ('SPA \u2014 2028 Strategic Planning Assumption', 12, True, 'a80000'),
            ('"By 2028, 40% of MDR contracts will include outcome-linked pricing, up from <5% today."', 11, False, '555555'),
        ])
        tb(s2, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 3: AI Influence Investment Mismatch ──
        s3 = prs.slides.add_slide(BLANK)
        box(s3, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s3, sx(50), sy(8), sw(1100), sh(28), '3. Investment Mismatch: AI-Significant vs AI-Minimal Vendors', 22, True, 'ffffff', 'center')
        tb(s3, sx(100), sy(36), sw(1000), sh(22), '1.02-point gap between AI-Significant (19) and AI-Minimal (27) vendors across pricing dimensions', 11, False, 'ccd8e8', 'center')

        tb(s3, sx(60), sy(72), sw(200), sh(18), '\u25a0 AI-Significant (19)', 11, True, '0078d4')
        tb(s3, sx(280), sy(72), sw(200), sh(18), '\u25a0 AI-Minimal (27)', 11, True, 'ca5010')
        tb(s3, sx(480), sy(72), sw(200), sh(18), '\u25a0 Recommended Target', 11, True, '107c10')

        mismatch = [
            ('Subscription Transparency', 3.72, 3.30, 3.5),
            ('Usage-Based Alignment', 2.58, 1.60, 3.0),
            ('Success & Outcome Fees', 2.41, 1.10, 3.0),
            ('Composability & Maturity', 2.88, 1.80, 3.0),
            ('Pricing-to-Outcomes', 2.58, 1.04, 3.0),
        ]
        for mi_, (name, sig, minn, ideal) in enumerate(mismatch):
            baseY = 100 + mi_ * 115
            tb(s3, sx(40), sy(baseY), sw(250), sh(22), name, 14, True, '333333')
            barW = 700; labW = 80
            for bi, (label, val, clr) in enumerate([('Significant', sig, '0078d4'), ('Minimal', minn, 'ca5010'), ('Target', ideal, '107c10')]):
                by_ = baseY + 28 + bi * 26
                tb(s3, sx(40), sy(by_), sw(labW), sh(18), label, 10, False, '777777')
                box(s3, sx(40 + labW), sy(by_ - 2), sw(barW), sh(22), 'f0f0ee', 'dddddd')
                box(s3, sx(40 + labW), sy(by_ - 2), sw(int(barW * val / 5.0)), sh(22), clr, clr)
                tb(s3, sx(40 + labW + barW + 10), sy(by_), sw(60), sh(18), f'{val:.2f}', 12, True, clr)

        box(s3, sx(40), sy(680), sw(1120), sh(50), 'f0faff', '0078d4')
        ml(s3, sx(55), sy(683), sw(1090), sh(44), [
            ('Key Insight', 11, True, '0078d4'),
            ('AI-Significant vendors score 55% higher than AI-Minimal. Gap widest on PRC-OUT (2.58 vs 1.04 = +148%).', 10, False, '555555'),
        ])
        tb(s3, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 4: Executive Summary Poster ──
        s4 = prs.slides.add_slide(BLANK)
        box(s4, sx(0), sy(0), sw(1200), sh(55), '0078d4', '0078d4')
        tb(s4, sx(50), sy(5), sw(1100), sh(28), '4. AI Is Reshaping MDR Pricing \u2014 But Most Vendors Aren\'t Ready', 20, True, 'ffffff', 'center')
        tb(s4, sx(100), sy(32), sw(1000), sh(20), 'Market Insight Executive Summary \u2014 95 vendors across 6 pricing dimensions', 11, False, 'b0d8ff', 'center')

        stats = [('95', 'Vendors Assessed', '0078d4'), ('86%', 'Score \u2264 2.0 PRC-SUC', 'a80000'), ('67%', 'Subscription-Only', 'ca5010'), ('2.86', 'AI-Significant Avg', '107c10'), ('0', 'Transformative-Tier', '8764b8')]
        for si, (sv, sl, sc) in enumerate(stats):
            bx = 40 + si * 232
            box(s4, sx(bx), sy(65), sw(215), sh(60), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(bx + 5), sy(68), sw(205), sh(28), sv, 22, True, sc, 'center')
            tb(s4, sx(bx + 5), sy(98), sw(205), sh(22), sl, 10, False, '666666', 'center')

        tb(s4, sx(40), sy(140), sw(300), sh(22), 'KEY FINDINGS', 13, True, '0078d4')
        findings = [
            ('\U0001f4c9 Outcome Adoption Vanishingly Low', '82 of 95 vendors (86%) score \u2264 2.0 on success-linked fees. Median is 1.0.'),
            ('\U0001f916 AI-Mature Benchmark 55% Higher', 'AI-Significant vendors avg 2.86 vs AI-Minimal 1.84. Most pronounced on PRC-OUT: 2.58 vs 1.04.'),
            ('\U0001f9e9 Composable Is Prerequisite', 'Composable vendors score 2.41 on PRC-SUC vs 1.34 for subscription-only \u2014 80% premium.'),
        ]
        for fi, (ft, fb) in enumerate(findings):
            fy = 168 + fi * 75
            box(s4, sx(40), sy(fy), sw(1120), sh(65), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(55), sy(fy + 5), sw(1090), sh(22), ft, 13, True, '333333')
            tb(s4, sx(55), sy(fy + 28), sw(1090), sh(34), fb, 11, False, '555555')

        tb(s4, sx(40), sy(400), sw(400), sh(22), 'RECOMMENDATIONS', 13, True, 'ca5010')
        recs = [
            ('01', 'Redefine Outcomes', 'MTTD/MTTR, risk-score, breach warranties', '0078d4'),
            ('02', 'AI Efficiency Sharing', 'Dashboards, cost-per-incident, AI utilization', '107c10'),
            ('03', 'Go Composable', 'Modular: detect + respond + hunt + outcome', 'ca5010'),
            ('04', 'Phased Roadmap', 'Transparency \u2192 Usage \u2192 Composable \u2192 Outcomes', '8764b8'),
        ]
        for ri, (rn, rt, rd, rc) in enumerate(recs):
            rx = 40 + ri * 285
            box(s4, sx(rx), sy(430), sw(260), sh(85), 'f8f8f5', rc)
            tb(s4, sx(rx + 10), sy(433), sw(240), sh(16), rn, 10, True, rc)
            tb(s4, sx(rx + 10), sy(450), sw(240), sh(22), rt, 14, True, rc, 'center')
            tb(s4, sx(rx + 10), sy(475), sw(240), sh(34), rd, 10, False, '555555', 'center')
            if ri < 3:
                tb(s4, sx(rx + 265), sy(460), sw(20), sh(22), '\u2192', 16, True, '888888', 'center')

        box(s4, sx(40), sy(540), sw(1120), sh(50), 'fff5f5', 'a80000')
        ml(s4, sx(55), sy(543), sw(1090), sh(44), [
            ('\u26a0 "By 2028, 40% of MDR contracts will include outcome-linked pricing \u2014 up from <5% today."', 12, True, 'a80000'),
            ('\u2014 Strategic Planning Assumption', 10, False, '888888'),
        ], 'center')
        tb(s4, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 5: Subscription vs Outcome VS poster ──
        s5 = prs.slides.add_slide(BLANK)
        box(s5, sx(0), sy(0), sw(1200), sh(55), 'a80000', 'a80000')
        tb(s5, sx(50), sy(5), sw(1100), sh(28), '5. The Great MDR Pricing Divide', 20, True, 'ffffff', 'center')
        tb(s5, sx(100), sy(32), sw(1000), sh(20), '67% cling to flat-rate subscriptions while the market demands outcome-based innovation', 11, False, 'ffcccc', 'center')

        # Left side
        box(s5, sx(30), sy(70), sw(540), sh(300), 'fff5f0', 'ca5010')
        tb(s5, sx(40), sy(75), sw(520), sh(24), '\u2699 Subscription-Only (The Old Model)', 16, True, 'ca5010', 'center')
        sub_items = ['\U0001f534 64 of 95 vendors (67%)', '\U0001f534 Average score: 2.16 / 5.0', '\U0001f534 PRC-SUC score: 1.34', '\U0001f534 Flat-rate, predictable, opaque', '\U0001f534 AI efficiency as vendor margin', '\U0001f534 Buyer bears all risk', '\U0001f534 No usage alignment']
        for i, t in enumerate(sub_items):
            tb(s5, sx(55), sy(108 + i * 35), sw(500), sh(30), t, 12, False, '660000')

        tb(s5, sx(575), sy(185), sw(50), sh(28), 'VS', 18, True, '888888', 'center')
        # Right side
        box(s5, sx(630), sy(70), sw(540), sh(300), 'f0faff', '0078d4')
        tb(s5, sx(640), sy(75), sw(520), sh(24), '\U0001f680 Composable / Outcome-Based', 16, True, '0078d4', 'center')
        comp_items = ['\U0001f535 27 of 95 vendors (28%)', '\U0001f535 Average score: 2.80 / 5.0 (+30%)', '\U0001f535 PRC-SUC score: 2.41 (+80%)', '\U0001f535 Modular: detect + respond + hunt', '\U0001f535 AI efficiency shared via dashboards', '\U0001f535 Vendor co-owns risk via warranties', '\U0001f535 Performance-linked, transparent']
        for i, t in enumerate(comp_items):
            tb(s5, sx(645), sy(108 + i * 35), sw(510), sh(30), t, 12, False, '003366')

        tb(s5, sx(40), sy(385), sw(500), sh(22), 'AI PRICING INFLUENCE SPECTRUM', 13, True, '0078d4')
        tiers = [('27', 'AI-Minimal', 'Avg 1.84', 'a80000', 'ffe0e0'), ('49', 'AI-Emerging', 'Avg 2.20', 'ca5010', 'fff0e0'), ('19', 'AI-Significant', 'Avg 2.86', '0078d4', 'e0f0ff'), ('0', 'AI-Transformative', 'No vendor yet', '107c10', 'e6f4e6')]
        for ti, (tv, tl, td, tc, tbg) in enumerate(tiers):
            tx = 40 + ti * 285
            box(s5, sx(tx), sy(415), sw(260), sh(80), tbg, tc)
            tb(s5, sx(tx + 5), sy(420), sw(250), sh(28), tv, 24, True, tc, 'center')
            tb(s5, sx(tx + 5), sy(450), sw(250), sh(18), tl, 12, True, tc, 'center')
            tb(s5, sx(tx + 5), sy(470), sw(250), sh(18), td, 10, False, '555555', 'center')

        box(s5, sx(40), sy(510), sw(1120), sh(35), 'fff0f0', 'a80000')
        tb(s5, sx(50), sy(514), sw(1100), sh(28), 'The question is not whether MDR pricing will change \u2014 but which vendors will lead the shift.', 12, True, 'a80000', 'center')
        tb(s5, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 6: Three Investments Model ──
        s6 = prs.slides.add_slide(BLANK)
        box(s6, sx(0), sy(0), sw(1200), sh(55), '107c10', '107c10')
        tb(s6, sx(50), sy(5), sw(1100), sh(28), '6. The Composable Pricing Imperative', 20, True, 'ffffff', 'center')
        tb(s6, sx(100), sy(32), sw(1000), sh(20), 'Three simultaneous investments \u2014 none is sufficient alone', 11, False, 'c0efc0', 'center')

        pillars = [
            ('\U0001f916 AI Capability Maturity', 'Underwrite outcome guarantees\nUsage-based metering\nOutcome confidence scoring\nCost-per-incident trending', '0078d4', 'f0faff', 'PRC-USG avg 2.14/5'),
            ('\U0001f9e9 Pricing Architecture', 'Composable, modular components\nBase detect + response + hunt\nOutcome guarantees layer\nEfficiency sharing layer', 'ca5010', 'fff8f0', '67% still subscription-only'),
            ('\U0001f4ca Measurement Infrastructure', 'Transparent outcome metrics\nMTTD/MTTR reduction proof\nRisk-score improvement tracking\nAI utilization dashboards', '8764b8', 'f5f0ff', 'PRC-OUT avg 1.77/5'),
        ]
        for pi, (pt, pd, pc, pbg, pg) in enumerate(pillars):
            px = 40 + pi * 380
            box(s6, sx(px), sy(70), sw(350), sh(210), pbg, pc)
            tb(s6, sx(px + 10), sy(78), sw(330), sh(24), pt, 14, True, pc, 'center')
            tb(s6, sx(px + 15), sy(108), sw(320), sh(130), pd, 11, False, '555555', 'center')
            box(s6, sx(px + 20), sy(240), sw(310), sh(28), pc, pc)
            tb(s6, sx(px + 25), sy(243), sw(300), sh(24), 'Gap: ' + pg, 10, True, 'ffffff', 'center')

        tb(s6, sx(40), sy(300), sw(500), sh(22), 'PHASED ROADMAP TO OUTCOME-BASED PRICING', 13, True, 'ca5010')
        phases = [
            ('Phase 1', 'Transparency', 'Publish tiers + metrics', '107c10'),
            ('Phase 2', 'Usage-Aligned', 'Consumption metering', '0078d4'),
            ('Phase 3', 'Composable', 'Modular components', 'ca5010'),
            ('Phase 4', 'Outcomes', 'Breach warranties', '8764b8'),
            ('Phase 5', 'Full Risk-Share', 'Co-owned outcomes', 'a80000'),
        ]
        for phi, (pn, pl, pd, pc) in enumerate(phases):
            px = 40 + phi * 228
            box(s6, sx(px), sy(330), sw(210), sh(75), pc, pc)
            tb(s6, sx(px + 5), sy(335), sw(200), sh(18), pn, 11, True, 'ffffff', 'center')
            tb(s6, sx(px + 5), sy(355), sw(200), sh(22), pl, 14, True, 'ffffff', 'center')
            tb(s6, sx(px + 5), sy(378), sw(200), sh(20), pd, 10, False, 'ffffff', 'center')
            if phi < 4:
                tb(s6, sx(px + 215), sy(355), sw(13), sh(22), '\u2192', 14, True, '888888', 'center')

        scorebox = [('2.35', 'Market Average \u2014 "Developing"', 'ca5010'), ('3.65', 'Top-10 Average \u2014 "Advancing"', '0078d4'), ('1.30', 'Leader / Market Gap (widening)', 'a80000')]
        for si, (sv, sl, sc) in enumerate(scorebox):
            bx = 40 + si * 380
            box(s6, sx(bx), sy(425), sw(350), sh(65), 'f8f8f5', 'e0ddd5')
            tb(s6, sx(bx + 5), sy(428), sw(340), sh(30), sv, 24, True, sc, 'center')
            tb(s6, sx(bx + 5), sy(463), sw(340), sh(22), sl, 11, False, '666666', 'center')

        box(s6, sx(40), sy(510), sw(1120), sh(35), 'f0fff0', '107c10')
        tb(s6, sx(50), sy(514), sw(1100), sh(28), 'Product leaders not invested in composable, AI-enabled pricing will face structural competitive disadvantage by 2027.', 11, True, '107c10', 'center')
        tb(s6, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 7: Hand-Drawn Workflow (reuse existing MDR infographic layout) ──
        s7 = prs.slides.add_slide(BLANK)
        def sx7(v): return int(v / 1200 * SW)
        def sy7(v): return int(v / 920 * SH)
        def sw7(v): return int(v / 1200 * SW)
        def sh7(v): return int(v / 920 * SH)

        box(s7, sx7(0), sy7(0), sw7(1200), sh7(70), '1a3a5c', '1a3a5c')
        tb(s7, sx7(50), sy7(5), sw7(1100), sh7(32), 'AI Is Reshaping MDR Pricing \u2014 But Most Vendors Aren\'t Ready', 22, True, 'ffffff', 'center')
        tb(s7, sx7(100), sy7(38), sw7(1000), sh7(28), 'From subscription-only models \u2192 outcome-based, AI-driven commercial innovation', 11, False, 'ccd8e8', 'center')

        box(s7, sx7(20), sy7(85), sw7(520), sh7(32), 'e0ddd5', 'b0a898')
        tb(s7, sx7(30), sy7(88), sw7(500), sh7(28), 'SUBSCRIPTION-ONLY: THE OLD MODEL', 11, True, '5a503c', 'center')

        box(s7, sx7(30), sy7(130), sw7(240), sh7(130), 'fff8f0', 'ca5010')
        ml(s7, sx7(40), sy7(140), sw7(220), sh7(120), [('67%', 22, True, 'ca5010'), ('Subscription-Only', 12, True, '5a503c'), ('', 4, False, None), ('64 of 95 Vendors', 11, False, '777777'), ('Flat-rate pricing with no', 10, False, '777777'), ('outcome alignment', 10, False, '777777')], 'center')

        box(s7, sx7(290), sy7(130), sw7(240), sh7(130), 'f5f0f0', 'a80000')
        ml(s7, sx7(300), sy7(140), sw7(220), sh7(120), [('1.63 / 5.0', 18, True, 'a80000'), ('Outcome Pricing Score', 12, True, '5a503c'), ('', 4, False, None), ('86% of vendors score 2.0 or', 10, False, '777777'), ('below \u2014 median is just 1.0', 10, False, '777777')], 'center')

        box(s7, sx7(30), sy7(280), sw7(240), sh7(120), 'fff5f5', 'a80000')
        ml(s7, sx7(40), sy7(288), sw7(220), sh7(108), [('\U0001f4b0 AI Dividend Captured', 13, True, 'a80000'), ('', 4, False, None), ('Vendors keep AI efficiency', 11, False, '777777'), ('gains as margin, not', 11, False, '777777'), ('sharing with buyers', 11, False, '777777')], 'center')

        box(s7, sx7(290), sy7(280), sw7(240), sh7(120), 'f8f5ff', '8764b8')
        ml(s7, sx7(300), sy7(288), sw7(220), sh7(108), [('2.35 / 5.0', 18, True, '8764b8'), ('Market Average', 12, True, '5a503c'), ('', 4, False, None), ('Overall pricing maturity:', 10, False, '777777'), ('"Developing" stage', 10, False, '777777')], 'center')

        box(s7, sx7(520), sy7(400), sw7(160), sh7(140), 'e8f4ff', '0078d4')
        ml(s7, sx7(525), sy7(420), sw7(150), sh7(110), [('\U0001f916 AI Pricing', 12, True, '1a3a5c'), ('Catalyst', 12, True, '1a3a5c'), ('', 4, False, None), ('1.02-point gap between', 10, False, '555555'), ('AI-mature & minimal', 10, False, '555555'), ('vendors', 10, False, '555555')], 'center')

        box(s7, sx7(505), sy7(560), sw7(190), sh7(90), 'f0fff0', '107c10')
        ml(s7, sx7(510), sy7(568), sw7(180), sh7(78), [('Composable Models', 12, True, '107c10'), ('2.80 vs 2.16', 14, True, '107c10'), ('', 4, False, None), ('Modular pricing beats', 10, False, '555555'), ('subscription-only by 30%', 10, False, '555555')], 'center')

        box(s7, sx7(660), sy7(85), sw7(520), sh7(32), '00b4d8', '0095b3')
        tb(s7, sx7(670), sy7(88), sw7(500), sh7(28), 'AI-DRIVEN: OUTCOME-BASED FUTURE', 11, True, 'ffffff', 'center')

        box(s7, sx7(670), sy7(130), sw7(240), sh7(130), 'f0faff', '0078d4')
        ml(s7, sx7(680), sy7(140), sw7(220), sh7(120), [('40%', 22, True, '0078d4'), ('Outcome-Linked', 12, True, '1a3a5c'), ('by 2028', 10, False, '0078d4'), ('', 4, False, None), ('MDR contracts with at least', 10, False, '555555'), ('one outcome component', 10, False, '555555')], 'center')

        box(s7, sx7(930), sy7(130), sw7(240), sh7(130), 'f0fff4', '107c10')
        ml(s7, sx7(940), sy7(140), sw7(220), sh7(120), [('2.86', 22, True, '107c10'), ('AI-Mature Average', 12, True, '1a3a5c'), ('', 4, False, None), ('vs 1.84 for AI-Minimal:', 10, False, '555555'), ('55% higher pricing', 10, False, '555555'), ('sophistication', 10, False, '555555')], 'center')

        box(s7, sx7(670), sy7(280), sw7(240), sh7(120), 'f5f0ff', '8764b8')
        ml(s7, sx7(680), sy7(288), sw7(220), sh7(108), [('\U0001f6e1 Breach Warranties', 13, True, '8764b8'), ('', 4, False, None), ('Performance credits, risk-', 11, False, '555555'), ('score pricing, shared', 11, False, '555555'), ('savings models', 11, False, '555555')], 'center')

        box(s7, sx7(930), sy7(280), sw7(240), sh7(120), 'fffff0', '0078d4')
        ml(s7, sx7(940), sy7(288), sw7(220), sh7(108), [('\U0001f4ca AI Efficiency Sharing', 13, True, '0078d4'), ('', 4, False, None), ('Cost-per-incident trending,', 10, False, '555555'), ('AI utilization dashboards,', 10, False, '555555'), ('graduated discounts', 10, False, '555555')], 'center')

        # Stats
        stat_boxes = [
            (65, 430, 170, 70, '86%', 'Score \u2264 2.0 on PRC-SUC', 'ca5010', 'fff5f0'),
            (325, 445, 150, 70, '19', 'AI-Significant Vendors', '0078d4', 'f0f8ff'),
            (80, 530, 140, 60, '0', 'Transformative-Tier', '8764b8', 'f8f0ff'),
            (960, 440, 180, 70, '95', 'Vendors Assessed', '107c10', 'f0fff4'),
            (765, 500, 170, 65, '3.52', 'Best: Sub Transparency', '0078d4', 'f0faff'),
        ]
        for bx, by, bw, bh, bv, bl, bc, bbg in stat_boxes:
            box(s7, sx7(bx), sy7(by), sw7(bw), sh7(bh), bbg, bc)
            ml(s7, sx7(bx + 5), sy7(by + 5), sw7(bw - 10), sh7(bh - 10), [(bv, 18, True, bc), (bl, 9, False, '5a503c')], 'center')

        # Bottom comparison bars
        box(s7, sx7(20), sy7(730), sw7(1160), sh7(170), 'f8f8f5', 'e0ddd5')
        tb(s7, sx7(50), sy7(738), sw7(200), sh7(20), 'SUBSCRIPTION', 11, True, 'ca5010')
        tb(s7, sx7(550), sy7(738), sw7(100), sh7(20), 'vs.', 12, True, '888888', 'center')
        tb(s7, sx7(950), sy7(738), sw7(200), sh7(20), 'OUTCOME-BASED', 11, True, '0078d4', 'right')
        bars7 = [
            (775, 'Revenue: Flat-rate, predictable, opaque', 'Revenue: Performance-linked, transparent'),
            (810, 'AI Efficiency: Captured as vendor margin', 'AI Efficiency: Shared with buyers'),
            (845, 'Risk: Buyer bears all risk', 'Risk: Vendor co-owns via breach warranties'),
            (880, 'Pricing Score: 2.16 avg (sub-only)', 'Pricing Score: 2.80 avg (composable)'),
        ]
        for y, left_txt, right_txt in bars7:
            box(s7, sx7(50), sy7(y), sw7(460), sh7(28), 'fce8d8', 'ca5010')
            tb(s7, sx7(55), sy7(y + 3), sw7(450), sh7(24), left_txt, 10, False, 'ca5010', 'center')
            tb(s7, sx7(520), sy7(y + 1), sw7(160), sh7(24), '\u2192', 14, True, '107c10', 'center')
            box(s7, sx7(690), sy7(y), sw7(460), sh7(28), 'd8eeff', '0078d4')
            tb(s7, sx7(695), sy7(y + 3), sw7(450), sh7(24), right_txt, 10, False, '0078d4', 'center')

        tb(s7, sx7(800), sy7(895), sw7(380), sh7(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 8: Whiteboard Session ──
        s8 = prs.slides.add_slide(BLANK)
        def sx8(v): return int(v / 1260 * SW)
        def sy8(v): return int(v / 1080 * SH)
        def sw8(v): return int(v / 1260 * SW)
        def sh8(v): return int(v / 1080 * SH)

        # Title banner
        box(s8, sx8(0), sy8(0), sw8(1260), sh8(70), '1a3a5c', '1a3a5c')
        tb(s8, sx8(50), sy8(5), sw8(1160), sh8(32), '8. AI Is Reshaping MDR Pricing \u2014 But Most Vendors Aren\'t Ready', 22, True, 'ffffff', 'center')
        tb(s8, sx8(100), sy8(40), sw8(1060), sh8(28), 'subscription-only models \u2192 outcome-based, AI-driven commercial innovation', 11, False, 'ccd8e8', 'center')

        # LEFT ZONE — old model header
        box(s8, sx8(30), sy8(90), sw8(540), sh8(32), 'f5d5d5', 'c0392b')
        tb(s8, sx8(40), sy8(93), sw8(520), sh8(28), '\u26a0 SUBSCRIPTION-ONLY: THE OLD MODEL', 11, True, 'c0392b', 'center')

        # Sticky notes — left zone
        # 67%
        box(s8, sx8(40), sy8(135), sw8(250), sh8(145), 'fce4ec', 'e57373')
        ml(s8, sx8(50), sy8(142), sw8(230), sh8(135), [
            ('67%', 28, True, 'c0392b'), ('Subscription-Only', 12, True, 'b71c1c'),
            ('64 of 95 Vendors', 10, False, '777777'), ('Flat-rate pricing, no outcome alignment', 9, False, '999999')
        ], 'center')

        # 1.63/5.0
        box(s8, sx8(310), sy8(135), sw8(250), sh8(145), 'fff9c4', 'ffd54f')
        ml(s8, sx8(320), sy8(142), sw8(230), sh8(135), [
            ('1.63 / 5.0', 22, True, 'd35400'), ('Outcome Pricing Score', 11, True, '8d4004'),
            ('86% score \u2264 2.0', 10, False, '777777'), ('Median is just 1.0', 9, False, '999999')
        ], 'center')

        # AI Dividend Captured
        box(s8, sx8(40), sy8(295), sw8(250), sh8(120), 'fff3e0', 'ffb74d')
        ml(s8, sx8(50), sy8(302), sw8(230), sh8(110), [
            ('AI Dividend Captured', 12, True, 'c0392b'),
            ('Vendors keep AI efficiency', 10, False, '777777'), ('gains as margin \u2014 not', 10, False, '777777'), ('sharing with buyers', 10, False, '777777')
        ], 'center')

        # Market Avg 2.35
        box(s8, sx8(310), sy8(295), sw8(250), sh8(120), 'e3f2fd', 'bbdefb')
        ml(s8, sx8(320), sy8(302), sw8(230), sh8(110), [
            ('2.35 / 5.0', 20, True, '7d3c98'), ('Market Average', 12, True, '4a148c'),
            ('Overall maturity:', 10, False, '777777'), ('"Developing" stage', 10, False, '777777')
        ], 'center')

        # CENTER — AI Pricing Catalyst
        box(s8, sx8(530), sy8(440), sw8(200), sh8(140), 'e8f4fd', '2471a3')
        ml(s8, sx8(535), sy8(450), sw8(190), sh8(125), [
            ('AI Pricing', 16, True, '1a3a5c'), ('Catalyst', 16, True, '1a3a5c'),
            ('', 6, False, None),
            ('1.02-point gap', 11, True, '2471a3'), ('between AI-mature &', 9, False, '555555'), ('minimal vendors', 9, False, '555555')
        ], 'center')

        # Composable Models
        box(s8, sx8(530), sy8(600), sw8(200), sh8(85), 'e8f5e9', '81c784')
        ml(s8, sx8(535), sy8(608), sw8(190), sh8(75), [
            ('Composable Models', 12, True, '1e8449'), ('2.80 vs 2.16', 16, True, '1e8449'),
            ('Modular beats sub-only 30%', 9, False, '555555')
        ], 'center')

        # RIGHT ZONE — future model header
        box(s8, sx8(690), sy8(90), sw8(540), sh8(32), 'd4eaf7', '2471a3')
        tb(s8, sx8(700), sy8(93), sw8(520), sh8(28), '\u2726 AI-DRIVEN: OUTCOME-BASED FUTURE', 11, True, '2471a3', 'center')

        # 40% by 2028
        box(s8, sx8(700), sy8(135), sw8(250), sh8(145), 'e3f2fd', '90caf9')
        ml(s8, sx8(710), sy8(142), sw8(230), sh8(135), [
            ('40%', 28, True, '2471a3'), ('Outcome-Linked', 12, True, '2471a3'),
            ('by 2028', 11, True, '1a3a5c'),
            ('MDR contracts with at least', 9, False, '555555'), ('one outcome component', 9, False, '555555')
        ], 'center')

        # AI-Mature 2.86
        box(s8, sx8(970), sy8(135), sw8(250), sh8(145), 'e8f5e9', 'a5d6a7')
        ml(s8, sx8(980), sy8(142), sw8(230), sh8(135), [
            ('2.86', 28, True, '1e8449'), ('AI-Mature Average', 11, True, '1e8449'),
            ('vs 1.84 for AI-Minimal:', 10, False, '555555'), ('55% higher pricing', 10, False, '555555'), ('sophistication', 10, False, '555555')
        ], 'center')

        # Breach Warranties
        box(s8, sx8(700), sy8(295), sw8(250), sh8(120), 'fff9c4', 'fff59d')
        ml(s8, sx8(710), sy8(302), sw8(230), sh8(110), [
            ('Breach Warranties', 12, True, '7d3c98'),
            ('Performance credits,', 10, False, '555555'), ('risk-score pricing,', 10, False, '555555'), ('shared savings models', 10, False, '555555')
        ], 'center')

        # AI Efficiency Sharing
        box(s8, sx8(970), sy8(295), sw8(250), sh8(120), 'fff3e0', 'ffe0b2')
        ml(s8, sx8(980), sy8(302), sw8(230), sh8(110), [
            ('AI Efficiency Sharing', 11, True, '2471a3'),
            ('Cost-per-incident trending,', 10, False, '555555'), ('AI utilization dashboards,', 10, False, '555555'), ('graduated discounts', 10, False, '555555')
        ], 'center')

        # Stat callouts
        stats8 = [
            (50, sy8(430), 160, 55, '86%', 'Score \u2264 2.0', 'c0392b', 'fce8d8'),
            (230, sy8(440), 140, 55, '19', 'AI-Significant Vendors', '2471a3', 'dbeeff'),
            (60, sy8(500), 140, 50, '0', 'Transformative-Tier', '7d3c98', 'f0e0ff'),
            (sx8(1050), sy8(430), 160, 55, '95', 'Vendors Assessed', '1e8449', 'd4edda'),
            (sx8(880), sy8(440), 150, 55, '3.52', 'Best: Sub Transparency', '2471a3', 'dbeeff'),
        ]
        for bx, by, bw, bh, bv, bl, bc, bbg in stats8:
            box(s8, bx, by, bw, bh, bbg, bc)
            ml(s8, bx + 5, by + 5, bw - 10, bh - 10, [(bv, 18, True, bc), (bl, 8, False, '5a503c')], 'center')

        # Bottom comparison bars
        box(s8, sx8(20), sy8(710), sw8(1220), sh8(200), 'f8f8f5', 'e0ddd5')
        tb(s8, sx8(50), sy8(718), sw8(200), sh8(20), 'SUBSCRIPTION', 11, True, 'c0392b')
        tb(s8, sx8(560), sy8(718), sw8(140), sh8(20), 'vs.', 12, True, '888888', 'center')
        tb(s8, sx8(960), sy8(718), sw8(250), sh8(20), 'OUTCOME-BASED', 11, True, '2471a3', 'right')
        bars8 = [
            (750, 'Revenue: Flat-rate, predictable, opaque', 'Revenue: Performance-linked, transparent'),
            (790, 'AI Efficiency: Captured as vendor margin', 'AI Efficiency: Shared with buyers'),
            (830, 'Risk: Buyer bears all risk', 'Risk: Vendor co-owns via breach warranties'),
            (870, 'Pricing Score: 2.16 avg (sub-only)', 'Pricing Score: 2.80 avg (composable)'),
        ]
        for y8, left_t, right_t in bars8:
            box(s8, sx8(50), sy8(y8), sw8(475), sh8(28), 'fce8d8', 'c0392b')
            tb(s8, sx8(55), sy8(y8 + 3), sw8(465), sh8(24), left_t, 9, False, 'c0392b', 'center')
            tb(s8, sx8(535), sy8(y8 + 1), sw8(190), sh8(24), '\u2192', 13, True, '1e8449', 'center')
            box(s8, sx8(735), sy8(y8), sw8(475), sh8(28), 'd8eeff', '2471a3')
            tb(s8, sx8(740), sy8(y8 + 3), sw8(465), sh8(24), right_t, 9, False, '2471a3', 'center')

        tb(s8, sx8(850), sy8(1050), sw8(380), sh8(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 9: Outcome Metric Framework — Cross-Section Matrix ──
        s9 = prs.slides.add_slide(BLANK)
        box(s9, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s9, sx(50), sy(6), sw(1100), sh(26), '9. Outcome Metric Framework \u2014 Cross-Section Matrix', 20, True, 'ffffff', 'center')
        tb(s9, sx(100), sy(32), sw(1000), sh(20), 'Six AI-enabled metric categories mapped to pricing mechanisms, market readiness, and deployment sequencing', 10, False, 'ccd8e8', 'center')

        # Table header row
        col_x = [0, 120, 310, 530, 735, 940]
        col_w = [120, 190, 220, 205, 205, 260]
        hdr_y = 65
        hdr_h = 38
        headers9 = ['METRIC\nCATEGORY', 'KEY\nMEASURES', 'PRICING MECHANISM\nUNLOCKED', 'AI ENABLEMENT\nFACTORS', 'MARKET\nREADINESS', 'DEPLOYMENT\nPHASE']
        for ci, hdr_text in enumerate(headers9):
            box(s9, sx(col_x[ci]), sy(hdr_y), sw(col_w[ci]), sh(hdr_h), '1a3a5c', '1a3a5c')
            tb(s9, sx(col_x[ci] + 4), sy(hdr_y + 4), sw(col_w[ci] - 8), sh(hdr_h - 8), hdr_text, 8, True, 'ffffff', 'center')

        # Metric data rows
        metrics9 = [
            {
                'cat': '\u26a1 Speed', 'cat_col': '107c10',
                'measures': 'MTTD reduction %\nMTTR reduction %\nSeverity-weighted RT',
                'pricing': 'MTTR/MTTD Guarantees\nwith Service Credits',
                'ai': 'Pipeline instrumentation\nAuto-timestamping\nPer-incident classification',
                'readiness': 'High', 'rd_col': '107c10', 'rd_bg': 'e6f4e6',
                'rd_note': '~16% of vendors publish\nspeed data',
                'phase': 'Phase 1', 'phase_time': 'Months 1\u20133',
                'ref': 'S3: Speed Metric Trap',
            },
            {
                'cat': '\U0001f916 Autonomy', 'cat_col': '0078d4',
                'measures': '% incidents auto-resolved\nAI vs human attribution\nCost-per-incident (AI/analyst)',
                'pricing': 'Efficiency-Sharing &\nCost-Per-Incident',
                'ai': 'Triage-to-containment auto\nResolution path tagging\nHuman/machine attribution',
                'readiness': 'Low', 'rd_col': 'ca5010', 'rd_bg': 'fff0e0',
                'rd_note': '<10% track per-customer\nautonomy rates',
                'phase': 'Phase 2', 'phase_time': 'Months 3\u20136',
                'ref': 'S4: Autonomy as Currency',
            },
            {
                'cat': '\U0001f6e1 Coverage', 'cat_col': '8764b8',
                'measures': 'Telemetry source breadth\nAttack surface mapped %\nVisibility gap trending',
                'pricing': 'Risk-Score-Tied\nSubscriptions',
                'ai': 'Automated asset discovery\nDynamic coverage scoring\nTelemetry gap analysis',
                'readiness': 'Medium', 'rd_col': 'ca5010', 'rd_bg': 'fff0e0',
                'rd_note': '28% composable lead;\n66% treat as static',
                'phase': 'Phase 3', 'phase_time': 'Months 6\u201312',
                'ref': 'S5: Coverage & Quality',
            },
            {
                'cat': '\U0001f3af Quality', 'cat_col': 'e3008c',
                'measures': 'False positive rate trend\nDetection confidence\nAI hallucination rate',
                'pricing': 'Detection Confidence\nWarranties',
                'ai': 'Per-alert confidence scoring\nQuality trending per cust.\nHallucination tracking',
                'readiness': 'Low', 'rd_col': 'ca5010', 'rd_bg': 'fff0e0',
                'rd_note': '~19% AI-mature have infra\nbut don\'t commercialize',
                'phase': 'Phase 3', 'phase_time': 'Months 6\u201312',
                'ref': 'S5: Coverage & Quality',
            },
            {
                'cat': '\U0001f464 Resource', 'cat_col': '00b7c3',
                'measures': 'Analyst time saved (hrs)\nSOC workload shift\nPre/post-AI baseline',
                'pricing': 'Analyst-Time-Saved\nRebates',
                'ai': 'SOC integration telemetry\nTime-per-incident measure\nWorkload attribution models',
                'readiness': 'Very Low', 'rd_col': 'a80000', 'rd_bg': 'ffe0e0',
                'rd_note': 'Requires customer SOC\nintegration',
                'phase': 'Phase 4', 'phase_time': '12+ months',
                'ref': 'S6: Metrics to Money',
            },
            {
                'cat': '\U0001f4b0 Impact', 'cat_col': 'a80000',
                'measures': 'Damage prevention ($)\nRegulatory penalty avoid.\nRisk score improvement',
                'pricing': 'Breach Warranties &\nDamage-Cap Guarantees',
                'ai': 'Predictive risk scoring\nAttack path analysis\nActuarial breach modeling',
                'readiness': 'Very Low', 'rd_col': 'a80000', 'rd_bg': 'ffe0e0',
                'rd_note': '~2% of vendors approach\noutcome maturity 4',
                'phase': 'Phase 4', 'phase_time': '12+ months',
                'ref': 'S6: Metrics to Money',
            },
        ]

        row_h = 80
        row_y0 = hdr_y + hdr_h
        for ri, m in enumerate(metrics9):
            ry = row_y0 + ri * row_h
            row_bg = 'fafaf8' if ri % 2 == 0 else 'f5f3f0'
            for ci in range(6):
                box(s9, sx(col_x[ci]), sy(ry), sw(col_w[ci]), sh(row_h), row_bg, 'e0ddd5')

            # Category
            ml(s9, sx(col_x[0] + 6), sy(ry + 6), sw(col_w[0] - 12), sh(row_h - 12), [
                (m['cat'], 12, True, m['cat_col']),
                (m['ref'], 7, False, '888888'),
            ])
            # Measures
            tb(s9, sx(col_x[1] + 4), sy(ry + 4), sw(col_w[1] - 8), sh(row_h - 8), m['measures'], 8, False, '333333')
            # Pricing mechanism unlock
            box(s9, sx(col_x[2] + 6), sy(ry + 8), sw(col_w[2] - 12), sh(row_h - 16), None, m['cat_col'])
            tb(s9, sx(col_x[2] + 10), sy(ry + 12), sw(col_w[2] - 20), sh(row_h - 24), m['pricing'], 9, True, m['cat_col'], 'center')
            # AI enablers
            tb(s9, sx(col_x[3] + 4), sy(ry + 4), sw(col_w[3] - 8), sh(row_h - 8), m['ai'], 8, False, '333333')
            # Readiness
            box(s9, sx(col_x[4] + 15), sy(ry + 6), sw(col_w[4] - 30), sh(22), m['rd_bg'], m['rd_col'])
            tb(s9, sx(col_x[4] + 15), sy(ry + 7), sw(col_w[4] - 30), sh(20), m['readiness'], 9, True, m['rd_col'], 'center')
            tb(s9, sx(col_x[4] + 6), sy(ry + 32), sw(col_w[4] - 12), sh(row_h - 36), m['rd_note'], 7, False, '555555', 'center')
            # Phase
            ml(s9, sx(col_x[5] + 4), sy(ry + 10), sw(col_w[5] - 8), sh(row_h - 12), [
                (m['phase'], 11, True, m['cat_col']),
                (m['phase_time'], 9, False, '555555'),
            ], 'center')

        # Deployment Sequence — 4 phase cards below table
        phase_y = row_y0 + 6 * row_h + 10
        phase_w = 275
        phase_gap = 13
        phases9 = [
            ('Phase 1', 'Months 1\u20133', 'Speed Metrics', '107c10', 'e6f4e6', '\u2192 MTTR guarantees'),
            ('Phase 2', 'Months 3\u20136', 'Autonomy Metrics', '0078d4', 'e0f0ff', '\u2192 Efficiency-sharing'),
            ('Phase 3', 'Months 6\u201312', 'Coverage + Quality', '8764b8', 'f0e8ff', '\u2192 Risk-score + warranties'),
            ('Phase 4', '12+ months', 'Resource + Impact', 'a80000', 'ffe0e0', '\u2192 Rebates + breach warranties'),
        ]
        for pi, (pl, pt, ptitle, pc, pbg, punlock) in enumerate(phases9):
            px = 12 + pi * (phase_w + phase_gap)
            box(s9, sx(px), sy(phase_y), sw(phase_w), sh(70), pbg, pc)
            ml(s9, sx(px + 5), sy(phase_y + 4), sw(phase_w - 10), sh(62), [
                (pl + ': ' + ptitle, 10, True, pc),
                (pt, 8, False, '555555'),
                (punlock, 8, True, pc),
            ], 'center')

        # Key Insight box
        insight_y = phase_y + 78
        box(s9, sx(0), sy(insight_y), sw(1200), sh(55), 'f0faff', '0078d4')
        ml(s9, sx(10), sy(insight_y + 4), sw(1180), sh(48), [
            ('Key Insight: Speed metrics are the foundation (measurable now), Autonomy metrics are the AI differentiator (highest commercial value),', 8, False, '333333'),
            ('Coverage + Quality are retention drivers. The 1.47-point outcome maturity gap between AI-Significant (2.58) and AI-Minimal (1.11) confirms AI enables measurement.', 8, False, '333333'),
        ])

        tb(s9, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='MDR_All_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/precyber-all-graphics-pptx', methods=['GET'])
def precyber_all_graphics_pptx():
    """Generate a multi-slide editable PowerPoint deck with ALL 7 PreCyber graphics."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from flask import send_file
        import io

        vendor_file = _get_precyber_vendor_file()
        if not vendor_file:
            return jsonify({'error': 'PreCyber vendor data not found'}), 404
        st = _compute_precyber_stats(vendor_file)
        n = st['vendor_count']
        pp = st['pillar_penetration']
        dm = st['delivery_models']
        fs_count = st['full_spectrum_count']
        fs_pct = st['full_spectrum_pct']
        maj_count = st['majority_spectrum_count']
        maj_pct = st['majority_spectrum_pct']
        narrow_count = st['narrow_count']
        narrow_pct = st['narrow_pct']
        blind_pct = st['blind_spot_pct']
        no_amt_pct = st['no_amt_pct']
        dm_ds = dm.get('direct_service', {})
        dm_pp = dm.get('platform_plus_partner', {})
        dm_po = dm.get('platform_only', {})
        ds_count = dm_ds.get('count', 0)
        ppp_count = dm_pp.get('count', 0)
        po_count = dm_po.get('count', 0)
        ds_pct = round(ds_count * 100 / n) if n else 0
        ppp_pct = round(ppp_count * 100 / n) if n else 0
        po_pct = round(po_count * 100 / n) if n else 0
        ds_pa = dm_ds.get('pillar_avgs', {})
        ppp_pa = dm_pp.get('pillar_avgs', {})
        po_pa = dm_po.get('pillar_avgs', {})
        ds_overall = dm_ds.get('overall_avg', 0)
        ppp_overall = dm_pp.get('overall_avg', 0)
        po_overall = dm_po.get('overall_avg', 0)
        po_svc_below_pct = dm_po.get('pillar_below_pct', {}).get('SVC', 0)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        def box(slide, l, t, w, h, fill=None, border=None):
            s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
            if fill:
                s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
            else:
                s.fill.background()
            if border:
                s.line.color.rgb = rgb(border); s.line.width = Pt(1.5)
            else:
                s.line.fill.background()
            return s

        def tb(slide, l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = text
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = rgb(col)
            if al == 'center': p.alignment = PP_ALIGN.CENTER
            elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        def ml(slide, l, t, w, h, lines, al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(11)
                p.font.bold = bld
                if col: p.font.color.rgb = rgb(col)
                if al == 'center': p.alignment = PP_ALIGN.CENTER
                elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        SW = Inches(13.333); SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 750 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 750 * SH)
        footer = '\u00a9 Gartner Research \u2022 PreCyber Market Insight 2026'

        # ── SLIDE 1: Five-Pillar Coverage Heatmap ──
        s1 = prs.slides.add_slide(BLANK)
        box(s1, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s1, sx(50), sy(8), sw(1100), sh(28), '1. Five-Pillar Coverage Heatmap by Delivery Model', 22, True, 'ffffff', 'center')
        tb(s1, sx(100), sy(36), sw(1000), sh(22), f'{n} preemptive cybersecurity vendors scored across 5 pillars by delivery model', 11, False, 'ccd8e8', 'center')

        pillars = [
            ('Exposure Management', 'EXM', f"{pp['EXM']['pct']}%", ds_pa.get('EXM', 0), ppp_pa.get('EXM', 0), po_pa.get('EXM', 0)),
            ('Posture & Policy Mgmt', 'PPM', f"{pp['PPM']['pct']}%", ds_pa.get('PPM', 0), ppp_pa.get('PPM', 0), po_pa.get('PPM', 0)),
            ('Detection & Response', 'ADR', f"{pp['ADR']['pct']}%", ds_pa.get('ADR', 0), ppp_pa.get('ADR', 0), po_pa.get('ADR', 0)),
            ('Services & Capability', 'SVC', f"{pp['SVC']['pct']}%", ds_pa.get('SVC', 0), ppp_pa.get('SVC', 0), po_pa.get('SVC', 0)),
            ('Adversary Management', 'AMT', f"{pp['AMT']['pct']}%", ds_pa.get('AMT', 0), ppp_pa.get('AMT', 0), po_pa.get('AMT', 0)),
        ]
        headers = [('Penetration', '1a3a5c'), (f'Direct Service\n{ds_count} vendors ({ds_pct}%)', '107c10'), (f'Platform+Partner\n{ppp_count} vendors ({ppp_pct}%)', '0078d4'), (f'Platform-Only\n{po_count} vendors ({po_pct}%)', 'ca5010')]
        colW = 195; startX = 340; rowH = 60
        for ci, (ht, hc) in enumerate(headers):
            box(s1, sx(startX + ci * colW), sy(75), sw(colW - 10), sh(50), hc, hc)
            tb(s1, sx(startX + ci * colW + 5), sy(78), sw(colW - 20), sh(46), ht, 10, True, 'ffffff', 'center')

        def pc_color(v):
            if v >= 3.0: return '107c10'
            if v >= 2.5: return '0078d4'
            if v >= 2.0: return 'ca5010'
            return 'a80000'
        def pc_bg(v):
            if v >= 3.0: return 'e6f4e6'
            if v >= 2.5: return 'e0f0ff'
            if v >= 2.0: return 'fff0e0'
            return 'ffe0e0'

        for ri, (name, code, pct, direct, partner, plat) in enumerate(pillars):
            y = 135 + ri * rowH
            box(s1, sx(30), sy(y), sw(300), sh(rowH - 5), 'f8f8f5', 'e0ddd5')
            ml(s1, sx(40), sy(y + 5), sw(280), sh(40), [(name, 12, True, '333333'), (code, 9, False, '888888')])
            # Penetration
            pct_val = int(pct.rstrip('%'))
            pct_col = '107c10' if pct_val >= 80 else 'ca5010' if pct_val >= 60 else 'a80000'
            box(s1, sx(startX), sy(y), sw(colW - 10), sh(rowH - 5), 'f8f8f5', pct_col)
            tb(s1, sx(startX + 5), sy(y + 10), sw(colW - 20), sh(32), pct, 18, True, pct_col, 'center')
            for ci, v in enumerate([direct, partner, plat]):
                box(s1, sx(startX + (ci + 1) * colW), sy(y), sw(colW - 10), sh(rowH - 5), pc_bg(v), pc_color(v))
                tb(s1, sx(startX + (ci + 1) * colW + 5), sy(y + 10), sw(colW - 20), sh(32), f'{v:.2f}', 18, True, pc_color(v), 'center')

        legend_y = 135 + 5 * rowH + 15
        for li, (lt, lc) in enumerate([('\u25a0 \u2265 3.0 Strong', '107c10'), ('\u25a0 2.5\u20132.99 Competitive', '0078d4'), ('\u25a0 2.0\u20132.49 Below Target', 'ca5010'), ('\u25a0 < 2.0 Critical', 'a80000')]):
            tb(s1, sx(40 + li * 270), sy(legend_y), sw(260), sh(22), lt, 11, True, lc)
        tb(s1, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 2: Delivery Model Maturity Roadmap ──
        s2 = prs.slides.add_slide(BLANK)
        box(s2, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s2, sx(50), sy(8), sw(1100), sh(28), '2. Path to Full-Spectrum: Delivery Model Maturity', 22, True, 'ffffff', 'center')
        tb(s2, sx(100), sy(36), sw(1000), sh(22), 'Current average score per delivery model relative to 2.0 baseline competency threshold', 11, False, 'ccd8e8', 'center')

        models = [
            ('Direct Service Providers', f'{ds_count} vendors ({ds_pct}%)', ds_overall, '107c10',
             [(f'\u2713 Highest SVC: {ds_pa.get("SVC", 0):.2f}', False), (f'\u2713 Highest ADR: {ds_pa.get("ADR", 0):.2f}', False), ('\u2713 Own SOCs + analysts', False), ('\u26a0 Limited platform depth', True)]),
            ('Platform + Partner', f'{ppp_count} vendors ({ppp_pct}%)', ppp_overall, '0078d4',
             [(f'\u2713 Highest AMT: {ppp_pa.get("AMT", 0):.2f}', False), (f'\u2713 Highest PPM: {ppp_pa.get("PPM", 0):.2f}', False), ('\u2713 Broadest coverage', False), ('\u26a0 Partner accountability gaps', True), (f'\u26a0 SVC via partners: {ppp_pa.get("SVC", 0):.2f}', True)]),
            ('Platform-Only', f'{po_count} vendors ({po_pct}%)', po_overall, 'ca5010',
             [(f'\u26a0 SVC: {po_pa.get("SVC", 0):.2f} (structural gap)', True), (f'\u26a0 AMT: {po_pa.get("AMT", 0):.2f} (weak intel)', True), (f'\u26a0 {po_svc_below_pct}% below 2.0 on SVC', True), ('\u26a0 No service delivery', True)]),
        ]
        for mi_, (name, sub, score, color, gaps) in enumerate(models):
            baseY = 80 + mi_ * 200
            tb(s2, sx(40), sy(baseY), sw(250), sh(24), name, 16, True, color)
            tb(s2, sx(40), sy(baseY + 26), sw(250), sh(18), f'{sub} \u2022 {score:.2f} avg', 10, False, '666666')
            barW = 880
            box(s2, sx(260), sy(baseY + 5), sw(barW), sh(38), 'e8e8e5', 'cccccc')
            fillW = int(barW * score / 5.0)
            box(s2, sx(260), sy(baseY + 5), sw(fillW), sh(38), color, color)
            tb(s2, sx(260 + fillW // 2 - 50), sy(baseY + 10), sw(100), sh(28), f'{score:.2f} / 5.00', 13, True, 'ffffff', 'center')
            targetX = int(barW * 2.0 / 5.0)
            box(s2, sx(260 + targetX - 1), sy(baseY + 2), sw(3), sh(44), 'a80000', 'a80000')
            tb(s2, sx(260 + targetX - 30), sy(baseY + 48), sw(60), sh(16), '2.0 Baseline', 8, True, 'a80000', 'center')
            for gi, (gtxt, critical) in enumerate(gaps):
                gx = 260 + (gi % 3) * 300
                gy = baseY + 72 + (gi // 3) * 30
                gc = 'a80000' if critical else '107c10'
                gb = 'fff0f0' if critical else 'f0fff0'
                box(s2, sx(gx), sy(gy), sw(280), sh(24), gb, gc)
                tb(s2, sx(gx + 5), sy(gy + 2), sw(270), sh(20), gtxt, 10, False, gc, 'center')

        box(s2, sx(40), sy(680), sw(1120), sh(50), 'fff5f5', 'a80000')
        ml(s2, sx(55), sy(683), sw(1090), sh(44), [
            ('Key Insight', 11, True, 'a80000'),
            ('No platform-only vendor achieves full-spectrum coverage. SVC is the structural missing link \u2014 88% of platform-only vendors score below 2.0.', 10, False, '555555'),
        ])
        tb(s2, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 3: Pillar Penetration Gap ──
        s3 = prs.slides.add_slide(BLANK)
        box(s3, sx(0), sy(0), sw(1200), sh(60), '1a3a5c', '1a3a5c')
        tb(s3, sx(50), sy(8), sw(1100), sh(28), '3. Pillar Penetration Gap: Where the Market Falls Short', 22, True, 'ffffff', 'center')
        tb(s3, sx(100), sy(36), sw(1000), sh(22), f'Score by delivery model across all five pillars; {pp["EXM"]["pct"] - pp["AMT"]["pct"]}-point gap between EXM ({pp["EXM"]["pct"]}%) and AMT ({pp["AMT"]["pct"]}%)', 11, False, 'ccd8e8', 'center')

        tb(s3, sx(60), sy(72), sw(160), sh(18), '\u25a0 Direct Service', 11, True, '107c10')
        tb(s3, sx(240), sy(72), sw(160), sh(18), '\u25a0 Platform+Partner', 11, True, '0078d4')
        tb(s3, sx(420), sy(72), sw(160), sh(18), '\u25a0 Platform-Only', 11, True, 'ca5010')

        gap_data = [
            ('Exposure Mgmt (EXM)', f"{pp['EXM']['pct']}%", ds_pa.get('EXM', 0), ppp_pa.get('EXM', 0), po_pa.get('EXM', 0)),
            ('Posture & Policy (PPM)', f"{pp['PPM']['pct']}%", ds_pa.get('PPM', 0), ppp_pa.get('PPM', 0), po_pa.get('PPM', 0)),
            ('Adversary Disruption (ADR)', f"{pp['ADR']['pct']}%", ds_pa.get('ADR', 0), ppp_pa.get('ADR', 0), po_pa.get('ADR', 0)),
            ('Services & Capability (SVC)', f"{pp['SVC']['pct']}%", ds_pa.get('SVC', 0), ppp_pa.get('SVC', 0), po_pa.get('SVC', 0)),
            ('Adversary Mgmt (AMT)', f"{pp['AMT']['pct']}%", ds_pa.get('AMT', 0), ppp_pa.get('AMT', 0), po_pa.get('AMT', 0)),
        ]
        for gi, (name, pct, direct, partner, plat) in enumerate(gap_data):
            baseY = 100 + gi * 115
            pct_val = int(pct.rstrip('%'))
            pct_col = '107c10' if pct_val >= 80 else 'ca5010' if pct_val >= 60 else 'a80000'
            tb(s3, sx(40), sy(baseY), sw(300), sh(22), name, 14, True, '333333')
            tb(s3, sx(350), sy(baseY), sw(80), sh(22), f'{pct} pen.', 12, True, pct_col, 'right')
            barW = 600; labW = 110
            for bi, (label, val, clr) in enumerate([('Direct', direct, '107c10'), ('Plat+Partner', partner, '0078d4'), ('Platform-Only', plat, 'ca5010')]):
                by_ = baseY + 28 + bi * 26
                tb(s3, sx(40), sy(by_), sw(labW), sh(18), label, 10, False, '777777')
                box(s3, sx(40 + labW), sy(by_ - 2), sw(barW), sh(22), 'f0f0ee', 'dddddd')
                box(s3, sx(40 + labW), sy(by_ - 2), sw(int(barW * val / 5.0)), sh(22), clr, clr)
                tb(s3, sx(40 + labW + barW + 10), sy(by_), sw(60), sh(18), f'{val:.2f}', 12, True, clr)

        box(s3, sx(40), sy(680), sw(1120), sh(50), 'f0faff', '0078d4')
        ml(s3, sx(55), sy(683), sw(1090), sh(44), [
            ('Key Insight', 11, True, '0078d4'),
            ('37-point penetration gap between EXM (92%) and AMT (55%). Vendors cluster around exposure scanning but ignore adversary intelligence and service delivery.', 10, False, '555555'),
        ])
        tb(s3, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 4: Executive Summary Poster ──
        s4 = prs.slides.add_slide(BLANK)
        box(s4, sx(0), sy(0), sw(1200), sh(55), 'a80000', 'a80000')
        tb(s4, sx(50), sy(5), sw(1100), sh(28), '4. The Preemptive Cybersecurity Market Is Dangerously Fragmented', 20, True, 'ffffff', 'center')
        tb(s4, sx(100), sy(32), sw(1000), sh(20), 'Market Insight \u2014 51 vendors across 5 capability pillars \u00d7 3 delivery models', 11, False, 'ffcccc', 'center')

        stats_s4 = [(str(n), 'Vendors Assessed', '0078d4'), (f'{blind_pct}%', '\u2265 1 Blind Spot', 'a80000'), (f'{no_amt_pct}%', 'No AMT Capability', 'ca5010'), (f'{fs_pct}%', 'Full-Spectrum', '107c10'), (f'{po_pct}%', 'Platform-Only', '8764b8')]
        for si, (sv, sl, sc) in enumerate(stats_s4):
            bx = 40 + si * 232
            box(s4, sx(bx), sy(65), sw(215), sh(60), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(bx + 5), sy(68), sw(205), sh(28), sv, 22, True, sc, 'center')
            tb(s4, sx(bx + 5), sy(98), sw(205), sh(22), sl, 10, False, '666666', 'center')

        tb(s4, sx(40), sy(140), sw(300), sh(22), 'KEY FINDINGS', 13, True, 'a80000')
        findings = [
            ('\U0001f534 Service Delivery Crisis', f'{po_count} of {n} vendors ({po_pct}%) are platform-only. Avg SVC: {po_pa.get("SVC", 0):.2f}. {po_svc_below_pct}% of platform-only vendors score below 2.0 on services.'),
            ('\u26a0 Adversary Intelligence Deficit', f'{no_amt_pct}% of vendors lack AMT capability. Among platform-only: {dm_po.get("pillar_below_pct", {}).get("AMT", 0)}% score below 2.0 on adversary management.'),
            ('\u2713 Rare Full-Spectrum Excellence', 'Only 3 vendors maintain min score \u2265 2.5 across all 5 pillars: Mandiant, SentinelOne, Fortinet.'),
        ]
        for fi, (ft, fb) in enumerate(findings):
            fy = 168 + fi * 75
            box(s4, sx(40), sy(fy), sw(1120), sh(65), 'f8f8f5', 'e0ddd5')
            tb(s4, sx(55), sy(fy + 5), sw(1090), sh(22), ft, 13, True, '333333')
            tb(s4, sx(55), sy(fy + 28), sw(1090), sh(34), fb, 11, False, '555555')

        tb(s4, sx(40), sy(400), sw(400), sh(22), 'RECOMMENDATIONS', 13, True, '0078d4')
        recs = [
            ('01', 'Map Gaps First', 'Prioritize AMT & SVC investments', 'a80000'),
            ('02', 'Build Service Layer', 'Platform vendors: build / buy / partner', '0078d4'),
            ('03', 'M&A for Speed', 'Acquire full-spectrum positioning', '107c10'),
            ('04', 'Outcome Pricing', 'Tie pricing to preemptive results', '8764b8'),
        ]
        for ri, (rn, rt, rd, rc) in enumerate(recs):
            rx = 40 + ri * 285
            box(s4, sx(rx), sy(430), sw(260), sh(85), 'f8f8f5', rc)
            tb(s4, sx(rx + 10), sy(433), sw(240), sh(16), rn, 10, True, rc)
            tb(s4, sx(rx + 10), sy(450), sw(240), sh(22), rt, 14, True, rc, 'center')
            tb(s4, sx(rx + 10), sy(475), sw(240), sh(34), rd, 10, False, '555555', 'center')

        box(s4, sx(40), sy(540), sw(1120), sh(40), 'fff0f0', 'a80000')
        tb(s4, sx(55), sy(545), sw(1090), sh(32), '\u26a0 Service delivery (SVC) is the missing link: no platform-only vendor achieves full-spectrum coverage.', 12, True, 'a80000', 'center')
        tb(s4, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 5: Delivery Model Trade-offs ──
        s5 = prs.slides.add_slide(BLANK)
        box(s5, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s5, sx(50), sy(5), sw(1100), sh(28), '5. No Single Delivery Model Solves Full-Spectrum', 20, True, 'ffffff', 'center')
        tb(s5, sx(100), sy(32), sw(1000), sh(20), 'Three models, three trade-offs \u2014 accountability vs breadth vs technical depth', 11, False, 'ccd8e8', 'center')

        model_cards = [
            ('\U0001f7e2 Direct Service', f'{ds_count} vendors ({ds_pct}%)', '107c10', 'f0fff4',
             ['\u2713 Own SOC + analyst teams', '\u2713 Single accountability point', f'\u2713 Highest SVC: {ds_pa.get("SVC", 0):.2f}', f'\u2713 Highest ADR: {ds_pa.get("ADR", 0):.2f}', '\u26a0 Limited platform depth'], 'Operational Accountability'),
            ('\U0001f535 Platform + Partner', f'{ppp_count} vendors ({ppp_pct}%)', '0078d4', 'f0faff',
             ['\u2713 Tech platform + MSSP delivery', '\u2713 Broadest pillar coverage', f'\u2713 Highest AMT: {ppp_pa.get("AMT", 0):.2f}', f'\u2713 Highest PPM: {ppp_pa.get("PPM", 0):.2f}', '\u26a0 Partner accountability gaps'], 'Breadth of Coverage'),
            ('\U0001f7e0 Platform-Only', f'{po_count} vendors ({po_pct}%)', 'ca5010', 'fff8f0',
             ['\u2713 Technology licensing model', f'\U0001f534 SVC: {po_pa.get("SVC", 0):.2f} (structural gap)', f'\U0001f534 AMT: {po_pa.get("AMT", 0):.2f} (weak intel)', f'\U0001f534 {po_svc_below_pct}% below 2.0 on SVC', '\U0001f534 No service delivery'], 'Structural Service Deficit'),
        ]
        for ci, (mt, ms, mc, mbg, items, strength) in enumerate(model_cards):
            cx = 30 + ci * 390
            box(s5, sx(cx), sy(70), sw(370), sh(310), mbg, mc)
            tb(s5, sx(cx + 10), sy(75), sw(350), sh(24), mt, 16, True, mc, 'center')
            tb(s5, sx(cx + 10), sy(100), sw(350), sh(18), ms, 11, False, '666666', 'center')
            for ii, it in enumerate(items):
                tb(s5, sx(cx + 20), sy(128 + ii * 32), sw(330), sh(28), it, 12, False, '333333')
            box(s5, sx(cx + 15), sy(295), sw(340), sh(30), mc, mc)
            tb(s5, sx(cx + 20), sy(298), sw(330), sh(26), strength, 11, True, 'ffffff', 'center')

        tb(s5, sx(40), sy(400), sw(500), sh(22), 'MARKET SPECTRUM SEGMENTATION', 13, True, '1a3a5c')
        segs = [
            (str(fs_count), f'Full-Spectrum ({fs_pct}%)', 'All 5 pillars \u2265 2.0', '107c10', 'e6f4e6'),
            (str(maj_count), f'Majority-Spectrum ({maj_pct}%)', '4 pillars \u2014 one from full', 'ca5010', 'fff0e0'),
            (str(narrow_count), f'Narrow-Spectrum ({narrow_pct}%)', '\u2264 3 pillars \u2014 niche specialists', 'a80000', 'ffe0e0'),
        ]
        for si, (sv, sl, sd, sc, sbg) in enumerate(segs):
            bx = 40 + si * 385
            box(s5, sx(bx), sy(430), sw(360), sh(85), sbg, sc)
            tb(s5, sx(bx + 5), sy(435), sw(350), sh(30), sv, 28, True, sc, 'center')
            tb(s5, sx(bx + 5), sy(468), sw(350), sh(18), sl, 13, True, sc, 'center')
            tb(s5, sx(bx + 5), sy(490), sw(350), sh(18), sd, 10, False, '555555', 'center')

        box(s5, sx(40), sy(530), sw(1120), sh(35), 'fff0f0', 'a80000')
        tb(s5, sx(50), sy(533), sw(1100), sh(28), f'{blind_pct}% of vendors have at least one structural blind spot. Buyers cannot assume any single vendor covers the full attack surface.', 11, True, 'a80000', 'center')
        tb(s5, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 6: 2030 Market Evolution ──
        s6 = prs.slides.add_slide(BLANK)
        box(s6, sx(0), sy(0), sw(1200), sh(55), '8764b8', '8764b8')
        tb(s6, sx(50), sy(5), sw(1100), sh(28), '6. From Fragmentation to Full-Spectrum: The 2030 Journey', 20, True, 'ffffff', 'center')
        tb(s6, sx(100), sy(32), sw(1000), sh(20), 'Full-spectrum vendors projected to grow from 27% to 50%+ by 2030', 11, False, 'e0d0ff', 'center')

        phases = [
            ('Phase 1\n2025\u201326', 'Platform\nExtension', 'M&A wave to fill pillar gaps', '0078d4'),
            ('Phase 2\n2026\u201328', 'Service Layer\nDevelopment', 'Managed services build-out', '107c10'),
            ('Phase 3\n2028\u201330', 'Outcome-Based\nDelivery', 'Full-spectrum outcome pricing', '8764b8'),
        ]
        for pi, (pn, pl, pd, pc) in enumerate(phases):
            px = 40 + pi * 380
            box(s6, sx(px), sy(70), sw(350), sh(120), pc, pc)
            tb(s6, sx(px + 10), sy(75), sw(330), sh(35), pn, 12, True, 'ffffff', 'center')
            tb(s6, sx(px + 10), sy(110), sw(330), sh(30), pl, 16, True, 'ffffff', 'center')
            tb(s6, sx(px + 10), sy(148), sw(330), sh(30), pd, 10, False, 'ffffff', 'center')
            if pi < 2:
                tb(s6, sx(px + 355), sy(115), sw(25), sh(28), '\u2192', 18, True, '888888', 'center')

        tb(s6, sx(40), sy(210), sw(500), sh(22), '2030 PROJECTIONS', 13, True, '8764b8')
        projections = [
            ('50%+', 'Full-Spectrum Vendors', 'Up from 27% today', '107c10'),
            ('$8B+', 'Addressable Market', 'Preemptive cybersecurity TAM', '0078d4'),
            ('30\u201340%', 'MSSP Premium', 'For full-spectrum managed delivery', '8764b8'),
        ]
        for si, (sv, sl, sd, sc) in enumerate(projections):
            bx = 40 + si * 385
            box(s6, sx(bx), sy(240), sw(360), sh(80), 'f8f8f5', sc)
            tb(s6, sx(bx + 5), sy(243), sw(350), sh(32), sv, 28, True, sc, 'center')
            tb(s6, sx(bx + 5), sy(278), sw(350), sh(18), sl, 12, True, '333333', 'center')
            tb(s6, sx(bx + 5), sy(298), sw(350), sh(16), sd, 10, False, '666666', 'center')

        tb(s6, sx(40), sy(340), sw(500), sh(22), 'WHAT MUST CHANGE', 13, True, 'a80000')
        actions = [
            ('\U0001f527 Platform Vendors Must Do:', 'Develop or acquire service delivery. No SVC = no full-spectrum = structural ceiling.', 'ca5010'),
            ('\U0001f3d7 Service Providers Must Do:', 'Invest in platform depth & adversary management technology. SOC teams alone aren\'t enough.', '107c10'),
            ('\U0001f91d MSSPs Must Do:', 'Position as the integration layer solving fragmentation. Multi-vendor orchestration = premium opportunity.', '0078d4'),
        ]
        for ai, (at, ad, ac) in enumerate(actions):
            ay = 368 + ai * 50
            box(s6, sx(40), sy(ay), sw(1120), sh(44), 'f8f8f5', ac)
            tb(s6, sx(55), sy(ay + 3), sw(350), sh(18), at, 12, True, ac)
            tb(s6, sx(410), sy(ay + 3), sw(740), sh(38), ad, 11, False, '555555')

        box(s6, sx(40), sy(530), sw(1120), sh(35), 'f5f0ff', '8764b8')
        tb(s6, sx(50), sy(533), sw(1100), sh(28), 'The preemptive cybersecurity market will consolidate around full-spectrum, outcome-based delivery by 2030.', 11, True, '8764b8', 'center')
        tb(s6, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 7: Hand-Drawn Fragmentation Infographic ──
        s7 = prs.slides.add_slide(BLANK)
        def sx7(v): return int(v / 1200 * SW)
        def sy7(v): return int(v / 920 * SH)
        def sw7(v): return int(v / 1200 * SW)
        def sh7(v): return int(v / 920 * SH)

        box(s7, sx7(0), sy7(0), sw7(1200), sh7(70), '1a3a5c', '1a3a5c')
        tb(s7, sx7(50), sy7(5), sw7(1100), sh7(32), 'The Preemptive Cybersecurity Market Is Dangerously Fragmented', 22, True, 'ffffff', 'center')
        tb(s7, sx7(100), sy7(38), sw7(1000), sh7(28), f'Only {fs_pct}% of {n} vendors achieve full-spectrum coverage across all five pillars', 11, False, 'ccd8e8', 'center')

        # Five pillars
        p_data = [
            ('EXM', 'Exposure Mgmt', f"{pp['EXM']['pct']}%", '107c10', 'f0fff4'),
            ('PPM', 'Posture & Policy', f"{pp['PPM']['pct']}%", '0078d4', 'f0faff'),
            ('ADR', 'Detection & Response', f"{pp['ADR']['pct']}%", '8764b8', 'f5f0ff'),
            ('SVC', 'Services & Capability', f"{pp['SVC']['pct']}%", 'ca5010', 'fff8f0'),
            ('AMT', 'Adversary Mgmt', f"{pp['AMT']['pct']}%", 'a80000', 'fff5f5'),
        ]
        for pi, (code, name, pct, color, bg_) in enumerate(p_data):
            px = 40 + pi * 224
            box(s7, sx7(px), sy7(85), sw7(210), sh7(100), bg_, color)
            ml(s7, sx7(px + 5), sy7(92), sw7(200), sh7(85), [
                (pct, 28, True, color),
                (code, 14, True, color),
                (name, 10, False, '555555'),
            ], 'center')

        # Delivery models
        dm_data = [
            ('\U0001f7e2 Direct Service', f'{ds_count} vendors ({ds_pct}%)', f'SVC avg: {ds_pa.get("SVC", 0):.2f} (highest)\nADR avg: {ds_pa.get("ADR", 0):.2f}', '107c10', 'f0fff4'),
            ('\U0001f535 Platform + Partner', f'{ppp_count} vendors ({ppp_pct}%)', f'AMT avg: {ppp_pa.get("AMT", 0):.2f} (highest)\nPPM avg: {ppp_pa.get("PPM", 0):.2f}', '0078d4', 'f0faff'),
            ('\U0001f7e0 Platform-Only', f'{po_count} vendors ({po_pct}%)', f'SVC avg: {po_pa.get("SVC", 0):.2f} (lowest)\nAMT avg: {po_pa.get("AMT", 0):.2f}', 'ca5010', 'fff8f0'),
        ]
        for di, (dt, ds, dd, dc, dbg) in enumerate(dm_data):
            dx = 40 + di * 390
            box(s7, sx7(dx), sy7(210), sw7(370), sh7(130), dbg, dc)
            tb(s7, sx7(dx + 10), sy7(215), sw7(350), sh7(22), dt, 14, True, dc, 'center')
            tb(s7, sx7(dx + 10), sy7(240), sw7(350), sh7(16), ds, 11, False, '666666', 'center')
            tb(s7, sx7(dx + 10), sy7(266), sw7(350), sh7(60), dd, 11, False, '333333', 'center')

        # Market segmentation
        seg_data = [
            ('Full Spectrum', f'{fs_count} vendors ({fs_pct}%)', 'All 5 pillars \u2265 2.0', '107c10', 'e6f4e6'),
            ('Majority Spectrum', f'{maj_count} vendors ({maj_pct}%)', '4 pillars covered', 'ca5010', 'fff0e0'),
            ('Narrow Spectrum', f'{narrow_count} vendors ({narrow_pct}%)', '\u2264 3 pillars', 'a80000', 'ffe0e0'),
        ]
        for si, (st, ss, sd, sc, sbg) in enumerate(seg_data):
            bx = 40 + si * 390
            box(s7, sx7(bx), sy7(370), sw7(370), sh7(90), sbg, sc)
            ml(s7, sx7(bx + 10), sy7(378), sw7(350), sh7(75), [
                (st, 16, True, sc),
                (ss, 11, False, '333333'),
                (sd, 10, False, '666666'),
            ], 'center')

        # Stat callouts
        stat_boxes = [
            (65, 490, 200, 70, f'{blind_pct}%', 'Have \u2265 1 Blind Spot', 'a80000', 'fff5f5'),
            (305, 490, 200, 70, f'{no_amt_pct}%', 'No AMT Capability', 'ca5010', 'fff8f0'),
            (545, 490, 200, 70, str(n), 'Vendors Assessed', '0078d4', 'f0faff'),
            (785, 490, 200, 70, f'{narrow_pct}%', 'Narrow (\u2264 3 pillars)', '8764b8', 'f5f0ff'),
        ]
        for bx, by, bw, bh, bv, bl, bc, bbg in stat_boxes:
            box(s7, sx7(bx), sy7(by), sw7(bw), sh7(bh), bbg, bc)
            ml(s7, sx7(bx + 5), sy7(by + 8), sw7(bw - 10), sh7(bh - 16), [(bv, 20, True, bc), (bl, 10, False, '555555')], 'center')

        # Key insight
        box(s7, sx7(40), sy7(585), sw7(1120), sh7(50), 'fff0f0', 'a80000')
        tb(s7, sx7(50), sy7(590), sw7(1100), sh7(40), '\u2192 Service delivery (SVC) is the missing link: no platform-only vendor achieves full-spectrum \u2190', 14, True, 'a80000', 'center')

        tb(s7, sx7(800), sy7(895), sw7(380), sh7(20), footer, 9, False, 'aaaaaa', 'right')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='PreCyber_All_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/precyber-killchain-pptx', methods=['GET'])
def precyber_killchain_pptx():
    """Generate a 7-slide editable PowerPoint deck for the Kill Chain Shift-Left report."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from flask import send_file
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        def box(slide, l, t, w, h, fill=None, border=None):
            s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
            if fill:
                s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
            else:
                s.fill.background()
            if border:
                s.line.color.rgb = rgb(border); s.line.width = Pt(1.5)
            else:
                s.line.fill.background()
            return s

        def tb(slide, l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = text
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = rgb(col)
            if al == 'center': p.alignment = PP_ALIGN.CENTER
            elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        def ml(slide, l, t, w, h, lines, al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            for i, (txt, sz, bld, col) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(sz) if sz else Pt(11)
                p.font.bold = bld
                if col: p.font.color.rgb = rgb(col)
                if al == 'center': p.alignment = PP_ALIGN.CENTER
                elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        SW = Inches(13.333); SH = Inches(7.5)
        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 750 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 750 * SH)
        footer = '\u00a9 Gartner Research \u2022 Kill Chain Shift-Left Analysis 2026'

        def pc_color(v):
            if v >= 3.0: return '107c10'
            if v >= 2.5: return '0078d4'
            if v >= 2.0: return 'ca5010'
            return 'a80000'
        def pc_bg(v):
            if v >= 3.0: return 'e6f4e6'
            if v >= 2.5: return 'e0f0ff'
            if v >= 2.0: return 'fff0e0'
            return 'ffe0e0'

        pillar_colors = {'EXM': '107c10', 'AMT': 'a80000', 'ADR': '8764b8', 'PPM': '0078d4', 'SVC': 'ca5010'}

        # Kill chain phase data
        kc_phases = [
            (1, 'Reconnaissance', ['EXM','AMT'], 'PREEMPTIVE', '107c10'),
            (2, 'Weaponization', ['AMT','PPM'], 'PREEMPTIVE', '107c10'),
            (3, 'Delivery', ['EXM','PPM','ADR'], 'PREEMPTIVE', '0078d4'),
            (4, 'Exploitation', ['AMT','PPM','ADR'], 'TRANSITION', 'ca5010'),
            (5, 'Installation', ['ADR'], 'REACTIVE', 'a80000'),
            (6, 'Command & Control', ['AMT','ADR'], 'REACTIVE', 'a80000'),
            (7, 'Actions on Obj.', ['ADR','SVC'], 'REACTIVE', '600000'),
        ]

        pillar_data = [
            ('EXM', 'Exposure Mgmt', 92, 3.33, 3.60, 3.69, 3.07),
            ('PPM', 'Posture & Policy', 86, 2.98, 2.97, 3.21, 2.84),
            ('ADR', 'Detection & Resp', 78, 2.85, 3.38, 2.83, 2.52),
            ('SVC', 'Services & Cap', 57, 2.18, 2.74, 2.32, 1.49),
            ('AMT', 'Adversary Mgmt', 55, 2.36, 2.45, 2.74, 1.87),
        ]
        pa = {p[0]: p[3] for p in pillar_data}

        # Phase scores
        phase_scores = [
            (pa['EXM'] + pa['AMT']) / 2,
            (pa['AMT'] + pa['PPM']) / 2,
            (pa['EXM'] + pa['PPM'] + pa['ADR']) / 3,
            (pa['AMT'] + pa['PPM'] + pa['ADR']) / 3,
            pa['ADR'],
            (pa['AMT'] + pa['ADR']) / 2,
            (pa['ADR'] + pa['SVC']) / 2,
        ]
        preempt_avg = sum(phase_scores[:3]) / 3
        react_avg = sum(phase_scores[4:]) / 3

        # ── SLIDE 1: Kill Chain Phase Mapping ──
        s1 = prs.slides.add_slide(BLANK)
        box(s1, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s1, sx(50), sy(5), sw(1100), sh(28), '1. Cyber Kill Chain \u2192 Preemptive Pillar Mapping', 22, True, 'ffffff', 'center')
        tb(s1, sx(100), sy(32), sw(1000), sh(20), '5 pillars mapped to 7 Lockheed Martin kill chain phases \u2014 shift-left = Phases 1-3', 11, False, 'ccd8e8', 'center')

        # Zone labels
        box(s1, sx(20), sy(65), sw(500), sh(22), '107c10', '107c10')
        tb(s1, sx(25), sy(66), sw(490), sh(20), '\u2190 PREEMPTIVE ZONE (Phases 1-3) \u2192', 10, True, 'ffffff', 'center')
        box(s1, sx(530), sy(65), sw(160), sh(22), 'ca5010', 'ca5010')
        tb(s1, sx(535), sy(66), sw(150), sh(20), 'Transition (4)', 10, True, 'ffffff', 'center')
        box(s1, sx(700), sy(65), sw(480), sh(22), 'a80000', 'a80000')
        tb(s1, sx(705), sy(66), sw(470), sh(20), '\u2190 REACTIVE ZONE (Phases 5-7) \u2192', 10, True, 'ffffff', 'center')

        # Table headers
        headers1 = ['Kill Chain Phase', 'Zone', 'Mapped Pillars', 'Phase Score', 'Defensive Purpose']
        col_w1 = [200, 100, 280, 100, 460]
        cx = 20
        for ci, ht in enumerate(headers1):
            box(s1, sx(cx), sy(95), sw(col_w1[ci]-5), sh(30), '1a3a5c', '1a3a5c')
            tb(s1, sx(cx+5), sy(98), sw(col_w1[ci]-15), sh(26), ht, 10, True, 'ffffff', 'center')
            cx += col_w1[ci]

        descs = [
            'Reduce attack surface, disrupt targeting',
            'Invalidate adversary preparation',
            'Close delivery vectors, deploy deception',
            'Block runtime exploitation, validate controls',
            'Detect persistence via deception & hunting',
            'Disrupt C2, rotate credentials',
            'Managed response, containment, remediation',
        ]
        for ri, (num, name, pillars, zone, color) in enumerate(kc_phases):
            ry = 130 + ri * 55
            zone_bg = 'f0fff4' if zone == 'PREEMPTIVE' else 'fff8f0' if zone == 'TRANSITION' else 'fff5f5'
            # Phase name
            box(s1, sx(20), sy(ry), sw(195), sh(48), zone_bg, color)
            tb(s1, sx(25), sy(ry+10), sw(185), sh(28), f'Phase {num}: {name}', 12, True, color, 'center')
            # Zone
            box(s1, sx(220), sy(ry), sw(95), sh(48), zone_bg, color)
            tb(s1, sx(225), sy(ry+12), sw(85), sh(24), zone, 9, True, color, 'center')
            # Pillars
            box(s1, sx(320), sy(ry), sw(275), sh(48), zone_bg, 'e0ddd5')
            pillar_str = ' + '.join([f'{p} ({[x for x in pillar_data if x[0]==p][0][2]}%)' for p in pillars])
            tb(s1, sx(325), sy(ry+12), sw(265), sh(24), pillar_str, 10, False, '333333', 'center')
            # Score
            score = phase_scores[ri]
            box(s1, sx(600), sy(ry), sw(95), sh(48), pc_bg(score), pc_color(score))
            tb(s1, sx(605), sy(ry+10), sw(85), sh(28), f'{score:.2f}', 16, True, pc_color(score), 'center')
            # Purpose
            box(s1, sx(700), sy(ry), sw(455), sh(48), zone_bg, 'e0ddd5')
            tb(s1, sx(705), sy(ry+12), sw(445), sh(24), descs[ri], 10, False, '555555')

        # Legend
        tb(s1, sx(30), sy(520), sw(200), sh(18), '\u25a0 \u2265 3.0 Strong', 10, True, '107c10')
        tb(s1, sx(230), sy(520), sw(200), sh(18), '\u25a0 2.5-2.99 Competitive', 10, True, '0078d4')
        tb(s1, sx(430), sy(520), sw(200), sh(18), '\u25a0 2.0-2.49 Below Target', 10, True, 'ca5010')
        tb(s1, sx(630), sy(520), sw(200), sh(18), '\u25a0 < 2.0 Critical', 10, True, 'a80000')

        # Key insight
        box(s1, sx(20), sy(550), sw(1160), sh(40), 'fff0f0', 'a80000')
        ml(s1, sx(30), sy(553), sw(1140), sh(34), [
            ('\u26a0 AMT bottleneck:', 11, True, 'a80000'),
            ('At 55% penetration, AMT undermines Phases 1, 2, 4, and 6 \u2014 the shift-left advantage collapses without adversary intelligence.', 10, False, '555555'),
        ])
        tb(s1, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 2: Phase Coverage Scores Bar Chart ──
        s2 = prs.slides.add_slide(BLANK)
        box(s2, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s2, sx(50), sy(5), sw(1100), sh(28), '2. Kill Chain Phase Coverage Scores (Market Average)', 22, True, 'ffffff', 'center')
        tb(s2, sx(100), sy(32), sw(1000), sh(20), f'Preemptive Avg: {preempt_avg:.2f} vs Reactive Avg: {react_avg:.2f} \u2014 Target: 3.0+ for competitive coverage', 11, False, 'ccd8e8', 'center')

        phase_labels = ['Phase 1: Recon', 'Phase 2: Weapon', 'Phase 3: Deliver', 'Phase 4: Exploit', 'Phase 5: Install', 'Phase 6: C2', 'Phase 7: Actions']
        phase_colors = ['107c10', '107c10', '0078d4', 'ca5010', 'a80000', 'a80000', '600000']
        barW = 750
        for pi, (label, score, color) in enumerate(zip(phase_labels, phase_scores, phase_colors)):
            by = 80 + pi * 75
            bg = 'f0fff4' if pi < 3 else 'fff8f0' if pi == 3 else 'fff5f5'
            box(s2, sx(30), sy(by), sw(1140), sh(62), bg, 'e0ddd5')
            tb(s2, sx(40), sy(by+15), sw(220), sh(28), label, 14, True, color)
            box(s2, sx(270), sy(by+10), sw(barW), sh(38), 'e8e8e5', 'cccccc')
            fillW = int(barW * score / 5.0)
            box(s2, sx(270), sy(by+10), sw(fillW), sh(38), color, color)
            tb(s2, sx(270 + fillW + 10), sy(by+15), sw(80), sh(28), f'{score:.2f}', 16, True, color)
            # Target line at 3.0
            targetX = int(barW * 3.0 / 5.0)
            box(s2, sx(270 + targetX), sy(by+8), sw(2), sh(42), '333333', '333333')

        # Averages
        box(s2, sx(30), sy(610), sw(560), sh(45), 'f0fff4', '107c10')
        tb(s2, sx(40), sy(615), sw(540), sh(35), f'\U0001f6e1 Preemptive Zone Average (Phases 1-3): {preempt_avg:.2f}', 14, True, '107c10', 'center')
        box(s2, sx(610), sy(610), sw(560), sh(45), 'fff5f5', 'a80000')
        tb(s2, sx(620), sy(615), sw(540), sh(35), f'\u2694 Reactive Zone Average (Phases 5-7): {react_avg:.2f}', 14, True, 'a80000', 'center')

        box(s2, sx(30), sy(665), sw(1140), sh(35), 'fff8f0', 'ca5010')
        tb(s2, sx(40), sy(668), sw(1120), sh(28), '\u26a0 Phase 7 (Actions on Objectives) is weakest: SVC gap (57% pen.) \u2014 platform-only vendors avg 1.49 on services', 11, True, 'ca5010', 'center')
        tb(s2, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 3: Shift-Left Readiness by Delivery Model ──
        s3 = prs.slides.add_slide(BLANK)
        box(s3, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s3, sx(50), sy(5), sw(1100), sh(28), '3. Shift-Left Readiness by Delivery Model', 22, True, 'ffffff', 'center')
        tb(s3, sx(100), sy(32), sw(1000), sh(20), 'Shift-Left = Phases 1-3 (EXM+AMT+PPM avg)  |  Reactive = Phases 5-7 (ADR+SVC avg)', 11, False, 'ccd8e8', 'center')

        dm_data = [
            ('\U0001f3af Direct Service', '11 vendors (22%)', '107c10', 'f0fff4',
             [('EXM', 3.60), ('AMT', 2.45), ('PPM', 2.97), ('ADR', 3.38), ('SVC', 2.74)]),
            ('\U0001f91d Platform + Partner', '15 vendors (29%)', '0078d4', 'f0faff',
             [('EXM', 3.69), ('AMT', 2.74), ('PPM', 3.21), ('ADR', 2.83), ('SVC', 2.32)]),
            ('\U0001f4bb Platform-Only', '25 vendors (49%)', 'ca5010', 'fff8f0',
             [('EXM', 3.07), ('AMT', 1.87), ('PPM', 2.84), ('ADR', 2.52), ('SVC', 1.49)]),
        ]
        for di, (dname, dsub, dcolor, dbg, dpillars) in enumerate(dm_data):
            dx = 20 + di * 390
            box(s3, sx(dx), sy(65), sw(375), sh(420), dbg, dcolor)
            tb(s3, sx(dx+10), sy(70), sw(355), sh(26), dname, 16, True, dcolor, 'center')
            tb(s3, sx(dx+10), sy(96), sw(355), sh(18), dsub, 11, False, '666666', 'center')
            # Shift-left vs reactive scores
            sl_avg = sum([v for c,v in dpillars[:3]]) / 3
            re_avg = sum([v for c,v in dpillars[3:]]) / 2
            box(s3, sx(dx+20), sy(120), sw(160), sh(60), 'f0fff4', '107c10')
            tb(s3, sx(dx+25), sy(125), sw(150), sh(30), f'{sl_avg:.2f}', 22, True, '107c10', 'center')
            tb(s3, sx(dx+25), sy(155), sw(150), sh(20), 'Shift-Left (1-3)', 10, False, '555555', 'center')
            box(s3, sx(dx+195), sy(120), sw(160), sh(60), 'fff5f5', 'a80000')
            tb(s3, sx(dx+200), sy(125), sw(150), sh(30), f'{re_avg:.2f}', 22, True, 'a80000', 'center')
            tb(s3, sx(dx+200), sy(155), sw(150), sh(20), 'Reactive (5-7)', 10, False, '555555', 'center')
            # Per-pillar bars
            for pi, (pcode, pval) in enumerate(dpillars):
                py = 195 + pi * 52
                is_left = pi < 3
                pct = int((pval / 5.0) * 100)
                barcolor = '107c10' if is_left else 'a80000'
                tb(s3, sx(dx+20), sy(py), sw(50), sh(20), pcode, 11, True, pillar_colors[pcode])
                box(s3, sx(dx+75), sy(py), sw(240), sh(22), 'e8e8e5', 'cccccc')
                box(s3, sx(dx+75), sy(py), sw(int(240*pval/5.0)), sh(22), barcolor, barcolor)
                tb(s3, sx(dx+320), sy(py), sw(40), sh(20), f'{pval:.2f}', 11, True, pc_color(pval))

        box(s3, sx(20), sy(500), sw(1160), sh(40), 'f0faff', '0078d4')
        ml(s3, sx(30), sy(503), sw(1140), sh(34), [
            ('Key Insight:', 11, True, '0078d4'),
            ('Platform-only vendors are structurally confined to Phases 3-5 (AMT: 1.87, SVC: 1.49) \u2014 cannot defend the earliest or latest kill chain phases.', 10, False, '555555'),
        ])
        tb(s3, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 4: Shift-Left Imperative Executive Summary ──
        s4 = prs.slides.add_slide(BLANK)
        box(s4, sx(0), sy(0), sw(1200), sh(750), '1a3a5c', '1a3a5c')
        tb(s4, sx(50), sy(10), sw(1100), sh(32), 'Reimagining Threat Defense Through Preemptive Cybersecurity', 24, True, 'ffffff', 'center')
        tb(s4, sx(100), sy(44), sw(1000), sh(20), 'Kill Chain Phase Coverage Analysis \u2014 51 Vendors, 5 Pillars, 7 Phases', 12, False, '8cb8e0', 'center')

        # Stat cards
        s4_stats = [
            ('51', 'Vendors Assessed', '0078d4'),
            ('78%', 'Preemptive Zone Pen.', '107c10'),
            ('68%', 'Reactive Zone Pen.', 'a80000'),
            ('55%', 'AMT (Bottleneck)', 'ff8c00'),
            ('27%', 'Full Kill Chain', '8764b8'),
        ]
        for si, (sv, sl, sc) in enumerate(s4_stats):
            bx = 40 + si * 225
            box(s4, sx(bx), sy(75), sw(210), sh(55), '2a4a6c', '3a5a7c')
            tb(s4, sx(bx+5), sy(78), sw(200), sh(28), sv, 22, True, sc, 'center')
            tb(s4, sx(bx+5), sy(108), sw(200), sh(18), sl, 9, False, 'aaaaaa', 'center')

        # Findings
        tb(s4, sx(40), sy(145), sw(200), sh(18), 'KEY FINDINGS', 12, True, 'ffcc00')
        s4_findings = [
            ('\U0001f4ca AMT Undermines Shift-Left', 'EXM 92% + PPM 86%, but AMT 55% = single point of failure at Phases 1-2.', 'ff8c00'),
            ('\U0001f4c9 Platform-Only: Phases 3-5 Only', '49% of market (25 vendors): AMT 1.87, SVC 1.49 \u2014 structurally unable to defend full chain.', 'ff8888'),
            ('\u2705 27% Full Kill Chain', '14 full-spectrum vendors maintain \u2265 2.0 across all pillars \u2014 genuine defense-in-depth.', '8eff8e'),
        ]
        for fi, (ft, fb, fc) in enumerate(s4_findings):
            fy = 170 + fi * 65
            box(s4, sx(40), sy(fy), sw(1120), sh(55), '2a4a6c', '3a5a7c')
            tb(s4, sx(55), sy(fy+5), sw(1090), sh(22), ft, 13, True, fc)
            tb(s4, sx(55), sy(fy+28), sw(1090), sh(22), fb, 10, False, 'cccccc')

        # Recommendations
        tb(s4, sx(40), sy(375), sw(300), sh(18), 'RECOMMENDATIONS FOR CPOs', 12, True, 'ffcc00')
        s4_recs = [
            ('01', 'Map Kill Chain', 'Score each phase, find < 2.0 gaps', '0078d4'),
            ('02', 'Invest in AMT', 'Polymorphic defense, MTD, rotation', '107c10'),
            ('03', 'Reframe GTM', 'Kill chain phases, not features', '8764b8'),
            ('04', 'Close SVC Gap', 'Build / acquire / partner', 'ca5010'),
        ]
        for ri, (rn, rt, rd, rc) in enumerate(s4_recs):
            rx = 40 + ri * 285
            box(s4, sx(rx), sy(400), sw(265), sh(80), '2a4a6c', rc)
            tb(s4, sx(rx+10), sy(403), sw(245), sh(14), rn, 9, True, rc)
            tb(s4, sx(rx+10), sy(418), sw(245), sh(22), rt, 14, True, 'ffffff', 'center')
            tb(s4, sx(rx+10), sy(445), sw(245), sh(28), rd, 10, False, 'aaaaaa', 'center')

        # SPA
        box(s4, sx(40), sy(500), sw(1120), sh(40), '4a2020', 'a80000')
        tb(s4, sx(50), sy(503), sw(1100), sh(34), '\u26a0 SPA: By 2028, 35% of evaluations will use kill chain phase coverage as primary vendor criterion \u2014 up from <10% today.', 11, True, 'ff8888', 'center')
        tb(s4, sx(800), sy(720), sw(380), sh(20), footer, 9, False, '888888', 'right')

        # ── SLIDE 5: Preemptive vs Reactive VS Poster ──
        s5 = prs.slides.add_slide(BLANK)
        box(s5, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s5, sx(50), sy(5), sw(1100), sh(28), '5. The Great Divide: Preemptive vs Reactive Coverage', 22, True, 'ffffff', 'center')
        tb(s5, sx(100), sy(32), sw(1000), sh(20), 'Two fundamentally different defensive postures across the kill chain', 11, False, 'ccd8e8', 'center')

        # Preemptive side
        box(s5, sx(30), sy(65), sw(530), sh(550), 'f0fff4', '107c10')
        tb(s5, sx(40), sy(70), sw(510), sh(28), '\U0001f6e1 PREEMPTIVE ZONE', 20, True, '107c10', 'center')
        tb(s5, sx(40), sy(98), sw(510), sh(18), 'Kill Chain Phases 1-3: Recon \u2192 Weaponize \u2192 Deliver', 11, False, '555555', 'center')
        tb(s5, sx(40), sy(130), sw(510), sh(40), '78%', 36, True, '107c10', 'center')
        tb(s5, sx(40), sy(170), sw(510), sh(18), 'Average vendor penetration', 10, False, '666666', 'center')

        pre_pillars = [
            ('EXM', 'Exposure Management', 92, 'Reduces what adversaries discover'),
            ('PPM', 'Posture & Policy', 86, 'Validates controls before attacks'),
            ('AMT', 'Adversary Management', 55, 'Disrupts adversary preparation'),
        ]
        for pi, (pc, pn, pv, pd) in enumerate(pre_pillars):
            py = 200 + pi * 65
            pbg = 'e6f4e6' if pv >= 80 else 'fff0e0'
            box(s5, sx(50), sy(py), sw(490), sh(55), pbg, pillar_colors[pc])
            tb(s5, sx(60), sy(py+5), sw(50), sh(20), pc, 12, True, pillar_colors[pc])
            tb(s5, sx(120), sy(py+5), sw(250), sh(18), pn, 11, True, '333333')
            tb(s5, sx(120), sy(py+28), sw(250), sh(18), pd, 9, False, '777777')
            pvc = '107c10' if pv >= 80 else '0078d4' if pv >= 60 else 'a80000'
            tb(s5, sx(430), sy(py+10), sw(100), sh(30), f'{pv}%', 18, True, pvc, 'center')

        box(s5, sx(50), sy(400), sw(490), sh(30), '107c10', '107c10')
        tb(s5, sx(55), sy(403), sw(480), sh(24), '\u2705 Prevents incidents before exploitation occurs', 11, True, 'ffffff', 'center')

        # VS divider
        tb(s5, sx(565), sy(300), sw(70), sh(40), 'VS', 32, True, '888888', 'center')

        # Reactive side
        box(s5, sx(640), sy(65), sw(530), sh(550), 'fff5f5', 'a80000')
        tb(s5, sx(650), sy(70), sw(510), sh(28), '\u2694 REACTIVE ZONE', 20, True, 'a80000', 'center')
        tb(s5, sx(650), sy(98), sw(510), sh(18), 'Kill Chain Phases 5-7: Install \u2192 C2 \u2192 Actions', 11, False, '555555', 'center')
        tb(s5, sx(650), sy(130), sw(510), sh(40), '68%', 36, True, 'a80000', 'center')
        tb(s5, sx(650), sy(170), sw(510), sh(18), 'Average vendor penetration', 10, False, '666666', 'center')

        react_pillars = [
            ('ADR', 'Detection & Response', 78, 'Detects after exploitation occurs'),
            ('SVC', 'Services & Capability', 57, 'Managed response after breach'),
        ]
        for pi, (pc, pn, pv, pd) in enumerate(react_pillars):
            py = 200 + pi * 65
            pbg = 'e0f0ff' if pv >= 70 else 'ffe0e0'
            box(s5, sx(660), sy(py), sw(490), sh(55), pbg, pillar_colors[pc])
            tb(s5, sx(670), sy(py+5), sw(50), sh(20), pc, 12, True, pillar_colors[pc])
            tb(s5, sx(730), sy(py+5), sw(250), sh(18), pn, 11, True, '333333')
            tb(s5, sx(730), sy(py+28), sw(250), sh(18), pd, 9, False, '777777')
            pvc = '107c10' if pv >= 80 else '0078d4' if pv >= 60 else 'a80000'
            tb(s5, sx(1040), sy(py+10), sw(100), sh(30), f'{pv}%', 18, True, pvc, 'center')

        box(s5, sx(660), sy(400), sw(490), sh(30), 'a80000', 'a80000')
        tb(s5, sx(665), sy(403), sw(480), sh(24), '\u26a0 Responds after damage has begun', 11, True, 'ffffff', 'center')

        # Paradox callout
        box(s5, sx(30), sy(640), sw(1140), sh(40), 'fff8f0', 'ca5010')
        tb(s5, sx(40), sy(643), sw(1120), sh(34), 'The Paradox: A market labeled "preemptive" that is structurally more reactive \u2014 because the AMT gap (55%) undermines the left side of the kill chain.', 11, True, 'ca5010', 'center')
        tb(s5, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 6: CPO Investment Roadmap ──
        s6 = prs.slides.add_slide(BLANK)
        box(s6, sx(0), sy(0), sw(1200), sh(55), '1a3a5c', '1a3a5c')
        tb(s6, sx(50), sy(5), sw(1100), sh(28), '6. CPO Roadmap: Full Kill Chain Coverage by 2028', 22, True, 'ffffff', 'center')
        tb(s6, sx(100), sy(32), sw(1000), sh(20), 'Three-phase investment strategy for Chief Product Officers', 11, False, 'ccd8e8', 'center')

        roadmap = [
            ('\U0001f5fa Phase 1: 2025-26', 'Assess & Map', 'f0faff', '0078d4',
             ['Conduct kill chain mapping', 'Score each phase 1-5', 'Identify gaps below 2.0', 'Develop Shift-Left Index', 'Benchmark vs market avg']),
            ('\U0001f3d7 Phase 2: 2026-27', 'Build Preemptive Core', 'f0fff4', '107c10',
             ['Invest in AMT (target > 3.0)', 'Polymorphic & MTD', 'Runtime app protection', 'Automated credential rotation', 'Embed BAS + validation']),
            ('\U0001f680 Phase 3: 2027-28', 'Complete Kill Chain', 'f5f0ff', '8764b8',
             ['Close SVC gap', 'Outcome-based delivery', '"Full Kill Chain" positioning', 'Kill chain GTM messaging', 'Target 2x win rate']),
        ]
        for ri, (rtitle, rsub, rbg, rcolor, ritems) in enumerate(roadmap):
            rx = 20 + ri * 390
            box(s6, sx(rx), sy(65), sw(375), sh(320), rbg, rcolor)
            tb(s6, sx(rx+10), sy(70), sw(355), sh(24), rtitle, 14, True, rcolor, 'center')
            tb(s6, sx(rx+10), sy(95), sw(355), sh(26), rsub, 18, True, '333333', 'center')
            for ii, it in enumerate(ritems):
                tb(s6, sx(rx+20), sy(130 + ii * 30), sw(335), sh(24), f'\u2022 {it}', 12, False, '333333')
            if ri < 2:
                tb(s6, sx(rx + 380), sy(200), sw(30), sh(28), '\u2192', 22, True, '888888', 'center')

        # Metrics
        tb(s6, sx(30), sy(400), sw(400), sh(22), '\U0001f3af SUCCESS METRICS BY 2028', 14, True, '1a3a5c')
        metrics = [
            ('\u2705', 'All 7 Phases \u2265 2.5', '107c10'),
            ('\U0001f6e1', 'AMT Score \u2265 3.0', '0078d4'),
            ('\U0001f4ca', 'Shift-Left Index', '8764b8'),
            ('\U0001f4e3', 'Kill Chain GTM', 'ca5010'),
            ('\U0001f3c6', '2x Win Rate', '107c10'),
        ]
        for mi, (micon, mlabel, mcolor) in enumerate(metrics):
            mx = 30 + mi * 232
            box(s6, sx(mx), sy(430), sw(220), sh(55), 'f8f8f5', mcolor)
            tb(s6, sx(mx+5), sy(433), sw(210), sh(22), micon, 18, False, mcolor, 'center')
            tb(s6, sx(mx+5), sy(460), sw(210), sh(20), mlabel, 11, True, mcolor, 'center')

        box(s6, sx(30), sy(500), sw(1140), sh(35), 'f0fff0', '107c10')
        tb(s6, sx(40), sy(503), sw(1120), sh(28), 'The window for differentiation is 2-3 years. By 2028, kill chain coverage transitions from advantage to minimum expectation.', 11, True, '107c10', 'center')
        tb(s6, sx(800), sy(720), sw(380), sh(20), footer, 9, False, 'aaaaaa', 'right')

        # ── SLIDE 7: Kill Chain Defense Architecture Infographic ──
        s7 = prs.slides.add_slide(BLANK)
        def sx7(v): return int(v / 1200 * SW)
        def sy7(v): return int(v / 920 * SH)
        def sw7(v): return int(v / 1200 * SW)
        def sh7(v): return int(v / 920 * SH)

        box(s7, sx7(0), sy7(0), sw7(1200), sh7(70), '1a3a5c', '1a3a5c')
        tb(s7, sx7(50), sy7(5), sw7(1100), sh7(32), 'Reimagining Threat Defense: The Shift-Left Kill Chain', 22, True, 'ffffff', 'center')
        tb(s7, sx7(100), sy7(38), sw7(1000), sh7(28), 'How preemptive cybersecurity pillars map to the Lockheed Martin Cyber Kill Chain', 11, False, 'ccd8e8', 'center')

        # Kill chain arrow phases
        kc_arrow = [
            ('1. Recon', '107c10', 'EXM + AMT', 'PREEMPTIVE'),
            ('2. Weapon', '107c10', 'AMT + PPM', 'PREEMPTIVE'),
            ('3. Deliver', '0078d4', 'EXM+PPM+ADR', 'PREEMPTIVE'),
            ('4. Exploit', 'ca5010', 'AMT+PPM+ADR', 'TRANSITION'),
            ('5. Install', 'a80000', 'ADR', 'REACTIVE'),
            ('6. C2', 'a80000', 'AMT + ADR', 'REACTIVE'),
            ('7. Actions', '600000', 'ADR + SVC', 'REACTIVE'),
        ]
        for ai, (aname, acolor, apillars, azone) in enumerate(kc_arrow):
            ax = 20 + ai * 168
            box(s7, sx7(ax), sy7(85), sw7(155), sh7(65), acolor, acolor)
            tb(s7, sx7(ax+5), sy7(88), sw7(145), sh7(20), aname, 12, True, 'ffffff', 'center')
            tb(s7, sx7(ax+5), sy7(112), sw7(145), sh7(14), apillars, 9, False, 'ffffffcc', 'center')
            tb(s7, sx7(ax+5), sy7(130), sw7(145), sh7(14), azone, 8, True, 'ffffffaa', 'center')

        # Zone labels
        box(s7, sx7(20), sy7(168), sw7(498), sh7(24), '107c10', '107c10')
        tb(s7, sx7(25), sy7(170), sw7(488), sh7(20), '\u2190 PREEMPTIVE ZONE (Phases 1-3) \u2192', 10, True, 'ffffff', 'center')
        box(s7, sx7(530), sy7(168), sw7(155), sh7(24), 'ca5010', 'ca5010')
        tb(s7, sx7(535), sy7(170), sw7(145), sh7(20), 'Transition', 10, True, 'ffffff', 'center')
        box(s7, sx7(698), sy7(168), sw7(498), sh7(24), 'a80000', 'a80000')
        tb(s7, sx7(703), sy7(170), sw7(488), sh7(20), '\u2190 REACTIVE ZONE (Phases 5-7) \u2192', 10, True, 'ffffff', 'center')

        # Pillar cards
        p7_data = [
            ('EXM', 'Exposure Mgmt', '92%', '3.33', '107c10', 'f0fff4', 'Phases: 1, 3'),
            ('AMT', 'Adversary Mgmt', '55%', '2.36', 'a80000', 'fff5f5', 'Phases: 1, 2, 4, 6'),
            ('PPM', 'Posture & Policy', '86%', '2.98', '0078d4', 'f0faff', 'Phases: 2, 3, 4'),
            ('ADR', 'Detection & Resp', '78%', '2.85', '8764b8', 'f5f0ff', 'Phases: 3-7'),
            ('SVC', 'Services & Cap', '57%', '2.18', 'ca5010', 'fff8f0', 'Phase: 7'),
        ]
        for pi, (pcode, pname, ppen, pavg, pcolor, pbg, pphases) in enumerate(p7_data):
            px = 20 + pi * 234
            box(s7, sx7(px), sy7(210), sw7(220), sh7(100), pbg, pcolor)
            tb(s7, sx7(px+5), sy7(215), sw7(210), sh7(18), f'{pcode}: {pname}', 12, True, pcolor, 'center')
            tb(s7, sx7(px+5), sy7(238), sw7(210), sh7(16), pphases, 9, False, '555555', 'center')
            tb(s7, sx7(px+5), sy7(260), sw7(210), sh7(22), f'Pen: {ppen}', 14, True, pcolor, 'center')
            tb(s7, sx7(px+5), sy7(283), sw7(210), sh7(16), f'Market Avg: {pavg}', 10, False, '333333', 'center')

        # Delivery model section
        box(s7, sx7(20), sy7(330), sw7(1160), sh7(24), '00b4d8', '0095b3')
        tb(s7, sx7(25), sy7(332), sw7(1150), sh7(20), 'KILL CHAIN COVERAGE BY DELIVERY MODEL', 11, True, 'ffffff', 'center')

        dm7 = [
            ('Direct Service', '11 (22%)', '3.01', '3.06', '107c10', 'f0fff4',
             ['EXM:3.60 AMT:2.45 PPM:2.97', 'ADR:3.38 SVC:2.74', '\u2713 Full kill chain coverage']),
            ('Platform + Partner', '15 (29%)', '3.21', '2.58', '0078d4', 'f0faff',
             ['EXM:3.69 AMT:2.74 PPM:3.21', 'ADR:2.83 SVC:2.32', '\u2713 Best shift-left scores']),
            ('Platform-Only', '25 (49%)', '2.59', '2.01', 'ca5010', 'fff8f0',
             ['EXM:3.07 AMT:1.87 PPM:2.84', 'ADR:2.52 SVC:1.49', '\u26a0 Phases 1-2 & 5-7 gaps']),
        ]
        for di, (dname, dcount, dsl, dre, dcolor, dbg, dlines) in enumerate(dm7):
            dx = 20 + di * 390
            box(s7, sx7(dx), sy7(365), sw7(375), sh7(160), dbg, dcolor)
            tb(s7, sx7(dx+10), sy7(370), sw7(355), sh7(20), dname, 13, True, dcolor, 'center')
            tb(s7, sx7(dx+10), sy7(392), sw7(355), sh7(16), dcount, 10, False, '666666', 'center')
            tb(s7, sx7(dx+20), sy7(415), sw7(160), sh7(16), f'Shift-Left: {dsl}', 10, True, '107c10')
            tb(s7, sx7(dx+195), sy7(415), sw7(160), sh7(16), f'Reactive: {dre}', 10, True, 'a80000')
            for li, lt in enumerate(dlines):
                tb(s7, sx7(dx+10), sy7(440 + li * 22), sw7(355), sh7(18), lt, 9, False, '555555', 'center')

        # Bottleneck callout
        box(s7, sx7(20), sy7(545), sw7(1160), sh7(50), 'fff0f0', 'a80000')
        tb(s7, sx7(30), sy7(550), sw7(1140), sh7(18), '\u26a0 THE SHIFT-LEFT BOTTLENECK: AMT at 55% undermines the entire preemptive value proposition', 12, True, 'a80000', 'center')
        tb(s7, sx7(30), sy7(572), sw7(1140), sh7(18), '49% of vendors (platform-only) cannot defend Phases 1-2 or Phases 5-7 \u2014 confined to the middle of the kill chain', 10, False, '666666', 'center')

        # Stat boxes
        for si, (sv, sl, sc, sbg) in enumerate([
            ('92%', 'EXM Pen.', '107c10', 'f0fff4'), ('55%', 'AMT Pen.', 'a80000', 'fff5f5'),
            ('86%', 'PPM Pen.', '0078d4', 'f0faff'), ('78%', 'ADR Pen.', '8764b8', 'f5f0ff'),
            ('57%', 'SVC Pen.', 'ca5010', 'fff8f0'),
        ]):
            bx = 30 + si * 230
            box(s7, sx7(bx), sy7(615), sw7(215), sh7(55), sbg, sc)
            tb(s7, sx7(bx+5), sy7(618), sw7(205), sh7(28), sv, 20, True, sc, 'center')
            tb(s7, sx7(bx+5), sy7(648), sw7(205), sh7(16), sl, 9, False, '555555', 'center')

        # Action plan strip
        box(s7, sx7(20), sy7(690), sw7(1160), sh7(120), 'f8f8f5', 'e0ddd5')
        tb(s7, sx7(30), sy7(695), sw7(1140), sh7(18), 'CPO ACTION PLAN: FROM FEATURE VENDOR TO KILL CHAIN DEFENDER', 11, True, '1a3a5c', 'center')
        for ai, (atitle, adesc, acolor) in enumerate([
            ('01 \u2013 Map', 'Kill chain scoring', '0078d4'),
            ('02 \u2013 Invest', 'AMT to 3.0+', '107c10'),
            ('03 \u2013 Reframe', 'Kill chain GTM', '8764b8'),
            ('04 \u2013 Close SVC', 'Build/acquire/partner', 'ca5010'),
            ('05 \u2013 Index', 'Shift-Left metric', 'a80000'),
        ]):
            ax = 40 + ai * 226
            box(s7, sx7(ax), sy7(720), sw7(210), sh7(50), acolor, acolor)
            tb(s7, sx7(ax+5), sy7(723), sw7(200), sh7(18), atitle, 11, True, 'ffffff', 'center')
            tb(s7, sx7(ax+5), sy7(747), sw7(200), sh7(16), adesc, 9, False, 'ffffffcc', 'center')

        tb(s7, sx7(30), sy7(785), sw7(1140), sh7(18), 'By 2028, kill chain phase coverage transitions from competitive advantage to minimum buyer expectation.', 11, True, '1a3a5c', 'center')
        tb(s7, sx7(30), sy7(808), sw7(1140), sh7(16), 'CPOs who act now define the market standard. Those who wait will be measured against it.', 10, False, '666666', 'center')
        tb(s7, sx7(750), sy7(895), sw7(430), sh7(16), '\u00a9 Gartner Research \u2022 PreCyber Kill Chain 2026', 9, False, 'aaaaaa', 'right')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='PreCyber_KillChain_All_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/precyber-kcmitre-pptx', methods=['GET'])
def precyber_kcmitre_pptx():
    """Generate a 7-slide PPTX deck for the Kill Chain + MITRE ATT&CK v2 report."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from flask import send_file
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
        def box(slide, l, t, w, h, fill=None, border=None):
            from pptx.util import Emu as E
            s = slide.shapes.add_shape(1, E(l), E(t), E(w), E(h))
            s.line.fill.background()
            if fill:
                s.fill.solid()
                s.fill.fore_color.rgb = rgb(fill)
            if border:
                s.line.fill.solid()
                s.line.fill.fore_color.rgb = rgb(border)
                s.line.width = Pt(1)
            return s
        def tb(slide, l, t, w, h, text, size=12, bold=False, color='333333', align='left'):
            from pptx.util import Emu as E
            txBox = slide.shapes.add_textbox(E(l), E(t), E(w), E(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = rgb(color)
            if align == 'center': p.alignment = PP_ALIGN.CENTER
            elif align == 'right': p.alignment = PP_ALIGN.RIGHT
            return txBox
        W = int(13.333 * 914400)
        H = int(7.5 * 914400)
        def sx(pct): return int(W * pct / 100)
        def sy(pct): return int(H * pct / 100)
        def sw(pct): return int(W * pct / 100)
        def sh(pct): return int(H * pct / 100)
        def ml(lines, size=11, bold=False, color='333333', align='left', spacing=1.2):
            """Create multi-line text helper returning list of (text, size, bold, color) tuples."""
            return lines

        pillar_key_text = 'Key: EXM = Exposure Mgmt | AMT = Adversary Mgmt | PPM = Posture & Policy | ADR = Detection & Response | SVC = Cross-Cutting Enabler'

        # ── Slide 1: Dual Framework Alignment Table ──
        s1 = prs.slides.add_slide(BLANK)
        box(s1, 0, 0, W, sy(12), '#1a3a5c')
        tb(s1, sx(3), sy(2), sw(94), sh(5), 'Dual Framework Alignment: Kill Chain + MITRE ATT&CK', 24, True, 'ffffff', 'center')
        tb(s1, sx(3), sy(7), sw(94), sh(4), 'How kill chain phases map to ATT&CK tactics and preemptive pillars', 13, False, 'aac8e8', 'center')

        mapping = [
            ('1. Reconnaissance', 'Reconnaissance (TA0043)', 'EXM, AMT', 'Preemptive', '107c10'),
            ('2. Weaponization', 'Resource Development (TA0042)', 'AMT, PPM', 'Preemptive', '107c10'),
            ('3. Delivery', 'Initial Access (TA0001)', 'EXM, PPM, ADR', 'Preemptive', '0078d4'),
            ('4. Exploitation', 'Execution / Priv. Esc. (TA0002/4)', 'AMT, PPM, ADR', 'Transition', 'ca5010'),
            ('5. Installation', 'Persistence / Def. Evasion (TA0003/5)', 'ADR', 'Reactive', 'a80000'),
            ('6. Command & Control', 'Command and Control (TA0011)', 'AMT, ADR', 'Reactive', 'a80000'),
            ('7. Actions on Obj.', 'Exfiltration / Impact (TA0010/40)', 'ADR', 'Reactive', '600000'),
        ]
        headers = ['Kill Chain Phase', 'MITRE ATT&CK Tactic(s)', 'Pillars', 'Zone']
        col_x = [sx(3), sx(25), sx(58), sx(82)]
        col_w = [sw(21), sw(32), sw(23), sw(15)]
        row_y = sy(16)
        row_h = sh(5)
        for i, hdr in enumerate(headers):
            box(s1, col_x[i], row_y, col_w[i], row_h, '1a3a5c')
            tb(s1, col_x[i]+sw(0.5), row_y+sh(0.5), col_w[i]-sw(1), row_h, hdr, 12, True, 'ffffff')
        for ri, (kc, att, pil, zone, col) in enumerate(mapping):
            ry = row_y + row_h + ri * row_h
            bg = 'f8f8f5' if ri % 2 == 0 else 'ffffff'
            box(s1, col_x[0], ry, col_w[0], row_h, bg)
            tb(s1, col_x[0]+sw(0.5), ry+sh(0.5), col_w[0]-sw(1), row_h, kc, 11, True, col)
            box(s1, col_x[1], ry, col_w[1], row_h, bg)
            tb(s1, col_x[1]+sw(0.5), ry+sh(0.5), col_w[1]-sw(1), row_h, att, 10, False, '333333')
            box(s1, col_x[2], ry, col_w[2], row_h, bg)
            tb(s1, col_x[2]+sw(0.5), ry+sh(0.5), col_w[2]-sw(1), row_h, pil, 10, True, col)
            box(s1, col_x[3], ry, col_w[3], row_h, bg)
            tb(s1, col_x[3]+sw(0.5), ry+sh(0.5), col_w[3]-sw(1), row_h, zone, 11, True, col, 'center')

        # Zone legend
        for i, (lbl, col) in enumerate([('Preemptive (Phases 1-3)','107c10'), ('Transition (Phase 4)','ca5010'), ('Reactive (Phases 5-7)','a80000')]):
            bx = sx(3) + i * sw(25)
            box(s1, bx, sy(58), sw(2), sh(2), col)
            tb(s1, bx + sw(3), sy(58), sw(20), sh(3), lbl, 10, True, col)
        tb(s1, sx(3), sy(64), sw(94), sh(3), pillar_key_text, 9, False, '888888')

        # ── Slide 2: ATT&CK Tactic Coverage ──
        s2 = prs.slides.add_slide(BLANK)
        box(s2, 0, 0, W, sy(12), '#1a3a5c')
        tb(s2, sx(3), sy(2), sw(94), sh(5), 'MITRE ATT&CK Tactic Coverage Assessment', 24, True, 'ffffff', 'center')
        tb(s2, sx(3), sy(7), sw(94), sh(4), 'Market-wide coverage strength for 14 enterprise tactics', 13, False, 'aac8e8', 'center')

        tactics = [
            ('Reconnaissance (TA0043)', 'EXM', 'Strong', '107c10'),
            ('Resource Development (TA0042)', 'AMT', 'Weak', 'a80000'),
            ('Initial Access (TA0001)', 'EXM+PPM+ADR', 'Strong', '107c10'),
            ('Execution (TA0002)', 'ADR', 'Moderate', '0078d4'),
            ('Persistence (TA0003)', 'ADR', 'Moderate', '0078d4'),
            ('Privilege Escalation (TA0004)', 'PPM', 'Moderate', '0078d4'),
            ('Defense Evasion (TA0005)', 'AMT', 'Weak', 'a80000'),
            ('Credential Access (TA0006)', 'AMT', 'Weak', 'a80000'),
            ('Discovery (TA0007)', 'EXM', 'Strong', '107c10'),
            ('Lateral Movement (TA0008)', 'ADR', 'Weak', 'a80000'),
            ('Collection (TA0009)', 'ADR', 'Moderate', '0078d4'),
            ('Command & Control (TA0011)', 'AMT+ADR', 'Moderate', '0078d4'),
            ('Exfiltration (TA0010)', 'ADR', 'Moderate', '0078d4'),
            ('Impact (TA0040)', 'ADR', 'Moderate', '0078d4'),
        ]
        per_col = 7
        for ti, (tac, pil, strength, col) in enumerate(tactics):
            c = ti // per_col
            r = ti % per_col
            bx = sx(3) + c * sw(48)
            by = sy(15) + r * sh(8)
            bg = 'ffe0e0' if strength == 'Weak' else ('e0f0ff' if strength == 'Moderate' else 'e6f4e6')
            box(s2, bx, by, sw(45), sh(7), bg, col)
            tb(s2, bx + sw(1), by + sh(0.5), sw(30), sh(3), tac, 10, True, col)
            tb(s2, bx + sw(1), by + sh(3.5), sw(15), sh(3), pil, 9, False, '888888')
            tb(s2, bx + sw(32), by + sh(1.5), sw(12), sh(4), strength, 11, True, col, 'right')

        tb(s2, sx(3), sy(73), sw(94), sh(3), pillar_key_text, 9, False, '888888')
        for i, (lbl, col) in enumerate([('Strong: Most vendors provide coverage','107c10'), ('Moderate: Partial coverage','0078d4'), ('Weak: Below threshold','a80000')]):
            bx = sx(3) + i * sw(30)
            box(s2, bx, sy(77), sw(1.5), sh(2), col)
            tb(s2, bx + sw(2), sy(77), sw(26), sh(3), lbl, 9, True, col)

        # ── Slide 3: Delivery Model Profiles ──
        s3 = prs.slides.add_slide(BLANK)
        box(s3, 0, 0, W, sy(12), '#1a3a5c')
        tb(s3, sx(3), sy(2), sw(94), sh(5), 'Delivery Model Adversary Lifecycle Profiles', 24, True, 'ffffff', 'center')
        tb(s3, sx(3), sy(7), sw(94), sh(4), 'Each model shows distinct coverage across adversary lifecycle phases', 13, False, 'aac8e8', 'center')

        model_cards = [
            ('Direct Service', '107c10', '11 vendors', 'Strong', 'Competitive', 'Broadest lifecycle coverage'),
            ('Platform + Partner', '0078d4', '15 vendors', 'Strong', 'Below Target', 'Strong early, partner-dependent late'),
            ('Platform-Only', 'ca5010', '25 vendors', 'Below Target', 'Critical Gap', 'Middle lifecycle only'),
        ]
        for mi, (name, col, count, early, late, summary) in enumerate(model_cards):
            cx = sx(3) + mi * sw(32)
            box(s3, cx, sy(15), sw(30), sh(65), 'f8f8f5', col)
            box(s3, cx, sy(15), sw(30), sh(4), col)
            tb(s3, cx + sw(1), sy(16), sw(28), sh(3), name, 16, True, 'ffffff', 'center')
            tb(s3, cx + sw(1), sy(21), sw(28), sh(3), count, 11, False, '888888', 'center')

            box(s3, cx + sw(2), sy(26), sw(12), sh(20), 'f0fff0', '107c10')
            tb(s3, cx + sw(2.5), sy(27), sw(11), sh(3), 'Early Phases', 10, True, '107c10', 'center')
            tb(s3, cx + sw(2.5), sy(30), sw(11), sh(3), 'KC 1-3', 9, False, '888888', 'center')
            early_col = '107c10' if early in ('Strong','Competitive') else 'a80000'
            tb(s3, cx + sw(2.5), sy(35), sw(11), sh(5), early, 14, True, early_col, 'center')

            box(s3, cx + sw(16), sy(26), sw(12), sh(20), 'fff0f0', 'a80000')
            tb(s3, cx + sw(16.5), sy(27), sw(11), sh(3), 'Late Phases', 10, True, 'a80000', 'center')
            tb(s3, cx + sw(16.5), sy(30), sw(11), sh(3), 'KC 5-7', 9, False, '888888', 'center')
            late_col = '107c10' if late in ('Strong','Competitive') else ('ca5010' if late == 'Below Target' else 'a80000')
            tb(s3, cx + sw(16.5), sy(35), sw(11), sh(5), late, 14, True, late_col, 'center')

            tb(s3, cx + sw(1), sy(50), sw(28), sh(5), summary, 11, False, '555555', 'center')

        tb(s3, sx(3), sy(83), sw(94), sh(3), pillar_key_text, 9, False, '888888')

        # ── Slide 4: Executive Summary ──
        s4 = prs.slides.add_slide(BLANK)
        box(s4, 0, 0, W, H, '1a3a5c')
        tb(s4, sx(3), sy(4), sw(94), sh(6), 'Reimagining Threat Defense Through Preemptive Cybersecurity', 26, True, 'ffffff', 'center')
        tb(s4, sx(3), sy(10), sw(94), sh(4), '51 vendors assessed across Kill Chain and MITRE ATT&CK frameworks', 14, False, 'aac8e8', 'center')

        stats_cards = [
            ('Preemptive Zone\nPenetration', '~78%', 'KC 1-3 / ATT&CK Recon-IA', '7eff7e'),
            ('Reactive Zone\nPenetration', '~68%', 'KC 5-7 / ATT&CK LM-Impact', 'ff8888'),
            ('Full Lifecycle\nVendors', '~1 in 4', 'All 7 KC phases covered', 'ffcc00'),
        ]
        for si, (lbl, val, sub, col) in enumerate(stats_cards):
            bx = sx(8) + si * sw(30)
            box(s4, bx, sy(18), sw(26), sh(22), None, col)
            from pptx.util import Emu as _E
            s = s4.shapes.add_shape(1, _E(bx), _E(sy(18)), _E(sw(26)), _E(sh(22)))
            s.fill.solid()
            s.fill.fore_color.rgb = rgb('0a2a4c')
            s.line.fill.solid()
            s.line.fill.fore_color.rgb = rgb(col)
            s.line.width = Pt(2)
            tb(s4, bx + sw(1), sy(20), sw(24), sh(6), val, 28, True, col, 'center')
            tb(s4, bx + sw(1), sy(28), sw(24), sh(4), lbl, 11, False, 'ffffff', 'center')
            tb(s4, bx + sw(1), sy(34), sw(24), sh(3), sub, 9, False, 'aac8e8', 'center')

        findings_box = [
            ('Strongest ATT&CK Tactic', 'Reconnaissance (TA0043)', '7eff7e'),
            ('Weakest ATT&CK Tactic', 'Resource Dev. (TA0042)', 'ff8888'),
            ('Critical Bottleneck', 'Adversary Management', 'ffcc00'),
            ('Largest Segment', 'Platform-Only (~49%)', 'ff8888'),
        ]
        for fi, (lbl, val, col) in enumerate(findings_box):
            bx = sx(5) + fi * sw(23)
            box(s4, bx, sy(46), sw(21), sh(12), '0a2a4c')
            tb(s4, bx + sw(1), sy(47), sw(19), sh(3), lbl, 9, False, 'aac8e8', 'center')
            tb(s4, bx + sw(1), sy(51), sw(19), sh(5), val, 12, True, col, 'center')

        box(s4, sx(5), sy(62), sw(90), sh(10), '0a2a4c')
        tb(s4, sx(6), sy(63), sw(88), sh(8), 'The market is transitioning from feature-list competition to adversary lifecycle phase coverage competition, measurable through both Kill Chain and ATT&CK frameworks.', 12, True, 'ffffff', 'center')

        tb(s4, sx(3), sy(93), sw(94), sh(3), pillar_key_text, 9, False, '667788')

        # ── Slide 5: Preemptive vs Reactive ──
        s5 = prs.slides.add_slide(BLANK)
        box(s5, 0, 0, W, sy(12), '#1a3a5c')
        tb(s5, sx(3), sy(2), sw(94), sh(5), 'Preemptive vs Reactive: The Adversary Lifecycle Divide', 24, True, 'ffffff', 'center')

        # Left: Preemptive
        box(s5, sx(3), sy(15), sw(44), sh(70), 'f0fff0', '107c10')
        tb(s5, sx(5), sy(17), sw(40), sh(5), 'Preemptive Zone', 20, True, '107c10', 'center')
        tb(s5, sx(5), sy(23), sw(40), sh(3), 'Kill Chain Phases 1-3', 11, False, '555555', 'center')
        tb(s5, sx(5), sy(26), sw(40), sh(3), 'ATT&CK: Recon, Resource Dev., Initial Access', 10, False, '888888', 'center')
        tb(s5, sx(5), sy(32), sw(40), sh(8), '~78%', 30, True, '107c10', 'center')
        tb(s5, sx(5), sy(40), sw(40), sh(3), 'average vendor penetration', 10, False, '888888', 'center')
        preempt_pillars = [('Exposure Management', 'Strong', '107c10'), ('Posture & Policy Mgmt', 'Strong', '107c10'), ('Adversary Management', 'Weak', 'a80000')]
        for pi, (plbl, pstr, pcol) in enumerate(preempt_pillars):
            py = sy(46) + pi * sh(6)
            tb(s5, sx(8), py, sw(25), sh(5), plbl, 11, False, '333333')
            tb(s5, sx(35), py, sw(10), sh(5), pstr, 11, True, pcol, 'right')
        tb(s5, sx(5), sy(66), sw(40), sh(6), 'ATT&CK Tactics: TA0043, TA0042*, TA0001', 10, True, '107c10', 'center')
        tb(s5, sx(5), sy(72), sw(40), sh(4), '*Weak: depends on Adversary Management', 9, False, '888888', 'center')

        # Right: Reactive
        box(s5, sx(53), sy(15), sw(44), sh(70), 'fff0f0', 'a80000')
        tb(s5, sx(55), sy(17), sw(40), sh(5), 'Reactive Zone', 20, True, 'a80000', 'center')
        tb(s5, sx(55), sy(23), sw(40), sh(3), 'Kill Chain Phases 4-7', 11, False, '555555', 'center')
        tb(s5, sx(55), sy(26), sw(40), sh(3), 'ATT&CK: Execution through Impact', 10, False, '888888', 'center')
        tb(s5, sx(55), sy(32), sw(40), sh(8), '~68%', 30, True, 'a80000', 'center')
        tb(s5, sx(55), sy(40), sw(40), sh(3), 'Detection & Response penetration', 10, False, '888888', 'center')
        react_pillars = [('Adversary Disruption', 'Covers Phases 4-7', '8764b8')]
        for pi, (plbl, pstr, pcol) in enumerate(react_pillars):
            py = sy(46) + pi * sh(6)
            tb(s5, sx(58), py, sw(25), sh(5), plbl, 11, False, '333333')
            tb(s5, sx(85), py, sw(10), sh(5), pstr, 11, True, pcol, 'right')
        tb(s5, sx(55), sy(54), sw(40), sh(6), 'Services Maturity: Cross-cutting enabler', 10, True, 'ca5010', 'center')
        tb(s5, sx(55), sy(60), sw(40), sh(6), 'ATT&CK: TA0002-TA0003, TA0008-TA0011, TA0009-TA0010, TA0040', 10, True, 'a80000', 'center')
        tb(s5, sx(55), sy(66), sw(40), sh(4), '*Weak: Lateral Movement poorly addressed', 9, False, '888888', 'center')

        tb(s5, sx(3), sy(88), sw(94), sh(3), pillar_key_text, 9, False, '888888')

        # ── Slide 6: Strategic Roadmap ──
        s6 = prs.slides.add_slide(BLANK)
        box(s6, 0, 0, W, sy(12), '#1a3a5c')
        tb(s6, sx(3), sy(2), sw(94), sh(5), 'Strategic Roadmap: Full Adversary Lifecycle Coverage by 2028', 24, True, 'ffffff', 'center')

        roadmap = [
            ('2025-2026', 'Assess & Map', '0078d4', ['Dual-framework portfolio mapping', 'Identify KC + ATT&CK gaps', 'Develop Shift-Left Index', 'Competitive benchmarking']),
            ('2026-2027', 'Build Preemptive Core', '107c10', ['Adversary management investment', 'Polymorphic + credential rotation', 'Runtime protection + MTD', 'Posture validation (BAS, CSPM)']),
            ('2027-2028', 'Complete Lifecycle', '8764b8', ['Close services gap', 'Outcome-based delivery models', 'Full lifecycle positioning', 'Phase coverage as standard']),
        ]
        # Timeline line
        box(s6, sx(5), sy(50), sw(90), sh(0.5), 'e0ddd5')
        for ri, (year, title, col, items) in enumerate(roadmap):
            cx = sx(5) + ri * sw(32)
            box(s6, cx, sy(18), sw(29), sh(60), 'ffffff', col)
            box(s6, cx, sy(18), sw(29), sh(5), col)
            tb(s6, cx + sw(1), sy(19), sw(27), sh(4), year, 14, True, 'ffffff', 'center')
            tb(s6, cx + sw(1), sy(25), sw(27), sh(5), title, 18, True, col, 'center')
            for ii, item in enumerate(items):
                tb(s6, cx + sw(3), sy(34) + ii * sh(6), sw(25), sh(5), '\u2022 ' + item, 11, False, '555555')

        targets = [
            ('All 7 KC phases covered', '107c10'),
            ('All 14 ATT&CK tactics addressed', '0078d4'),
            ('Shift-Left Readiness Index', '8764b8'),
        ]
        for ti, (lbl, col) in enumerate(targets):
            bx = sx(5) + ti * sw(32)
            box(s6, bx, sy(82), sw(29), sh(5), col)
            tb(s6, bx + sw(1), sy(82.5), sw(27), sh(4), lbl, 10, True, 'ffffff', 'center')

        # ── Slide 7: Infographic Summary ──
        s7 = prs.slides.add_slide(BLANK)
        box(s7, 0, 0, W, H, 'fffef8')
        box(s7, 0, 0, W, sy(10), '1a3a5c')
        tb(s7, sx(3), sy(2), sw(94), sh(4), 'Dual Framework Defense: Kill Chain + MITRE ATT&CK', 22, True, 'ffffff', 'center')
        tb(s7, sx(3), sy(6), sw(94), sh(3), 'Preemptive Cybersecurity Pillar Mapping Across Both Adversary Models', 11, False, 'aac8e8', 'center')

        # Zone bar
        zones = [('PREEMPTIVE (Phases 1-3)', '107c10', 33), ('TRANSITION (Phase 4)', 'ca5010', 15), ('REACTIVE (Phases 5-7)', 'a80000', 47)]
        zx = sx(3)
        for zlbl, zcol, zpct in zones:
            zw = sw(zpct)
            box(s7, zx, sy(13), zw, sh(4), zcol)
            tb(s7, zx + sw(0.5), sy(13.5), zw - sw(1), sh(3), zlbl, 10, True, 'ffffff', 'center')
            zx += zw + sw(0.5)

        # Kill chain phases
        kcphases = [
            ('1. Recon', '107c10'), ('2. Weapon.', '107c10'), ('3. Delivery', '0078d4'), ('4. Exploit', 'ca5010'),
            ('5. Install', 'a80000'), ('6. C2', 'a80000'), ('7. Actions', '600000'),
        ]
        tb(s7, sx(3), sy(18), sw(10), sh(3), 'KILL CHAIN:', 9, True, '1a3a5c')
        for ki, (klbl, kcol) in enumerate(kcphases):
            kx = sx(14) + ki * sw(12)
            box(s7, kx, sy(18), sw(11), sh(4), kcol)
            tb(s7, kx + sw(0.5), sy(18.5), sw(10), sh(3), klbl, 9, True, 'ffffff', 'center')

        # ATT&CK row
        att_items = [
            ('TA0043', '107c10'), ('TA0042', '107c10'), ('TA0001', '0078d4'), ('TA0002/4', 'ca5010'),
            ('TA0003/5', 'a80000'), ('TA0011', 'a80000'), ('TA0010/40', '600000'),
        ]
        tb(s7, sx(3), sy(24), sw(10), sh(3), 'ATT&CK:', 9, True, '1a3a5c')
        for ai, (albl, acol) in enumerate(att_items):
            ax = sx(14) + ai * sw(12)
            box(s7, ax, sy(24), sw(11), sh(4), acol)
            tb(s7, ax + sw(0.5), sy(24.5), sw(10), sh(3), albl, 9, True, 'ffffff', 'center')

        # Pillar bars (4 defensive pillars only - SVC is cross-cutting)
        pillar_bars = [
            ('EXM', '107c10', 3, 40), ('AMT', 'd32f2f', 3, 55), ('PPM', '0078d4', 16, 55),
            ('ADR', '8764b8', 28, 95),
        ]
        tb(s7, sx(3), sy(30), sw(10), sh(3), 'PILLARS:', 9, True, '1a3a5c')
        for pi, (plbl, pcol, p_start, p_end) in enumerate(pillar_bars):
            py = sy(33) + pi * sh(4.5)
            box(s7, sx(p_start), py, sw(p_end - p_start), sh(3.5), pcol)
            tb(s7, sx(p_start) + sw(0.5), py + sh(0.3), sw(10), sh(3), plbl, 9, True, 'ffffff')

        # SVC cross-cutting enabler bar (spans full width)
        svc_y = sy(33) + 4 * sh(4.5)
        box(s7, sx(3), svc_y, sw(92), sh(3.5), 'fff3e0', 'ca5010')
        tb(s7, sx(3) + sw(0.5), svc_y + sh(0.3), sw(90), sh(3), 'SVC \u2022 Cross-Cutting Enabler (supports all pillars)', 9, True, 'ca5010', 'center')

        # Key findings
        kf_items = [
            ('Strongest', 'Recon (TA0043)', '107c10'),
            ('Weakest', 'Resource Dev. (TA0042)', 'a80000'),
            ('Convergence', 'Initial Access (TA0001)', '0078d4'),
            ('Market Gap', 'Platform-Only ~49%', 'ca5010'),
        ]
        tb(s7, sx(3), sy(58), sw(20), sh(3), 'KEY FINDINGS:', 10, True, '1a3a5c')
        for fi, (flbl, fval, fcol) in enumerate(kf_items):
            fx = sx(3) + fi * sw(24)
            box(s7, fx, sy(62), sw(22), sh(10), 'f8f8f5', fcol)
            box(s7, fx, sy(62), sw(22), sh(0.5), fcol)
            tb(s7, fx + sw(1), sy(63), sw(20), sh(3), flbl, 11, True, fcol, 'center')
            tb(s7, fx + sw(1), sy(67), sw(20), sh(3), fval, 10, False, '555555', 'center')

        # Shift-left box
        box(s7, sx(3), sy(75), sw(94), sh(10), 'f0fff0', '107c10')
        tb(s7, sx(5), sy(76), sw(90), sh(3), 'THE SHIFT-LEFT IMPERATIVE', 14, True, '107c10', 'center')
        tb(s7, sx(5), sy(80), sw(90), sh(4), 'Invest in the earliest adversary phases (KC 1-3 / ATT&CK TA0043, TA0042, TA0001) where interventions prevent attacks.', 11, False, '333333', 'center')

        # Pillar key
        box(s7, sx(3), sy(88), sw(94), sh(5), 'f8f8f5')
        tb(s7, sx(5), sy(89), sw(90), sh(4), pillar_key_text, 9, False, '888888', 'center')

        # Footer
        tb(s7, sx(3), sy(95), sw(94), sh(3), '\u00a9 Gartner Research \u2022 Preemptive Cybersecurity Dual Framework Analysis 2026', 8, False, 'aaaaaa', 'center')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='PreCyber_KillChain_MITRE_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


# ═══════════════════════════════════════════════════════════════════════
#  PMR Credibility Gap — All Graphics PPTX Export
# ═══════════════════════════════════════════════════════════════════════

@app.route('/api/pmr-all-graphics-pptx', methods=['GET'])
def pmr_all_graphics_pptx():
    """Generate a multi-slide editable PowerPoint deck with PMR credibility gap graphics."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from flask import send_file
        import io, statistics as stats_mod

        # Load vendor data
        vendor_file = os.path.join(os.path.dirname(__file__),
                                   'Product Market Readiness Vendor 1-0 Seed.json')
        vdata = read_dataset(
            'Product Market Readiness Vendor 1-0 Seed.json'
        )
        vendors = vdata.get('vendors', [])
        n = len(vendors)
        pillars = ['PPD', 'PCS', 'TDT', 'PCM', 'CTL']
        pillar_labels = {
            'PPD': 'Product Positioning & Differentiation',
            'PCS': 'Proof Points & Case Studies',
            'TDT': 'Technical Depth & Transparency',
            'PCM': 'Pricing & Commercial Model Clarity',
            'CTL': 'Content & Thought Leadership',
        }

        # Compute stats
        pillar_stats = {}
        for p in pillars:
            gtm_scores = [v.get('pillar_gtm_scores', {}).get(p, 0) for v in vendors]
            proof_scores = [v.get('pillar_proof_scores', {}).get(p, 0) for v in vendors]
            pillar_stats[p] = {
                'gtm_avg': round(stats_mod.mean(gtm_scores), 2) if gtm_scores else 0,
                'proof_avg': round(stats_mod.mean(proof_scores), 2) if proof_scores else 0,
                'gap_avg': round(stats_mod.mean(gtm_scores), 2) - round(stats_mod.mean(proof_scores), 2),
            }

        all_gtm_mean = round(stats_mod.mean([v.get('overall_gtm_score', 0) for v in vendors]), 2) if vendors else 0
        all_proof_mean = round(stats_mod.mean([v.get('overall_proof_score', 0) for v in vendors]), 2) if vendors else 0
        all_gap_mean = round(all_gtm_mean - all_proof_mean, 2)

        grade_dist = {}
        for v in vendors:
            g = v.get('coverage_grade', 'F')
            grade_dist[g] = grade_dist.get(g, 0) + 1

        # Over-claimers / best-aligned
        vendor_gaps = [(v.get('vendor', ''), v.get('overall_gtm_score', 0), v.get('overall_proof_score', 0),
                        v.get('overall_credibility_gap', 0), v.get('coverage_grade', 'F')) for v in vendors]
        over_claimers = sorted(vendor_gaps, key=lambda x: x[3], reverse=True)[:8]
        best_aligned = sorted(vendor_gaps, key=lambda x: abs(x[3]))[:8]

        # Vendor type stats
        type_counts = {}
        type_gtm = {}
        type_proof = {}
        for v in vendors:
            vt = v.get('vendor_type', 'Unknown')
            type_counts[vt] = type_counts.get(vt, 0) + 1
            if vt not in type_gtm:
                type_gtm[vt] = []
                type_proof[vt] = []
            type_gtm[vt].append(v.get('overall_gtm_score', 0))
            type_proof[vt].append(v.get('overall_proof_score', 0))

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        SW = prs.slide_width
        SH = prs.slide_height
        BLANK = prs.slide_layouts[6]

        def rgb(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        def box(slide, l, t, w, h, fill=None, border=None):
            s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
            if fill:
                s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
            else:
                s.fill.background()
            if border:
                s.line.color.rgb = rgb(border); s.line.width = Pt(1.5)
            else:
                s.line.fill.background()
            return s

        def tb(slide, l, t, w, h, text, sz=11, bold=False, col='333333', al='left'):
            tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
            tf = tx.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = str(text)
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = rgb(col)
            if al == 'center': p.alignment = PP_ALIGN.CENTER
            elif al == 'right': p.alignment = PP_ALIGN.RIGHT
            return tx

        def sx(v): return int(v / 1200 * SW)
        def sy(v): return int(v / 750 * SH)
        def sw(v): return int(v / 1200 * SW)
        def sh(v): return int(v / 750 * SH)

        # ── Slide 1: Title Slide ──
        s1 = prs.slides.add_slide(BLANK)
        box(s1, 0, 0, SW, SH, fill='0d1117')
        box(s1, sx(40), sy(30), sw(1120), sh(690), border='333333')
        tb(s1, sx(80), sy(80), sw(1040), sh(60), 'Product Market Readiness', sz=32, bold=True, col='60a5fa', al='center')
        tb(s1, sx(80), sy(150), sw(1040), sh(40), 'Credibility Gap Analysis — GTM Claims vs. Proof of Execution', sz=18, col='a78bfa', al='center')
        tb(s1, sx(80), sy(210), sw(1040), sh(40), f'{n} Vendors \u00b7 5 Pillars \u00b7 25 Sub-Pillars \u00b7 Dual Scoring Methodology', sz=14, col='999999', al='center')
        # Key metrics boxes
        metrics = [
            ('Vendors', str(n), '60a5fa'),
            ('Avg GTM', str(all_gtm_mean), '34d399'),
            ('Avg Proof', str(all_proof_mean), 'f59e0b'),
            ('Avg Gap', str(all_gap_mean), 'ef4444' if all_gap_mean > 0.5 else '34d399'),
        ]
        mw = 220; gap_x = 30; start_x = (1200 - (mw * 4 + gap_x * 3)) // 2
        for i, (label, val, color) in enumerate(metrics):
            mx = start_x + i * (mw + gap_x)
            box(s1, sx(mx), sy(300), sw(mw), sh(120), fill='1a1a2e', border='333333')
            tb(s1, sx(mx), sy(310), sw(mw), sh(60), val, sz=36, bold=True, col=color, al='center')
            tb(s1, sx(mx), sy(380), sw(mw), sh(30), label, sz=12, col='999999', al='center')
        # Pillar legend
        for i, p in enumerate(pillars):
            tb(s1, sx(80 + i * 210), sy(470), sw(200), sh(40), f'{p}: {pillar_labels[p]}', sz=10, col='cccccc')

        # ── Slide 2: Pillar Gap Comparison ──
        s2 = prs.slides.add_slide(BLANK)
        box(s2, 0, 0, SW, SH, fill='0d1117')
        tb(s2, sx(40), sy(20), sw(1120), sh(40), 'Slide 2: Pillar-Level Credibility Gap — GTM vs. Proof', sz=20, bold=True, col='60a5fa')
        bar_top = 90
        bar_h = 110
        for i, p in enumerate(pillars):
            ps = pillar_stats[p]
            y = bar_top + i * bar_h
            tb(s2, sx(40), sy(y), sw(300), sh(25), f'{p} — {pillar_labels[p]}', sz=12, bold=True, col='e6e6e6')
            # GTM bar
            gtm_w = max(10, int(ps['gtm_avg'] / 5 * 700))
            box(s2, sx(340), sy(y + 28), sw(gtm_w), sh(22), fill='34d399')
            tb(s2, sx(340 + gtm_w + 8), sy(y + 28), sw(80), sh(22), f'GTM {ps["gtm_avg"]}', sz=10, col='34d399')
            # Proof bar
            proof_w = max(10, int(ps['proof_avg'] / 5 * 700))
            box(s2, sx(340), sy(y + 55), sw(proof_w), sh(22), fill='f59e0b')
            tb(s2, sx(340 + proof_w + 8), sy(y + 55), sw(80), sh(22), f'Proof {ps["proof_avg"]}', sz=10, col='f59e0b')
            # Gap label
            gap_val = round(ps['gap_avg'], 2)
            gap_col = '34d399' if abs(gap_val) < 0.5 else 'f59e0b' if abs(gap_val) < 1.0 else 'ef4444'
            tb(s2, sx(1060), sy(y + 35), sw(100), sh(25), f'Gap: {"+" if gap_val > 0 else ""}{gap_val}', sz=11, bold=True, col=gap_col, al='right')

        # ── Slide 3: Coverage Grade Distribution ──
        s3 = prs.slides.add_slide(BLANK)
        box(s3, 0, 0, SW, SH, fill='0d1117')
        tb(s3, sx(40), sy(20), sw(1120), sh(40), 'Slide 3: Coverage Grade Distribution', sz=20, bold=True, col='60a5fa')
        tb(s3, sx(40), sy(60), sw(1120), sh(30), 'Letter grades based on percentage of 25 sub-pillars with non-zero GTM scores', sz=12, col='999999')
        grades = ['A', 'B', 'C', 'D', 'F']
        grade_colors = {'A': '34d399', 'B': '60a5fa', 'C': 'f59e0b', 'D': 'f97316', 'F': 'ef4444'}
        max_grade = max(1, max(grade_dist.get(g, 0) for g in grades))
        gbar_w = 160; gbar_gap = 30; gbar_start = (1200 - (gbar_w * 5 + gbar_gap * 4)) // 2
        for i, g in enumerate(grades):
            count = grade_dist.get(g, 0)
            bar_height = max(10, int(count / max_grade * 400))
            gx = gbar_start + i * (gbar_w + gbar_gap)
            gy = 600 - bar_height
            box(s3, sx(gx), sy(gy), sw(gbar_w), sh(bar_height), fill=grade_colors[g])
            tb(s3, sx(gx), sy(gy - 35), sw(gbar_w), sh(30), str(count), sz=20, bold=True, col=grade_colors[g], al='center')
            tb(s3, sx(gx), sy(610), sw(gbar_w), sh(30), f'Grade {g}', sz=14, bold=True, col=grade_colors[g], al='center')

        # ── Slide 4: Top Over-Claimers ──
        s4 = prs.slides.add_slide(BLANK)
        box(s4, 0, 0, SW, SH, fill='0d1117')
        tb(s4, sx(40), sy(20), sw(1120), sh(40), 'Slide 4: Top Over-Claimers — Largest Positive Credibility Gaps', sz=20, bold=True, col='ef4444')
        # Table header
        cols = [('Vendor', 400, 'left'), ('GTM', 120, 'center'), ('Proof', 120, 'center'), ('Gap', 120, 'center'), ('Grade', 120, 'center')]
        cx = 60
        for label, cw, ca in cols:
            box(s4, sx(cx), sy(80), sw(cw), sh(35), fill='1a1a2e')
            tb(s4, sx(cx + 5), sy(82), sw(cw - 10), sh(30), label, sz=11, bold=True, col='e6e6e6', al=ca)
            cx += cw + 5
        for ri, (vname, gtm, proof, gap_v, grade) in enumerate(over_claimers):
            ry = 120 + ri * 38
            cx = 60
            for val, cw, ca, vc in [(vname, 400, 'left', 'e6e6e6'), (str(gtm), 120, 'center', '34d399'),
                                     (str(proof), 120, 'center', 'f59e0b'), (f'+{gap_v}', 120, 'center', 'ef4444'),
                                     (grade, 120, 'center', 'cccccc')]:
                tb(s4, sx(cx + 5), sy(ry), sw(cw - 10), sh(30), val, sz=10, col=vc, al=ca)
                cx += cw + 5
            box(s4, sx(60), sy(ry + 32), sw(885), sh(1), fill='333333')

        # ── Slide 5: Best Aligned Vendors ──
        s5 = prs.slides.add_slide(BLANK)
        box(s5, 0, 0, SW, SH, fill='0d1117')
        tb(s5, sx(40), sy(20), sw(1120), sh(40), 'Slide 5: Best-Aligned Vendors — Smallest Credibility Gaps', sz=20, bold=True, col='34d399')
        cx = 60
        for label, cw, ca in cols:
            box(s5, sx(cx), sy(80), sw(cw), sh(35), fill='1a1a2e')
            tb(s5, sx(cx + 5), sy(82), sw(cw - 10), sh(30), label, sz=11, bold=True, col='e6e6e6', al=ca)
            cx += cw + 5
        for ri, (vname, gtm, proof, gap_v, grade) in enumerate(best_aligned):
            ry = 120 + ri * 38
            cx = 60
            gap_col = '34d399' if abs(gap_v) < 0.5 else 'f59e0b'
            for val, cw, ca, vc in [(vname, 400, 'left', 'e6e6e6'), (str(gtm), 120, 'center', '34d399'),
                                     (str(proof), 120, 'center', 'f59e0b'), (str(gap_v), 120, 'center', gap_col),
                                     (grade, 120, 'center', 'cccccc')]:
                tb(s5, sx(cx + 5), sy(ry), sw(cw - 10), sh(30), val, sz=10, col=vc, al=ca)
                cx += cw + 5
            box(s5, sx(60), sy(ry + 32), sw(885), sh(1), fill='333333')

        # ── Slide 6: Vendor Type Gap Comparison ──
        s6 = prs.slides.add_slide(BLANK)
        box(s6, 0, 0, SW, SH, fill='0d1117')
        tb(s6, sx(40), sy(20), sw(1120), sh(40), 'Slide 6: Credibility Gap by Vendor Type', sz=20, bold=True, col='a78bfa')
        vt_sorted = sorted(type_counts.items(), key=lambda x: x[1], reverse=True)
        vt_colors = ['60a5fa', '34d399', 'f59e0b', 'a78bfa', 'f472b6', 'f97316']
        card_w = 340; card_h = 150; card_gap = 20
        cards_per_row = 3
        for idx, (vt, cnt) in enumerate(vt_sorted[:6]):
            col_idx = idx % cards_per_row
            row_idx = idx // cards_per_row
            cx = 60 + col_idx * (card_w + card_gap)
            cy = 80 + row_idx * (card_h + card_gap)
            color = vt_colors[idx % len(vt_colors)]
            box(s6, sx(cx), sy(cy), sw(card_w), sh(card_h), fill='1a1a2e', border='333333')
            tb(s6, sx(cx + 12), sy(cy + 8), sw(card_w - 24), sh(25), vt, sz=13, bold=True, col=color)
            tb(s6, sx(cx + 12), sy(cy + 38), sw(card_w - 24), sh(40), str(cnt), sz=30, bold=True, col='e6e6e6')
            gtm_a = round(stats_mod.mean(type_gtm[vt]), 2) if type_gtm.get(vt) else 0
            proof_a = round(stats_mod.mean(type_proof[vt]), 2) if type_proof.get(vt) else 0
            gap_a = round(gtm_a - proof_a, 2)
            tb(s6, sx(cx + 12), sy(cy + 90), sw(card_w - 24), sh(20), f'GTM: {gtm_a}  |  Proof: {proof_a}', sz=10, col='cccccc')
            gap_col = '34d399' if abs(gap_a) < 0.5 else 'ef4444'
            tb(s6, sx(cx + 12), sy(cy + 115), sw(card_w - 24), sh(20), f'Gap: {"+" if gap_a > 0 else ""}{gap_a}', sz=11, bold=True, col=gap_col)

        # ── Slide 7: Methodology & Framework ──
        s7 = prs.slides.add_slide(BLANK)
        box(s7, 0, 0, SW, SH, fill='0d1117')
        tb(s7, sx(40), sy(20), sw(1120), sh(40), 'Slide 7: PMR Dual-Scoring Methodology', sz=20, bold=True, col='60a5fa')
        # Two-column layout: GTM scale left, Proof scale right
        box(s7, sx(40), sy(80), sw(540), sh(600), fill='1a1a2e', border='333333')
        tb(s7, sx(60), sy(90), sw(500), sh(30), 'GTM Messaging Score (0–5)', sz=16, bold=True, col='34d399')
        gtm_scale = [
            ('5', 'Market-Defining', 'Sets category narrative'),
            ('4', 'Compelling Narrative', 'Audience-aware, competitive context'),
            ('3', 'Specific Claims', 'Named features, differentiated'),
            ('2', 'Generic Positioning', 'Category-level, no differentiation'),
            ('1', 'Vague Mention', 'Buzzwords, no product tie-in'),
            ('0', 'No Messaging', 'No public content'),
        ]
        for i, (score, name, desc) in enumerate(gtm_scale):
            gy = 130 + i * 85
            tb(s7, sx(60), sy(gy), sw(40), sh(25), score, sz=22, bold=True, col='34d399')
            tb(s7, sx(110), sy(gy), sw(460), sh(25), name, sz=12, bold=True, col='e6e6e6')
            tb(s7, sx(110), sy(gy + 22), sw(460), sh(50), desc, sz=10, col='999999')

        box(s7, sx(620), sy(80), sw(540), sh(600), fill='1a1a2e', border='333333')
        tb(s7, sx(640), sy(90), sw(500), sh(30), 'Proof of Execution Score (0–5)', sz=16, bold=True, col='f59e0b')
        proof_scale = [
            ('5', 'Definitive Proof', 'Comprehensive, multi-source, replicable'),
            ('4', 'Strong Evidence', 'Multiple independent validations'),
            ('3', 'Demonstrated Proof', 'Named customers, third-party validation'),
            ('2', 'Basic Evidence', 'Internal case study, self-reported'),
            ('1', 'Minimal Proof', 'Single anecdote, outdated'),
            ('0', 'No Evidence', 'No verifiable proof'),
        ]
        for i, (score, name, desc) in enumerate(proof_scale):
            gy = 130 + i * 85
            tb(s7, sx(640), sy(gy), sw(40), sh(25), score, sz=22, bold=True, col='f59e0b')
            tb(s7, sx(690), sy(gy), sw(460), sh(25), name, sz=12, bold=True, col='e6e6e6')
            tb(s7, sx(690), sy(gy + 22), sw(460), sh(50), desc, sz=10, col='999999')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='PMR_Credibility_Gap_Graphics.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/innovation-profiles', methods=['GET'])
def get_innovation_profiles():
    """Return innovation profile(s). Optional ?id= filter."""
    json_file = os.path.join(os.path.dirname(__file__), 'innovation_profiles.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Innovation profiles file not found'}), 404
    try:
        data = read_dataset('innovation_profiles.json')
        profiles = data.get('profiles', [])
        profile_id = request.args.get('id', '')
        if profile_id:
            profile = next((p for p in profiles if p['id'] == profile_id), None)
            if not profile:
                return jsonify({'error': f'Profile "{profile_id}" not found'}), 404
            return jsonify(profile)
        return jsonify({'profiles': [{'id': p['id'], 'title': p['title'], 'schema': p.get('schema', '')} for p in profiles]})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/innovation-profiles', methods=['POST'])
def save_innovation_profile():
    """Save edits to an innovation profile."""
    json_file = os.path.join(os.path.dirname(__file__), 'innovation_profiles.json')
    if not os.path.exists(json_file):
        return jsonify({'error': 'Innovation profiles file not found'}), 404
    try:
        data = read_dataset('innovation_profiles.json')
        updated = request.get_json()
        if not updated or 'id' not in updated:
            return jsonify({'error': 'Missing profile id'}), 400
        profiles = data.get('profiles', [])
        idx = next((i for i, p in enumerate(profiles) if p['id'] == updated['id']), None)
        if idx is None:
            return jsonify({'error': f'Profile "{updated["id"]}" not found'}), 404
        profiles[idx] = updated
        write_error = persist_dataset('innovation_profiles.json', data)
        if write_error is not None:
            return write_error
        return jsonify({'status': 'ok'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyst-take-graphics-pptx', methods=['GET'])
def analyst_take_graphics_pptx():
    """Generate editable PowerPoint slides from Analyst Take SVG graphics."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from flask import send_file
        import io
        import xml.etree.ElementTree as ET
        import re as _re
        import html as _html

        json_file = os.path.join(os.path.dirname(__file__), 'analyst_take_reports.json')
        if not os.path.exists(json_file):
            return jsonify({'error': 'Analyst Take reports file not found'}), 404

        data = read_dataset('analyst_take_reports.json')

        perspective_id = request.args.get('perspective', '')
        reports = data.get('reports', [])
        report = None
        if perspective_id:
            report = next((r for r in reports if r['id'] == perspective_id), None)
        if not report:
            # Fall back to first report with graphics
            report = next((r for r in reports if r.get('graphics')), None)
        if not report or not report.get('graphics'):
            return jsonify({'error': 'No graphics found for this Analyst Take'}), 404

        graphics = report['graphics']
        report_title = report.get('title', 'Analyst Take')

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]
        SW = prs.slide_width
        SH = prs.slide_height

        def rgb(h):
            h = h.lstrip('#')
            if len(h) == 3:
                h = h[0]*2 + h[1]*2 + h[2]*2
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def _parse_color(val):
            """Parse hex color from SVG attribute; return RGBColor or None."""
            if not val:
                return None
            val = val.strip()
            if val.lower() in ('none', 'transparent', ''):
                return None
            if val.startswith('#'):
                return rgb(val)
            # Named colors
            named = {
                'white': 'ffffff', 'black': '000000', 'red': 'ff0000',
                'green': '008000', 'blue': '0000ff', 'yellow': 'ffff00',
                'gray': '808080', 'grey': '808080', 'orange': 'ffa500',
            }
            if val.lower() in named:
                return rgb(named[val.lower()])
            return None

        def _fval(el, attr, default=0):
            """Get a float attribute from an element."""
            v = el.get(attr)
            if v is None:
                return default
            try:
                return float(v)
            except (ValueError, TypeError):
                return default

        def _ival(el, attr, default=0):
            return int(_fval(el, attr, default))

        def _get_font_size(el):
            fs = el.get('font-size')
            if not fs:
                style = el.get('style', '')
                m = _re.search(r'font-size:\s*([\d.]+)', style)
                if m:
                    return float(m.group(1))
                return 11
            return float(fs.replace('px', '').replace('pt', ''))

        def _is_bold(el):
            fw = el.get('font-weight', '')
            if fw in ('bold', '700', '800', '900', '600'):
                return True
            style = el.get('style', '')
            if 'font-weight' in style:
                m = _re.search(r'font-weight:\s*(\w+)', style)
                if m and m.group(1) in ('bold', '700', '800', '900', '600'):
                    return True
            return False

        def _is_italic(el):
            fs = el.get('font-style', '')
            if fs == 'italic':
                return True
            style = el.get('style', '')
            return 'font-style:italic' in style or 'font-style: italic' in style

        def _text_anchor(el):
            ta = el.get('text-anchor', 'start')
            style = el.get('style', '')
            if 'text-anchor' in style:
                m = _re.search(r'text-anchor:\s*(\w+)', style)
                if m:
                    ta = m.group(1)
            return ta

        def _svg_to_slide(slide, svg_str, vb_w, vb_h, offset_left=0, offset_top=0, scale_w=1.0, scale_h=1.0):
            """Parse SVG and create native PPTX shapes on the slide."""
            svg_clean = svg_str.strip()
            # Strip style attributes from root <svg> that confuse ET (max-width etc)
            svg_clean = _re.sub(r'(<svg[^>]*)\s+style="[^"]*"', r'\1', svg_clean)

            try:
                root = ET.fromstring(svg_clean)
            except ET.ParseError:
                # Try wrapping in a root
                try:
                    root = ET.fromstring('<root>' + svg_clean + '</root>')
                except Exception:
                    return

            ns = {'svg': 'http://www.w3.org/2000/svg'}

            def _process_elements(parent):
                for el in parent:
                    tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag

                    if tag == 'rect':
                        _add_rect(el)
                    elif tag == 'text':
                        _add_text(el)
                    elif tag == 'circle':
                        _add_circle(el)
                    elif tag == 'line':
                        _add_line(el)
                    elif tag == 'path':
                        _add_path_as_line(el)
                    elif tag == 'g':
                        _process_elements(el)
                    elif tag in ('defs', 'marker', 'clipPath', 'style'):
                        pass  # skip
                    else:
                        # Recurse into unknown containers
                        if len(el) > 0:
                            _process_elements(el)

            def _sx(v):
                return int(offset_left + v * scale_w)

            def _sy(v):
                return int(offset_top + v * scale_h)

            def _sw(v):
                return int(v * scale_w)

            def _sh(v):
                return int(v * scale_h)

            def _add_rect(el):
                x = _fval(el, 'x', 0)
                y = _fval(el, 'y', 0)
                w = _fval(el, 'width', 0)
                h = _fval(el, 'height', 0)
                rx = _fval(el, 'rx', 0)
                if w <= 0 or h <= 0:
                    return

                if rx > 0:
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        _sx(x), _sy(y), _sw(w), _sh(h)
                    )
                    # Set corner radius as proportion
                    try:
                        adj = min(rx / min(w, h) * 2, 1.0)
                        shape.adjustments[0] = adj * 0.5
                    except Exception:
                        pass
                else:
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        _sx(x), _sy(y), _sw(w), _sh(h)
                    )

                fill_val = el.get('fill')
                fill_color = _parse_color(fill_val)
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill_color
                else:
                    shape.fill.background()

                stroke_val = el.get('stroke')
                stroke_color = _parse_color(stroke_val)
                if stroke_color:
                    shape.line.color.rgb = stroke_color
                    sw_val = _fval(el, 'stroke-width', 1)
                    shape.line.width = Pt(sw_val)
                    # Dashed lines
                    dash = el.get('stroke-dasharray', '')
                    if dash:
                        from pptx.oxml.ns import qn
                        from lxml.etree import SubElement as lxml_SubElement
                        ln = shape.line._ln
                        prstDash = ln.find(qn('a:prstDash'))
                        if prstDash is None:
                            prstDash = lxml_SubElement(ln, qn('a:prstDash'))
                        prstDash.set('val', 'dash')
                else:
                    shape.line.fill.background()

            def _add_text(el):
                x = _fval(el, 'x', 0)
                y = _fval(el, 'y', 0)
                # Get all text content (including nested tspan)
                text = ''.join(el.itertext()).strip()
                if not text:
                    return

                # Unescape HTML entities
                text = _html.unescape(text)

                font_size = _get_font_size(el)
                bold = _is_bold(el)
                italic = _is_italic(el)
                fill = _parse_color(el.get('fill', '#000000'))
                anchor = _text_anchor(el)

                # Estimate text width based on character count and font size
                char_w = font_size * 0.55
                text_w = max(len(text) * char_w, font_size * 3)
                text_h = font_size * 1.8

                # Position adjustment based on anchor
                if anchor == 'middle':
                    tx = x - text_w / 2
                elif anchor == 'end':
                    tx = x - text_w
                else:
                    tx = x

                # SVG y is baseline; shift up for PPTX top-left
                ty = y - font_size * 1.1

                # Check for rotation
                transform = el.get('transform', '')
                rotation = 0
                if 'rotate' in transform:
                    m = _re.search(r'rotate\(([-\d.]+)', transform)
                    if m:
                        rotation = float(m.group(1))

                txBox = slide.shapes.add_textbox(
                    _sx(tx), _sy(ty), _sw(text_w), _sh(text_h)
                )
                tf = txBox.text_frame
                tf.word_wrap = False
                tf.auto_size = None
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.italic = italic
                if fill:
                    p.font.color.rgb = fill
                p.font.name = 'Segoe UI'

                if anchor == 'middle':
                    p.alignment = PP_ALIGN.CENTER
                elif anchor == 'end':
                    p.alignment = PP_ALIGN.RIGHT

                if rotation != 0:
                    txBox.rotation = rotation

            def _add_circle(el):
                cx = _fval(el, 'cx', 0)
                cy = _fval(el, 'cy', 0)
                r = _fval(el, 'r', 5)

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    _sx(cx - r), _sy(cy - r), _sw(r * 2), _sh(r * 2)
                )

                fill_color = _parse_color(el.get('fill'))
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill_color
                else:
                    shape.fill.background()

                stroke_color = _parse_color(el.get('stroke'))
                if stroke_color:
                    shape.line.color.rgb = stroke_color
                    shape.line.width = Pt(_fval(el, 'stroke-width', 1))
                else:
                    shape.line.fill.background()

            def _add_line(el):
                x1 = _fval(el, 'x1', 0)
                y1 = _fval(el, 'y1', 0)
                x2 = _fval(el, 'x2', 0)
                y2 = _fval(el, 'y2', 0)

                connector = slide.shapes.add_connector(
                    1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                    _sx(x1), _sy(y1), _sx(x2), _sy(y2)
                )
                stroke_color = _parse_color(el.get('stroke'))
                if stroke_color:
                    connector.line.color.rgb = stroke_color
                    connector.line.width = Pt(_fval(el, 'stroke-width', 1))

                dash = el.get('stroke-dasharray', '')
                if dash:
                    from pptx.oxml.ns import qn
                    from lxml.etree import SubElement as lxml_SubElement
                    ln = connector.line._ln
                    prstDash = ln.find(qn('a:prstDash'))
                    if prstDash is None:
                        prstDash = lxml_SubElement(ln, qn('a:prstDash'))
                    prstDash.set('val', 'dash')

            def _add_path_as_line(el):
                """Convert simple SVG paths (M/L commands) to PPTX connectors."""
                d = el.get('d', '')
                if not d:
                    return
                # Extract M and L coordinates
                coords = _re.findall(r'([ML])\s*([-\d.]+)[\s,]+([-\d.]+)', d)
                if len(coords) < 2:
                    return
                # Draw line segments
                stroke_color = _parse_color(el.get('stroke'))
                for i in range(len(coords) - 1):
                    _, x1s, y1s = coords[i]
                    _, x2s, y2s = coords[i + 1]
                    x1, y1 = float(x1s), float(y1s)
                    x2, y2 = float(x2s), float(y2s)
                    connector = slide.shapes.add_connector(
                        1, _sx(x1), _sy(y1), _sx(x2), _sy(y2)
                    )
                    if stroke_color:
                        connector.line.color.rgb = stroke_color
                        connector.line.width = Pt(_fval(el, 'stroke-width', 1))

            _process_elements(root)

        # ── Build slides ──
        for idx, g in enumerate(graphics):
            slide = prs.slides.add_slide(BLANK)

            # Parse SVG viewBox for scaling
            svg_str = g.get('svg', '')
            vb_w, vb_h = 700, 420  # defaults
            m = _re.search(r'viewBox=["\'](\d+)\s+(\d+)\s+(\d+)\s+(\d+)["\']', svg_str)
            if m:
                vb_w = int(m.group(3))
                vb_h = int(m.group(4))

            # Slide layout:
            #  - Header bar (title): top 60px
            #  - SVG graphic: center, scaled to fill
            #  - Footer bar (caption + takeaway): bottom area

            header_h = Inches(0.85)
            footer_h = Inches(1.6)
            graphic_top = header_h
            graphic_h = SH - header_h - footer_h

            # Scale SVG to fit the graphic area (centered, maintain aspect ratio)
            svg_aspect = vb_w / vb_h if vb_h else 1.5
            avail_w = SW
            avail_h = graphic_h

            if avail_w / avail_h > svg_aspect:
                # Height-constrained
                render_h = avail_h
                render_w = int(avail_h * svg_aspect)
                offset_left = int((avail_w - render_w) / 2)
                offset_top = int(graphic_top)
            else:
                # Width-constrained
                render_w = avail_w
                render_h = int(avail_w / svg_aspect)
                offset_left = 0
                offset_top = int(graphic_top + (avail_h - render_h) / 2)

            scale_w = render_w / vb_w
            scale_h = render_h / vb_h

            # Header bar
            hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, header_h)
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = rgb('005a9e')
            hdr.line.fill.background()

            # Title text in header
            title_text = g.get('title', f'Graphic {idx + 1}')
            tb = slide.shapes.add_textbox(Inches(0.5), Emu(int(header_h * 0.1)), SW - Inches(1), Emu(int(header_h * 0.55)))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = rgb('ffffff')
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.CENTER

            # Purpose text in header
            purpose = g.get('purpose', '')
            if purpose:
                tb2 = slide.shapes.add_textbox(Inches(1), Emu(int(header_h * 0.58)), SW - Inches(2), Emu(int(header_h * 0.38)))
                tf2 = tb2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = purpose
                p2.font.size = Pt(11)
                p2.font.color.rgb = rgb('b8d4f0')
                p2.font.name = 'Segoe UI'
                p2.alignment = PP_ALIGN.CENTER

            # Render SVG elements as native shapes
            _svg_to_slide(slide, svg_str, vb_w, vb_h, offset_left, offset_top, scale_w, scale_h)

            # Footer area
            footer_top = SH - footer_h
            ftr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_top, SW, footer_h)
            ftr.fill.solid()
            ftr.fill.fore_color.rgb = rgb('f5f5f5')
            ftr.line.fill.background()

            # Caption
            caption = g.get('caption', '')
            if caption:
                ctb = slide.shapes.add_textbox(Inches(0.8), footer_top + Emu(Inches(0.15)), SW - Inches(1.6), Inches(0.7))
                ctf = ctb.text_frame
                ctf.word_wrap = True
                cp = ctf.paragraphs[0]
                cp.text = f'Figure {idx + 1}: {caption}'
                cp.font.size = Pt(10)
                cp.font.color.rgb = rgb('555555')
                cp.font.name = 'Segoe UI'
                cp.font.italic = True

            # Key Takeaway
            takeaway = g.get('takeaway', '')
            if takeaway:
                tk_top = footer_top + Inches(0.8)
                tk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), tk_top, SW - Inches(1.6), Inches(0.65))
                tk_box.fill.solid()
                tk_box.fill.fore_color.rgb = rgb('e8f0fe')
                tk_box.line.color.rgb = rgb('005a9e')
                tk_box.line.width = Pt(1)

                ttb = slide.shapes.add_textbox(Inches(1.0), tk_top + Emu(Inches(0.08)), SW - Inches(2.0), Inches(0.55))
                ttf = ttb.text_frame
                ttf.word_wrap = True
                tp = ttf.paragraphs[0]
                tp.text = 'KEY TAKEAWAY'
                tp.font.size = Pt(8)
                tp.font.bold = True
                tp.font.color.rgb = rgb('005a9e')
                tp.font.name = 'Segoe UI'
                tp2 = ttf.add_paragraph()
                tp2.text = takeaway
                tp2.font.size = Pt(10)
                tp2.font.color.rgb = rgb('333333')
                tp2.font.name = 'Segoe UI'

            # Slide number
            sn = slide.shapes.add_textbox(SW - Inches(1), SH - Inches(0.3), Inches(0.8), Inches(0.25))
            snf = sn.text_frame
            snp = snf.paragraphs[0]
            snp.text = f'{idx + 1}/{len(graphics)}'
            snp.font.size = Pt(8)
            snp.font.color.rgb = rgb('999999')
            snp.font.name = 'Segoe UI'
            snp.alignment = PP_ALIGN.RIGHT

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        title_slug = report_title.replace(' ', '_')[:40]
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name=f'Analyst_Take_Graphics_{title_slug}.pptx')
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/shutdown', methods=['POST'])
def shutdown_server():
    """Gracefully shut down the Flask server."""
    if (
        request.remote_addr not in {'127.0.0.1', '::1'}
        or os.getenv('ENABLE_LOCAL_SHUTDOWN', 'false').lower() != 'true'
    ):
        return jsonify({'error': 'Shutdown endpoint is disabled'}), 403
    func = request.environ.get('werkzeug.server.shutdown')
    if func is not None:
        func()
        return jsonify({'success': True, 'message': 'Server shutting down...'})
    # For newer Werkzeug versions, use os._exit
    import threading, os as _os
    def _exit():
        import time; time.sleep(0.5)
        _os._exit(0)
    threading.Thread(target=_exit, daemon=True).start()
    return jsonify({'success': True, 'message': 'Server shutting down...'})


# ── Blumira Vendor Deep-Dive Report ─────────────────────────────────────────
@app.route('/api/blumira-deep-dive', methods=['GET'])
def blumira_deep_dive():
    """Return Blumira data from every schema for the vendor deep-dive report."""
    result = {}

    # MDR
    try:
        mdr = read_dataset('MDR Services Vendor 2-1 Consolidated.json')
        b = next((v for v in mdr['vendors'] if v['vendor'] == 'Blumira'), None)
        result['mdr'] = b
    except Exception:
        result['mdr'] = None

    # MDR Pricing
    try:
        prc = read_dataset('MDR Services Vendor Pricing 2-1 AI Enriched.json')
        b = next((v for v in prc['vendors'] if v['vendor'] == 'Blumira'), None)
        result['mdr_pricing'] = b
    except Exception:
        result['mdr_pricing'] = None

    # Preemptive Cyber
    try:
        pc = read_dataset('Preemptive Cybersecurity Vendor 6-0 v3.json')
        b = next((v for v in pc['vendors'] if v['vendor'] == 'Blumira'), None)
        result['precyber'] = b
    except Exception:
        result['precyber'] = None

    # DFIR
    try:
        dfir = read_dataset('Vendor 3-7.json')
        b = next((v for v in dfir['vendors'] if v['vendor'] == 'Blumira'), None)
        result['dfir'] = b
    except Exception:
        result['dfir'] = None

    # PMR
    try:
        pmr = read_dataset('Product Market Readiness Vendor 1-0 Seed.json')
        b = next((v for v in pmr['vendors'] if v['vendor'] == 'Blumira'), None)
        result['pmr'] = b
    except Exception:
        result['pmr'] = None

    # Unmapped capabilities – things Blumira does that none of the schemas capture
    result['unmapped'] = {
        'vendor': 'Blumira',
        'capabilities': [
            {
                'id': 'UNM-01',
                'name': 'Compliance Reporting & Automation',
                'category': 'Governance, Risk & Compliance',
                'description': 'Built-in compliance reporting across 13+ regulatory frameworks including HIPAA, SOC 2, NIST CSF, PCI DSS, CMMC, CIS, and ISO 27001. Automated monitoring generates ready-made compliance reports and dashboards without manual evidence collection.',
                'evidence': 'Homepage: "Compliance and cyber insurance requirements built in with 24/7 automated monitoring, 365-day retention, and ready-made reports for 13+ frameworks." Compliance page details framework-specific report templates.',
                'maturity': 4,
                'schema_gap': 'MDR SOG-03 partially captures compliance alignment but not the automated report generation, multi-framework templating, or continuous compliance monitoring capabilities.',
                'market_category': 'GRC / Compliance Automation'
            },
            {
                'id': 'UNM-02',
                'name': 'MSP Multi-Tenant Platform & Channel Program',
                'category': 'Channel & Partner Ecosystem',
                'description': 'Purpose-built MSP program with multi-tenant management console, tiered partner pricing, dedicated partner support, and a platform designed for profitable resale. Enables MSPs to deliver enterprise-grade security to SMB clients without growing overhead.',
                'evidence': 'Homepage: "Deliver enterprise-grade security to all your clients with one platform that\'s easy to deploy, simple to manage, and profitable to sell." MSP program includes tiered pricing and dedicated partner support.',
                'maturity': 4,
                'schema_gap': 'No schema captures channel/partner ecosystem maturity, multi-tenant architecture, or MSP-specific go-to-market capabilities.',
                'market_category': 'MSP Security Platform'
            },
            {
                'id': 'UNM-03',
                'name': 'Executive ROI Reporting & Business Value Communication',
                'category': 'Security Program Communication',
                'description': 'Automated Executive Summaries showing leadership ROI metrics: cost savings, risk mitigation impact, and threat resolution statistics. Colorful, non-technical reports designed for board-level consumption.',
                'evidence': '"Automated Executive Summaries show leadership exactly how much you\'re saving, what risks you\'ve mitigated, and how threats are being resolved with colorful reports they\'ll actually understand."',
                'maturity': 3,
                'schema_gap': 'MDR SOG-04 touches reporting quality but none of the schemas assess executive communication, ROI quantification, or business-value storytelling capabilities.',
                'market_category': 'Security Analytics & BI'
            },
            {
                'id': 'UNM-04',
                'name': 'Cloud-Native Log Management & Long-Term Retention',
                'category': 'Data Infrastructure',
                'description': '365-day cloud-native log retention with flat-rate pricing independent of data volume. Eliminates the cost-vs-coverage tradeoff that forces organizations to disable critical security logs.',
                'evidence': '"Flat-rate pricing based on employee count, not data volume. So you don\'t have to choose between budget overages and turning off critical security logs." 365-day retention included across all editions.',
                'maturity': 4,
                'schema_gap': 'No schema evaluates log management architecture, retention policies, or data-volume-independent pricing models as distinct capabilities.',
                'market_category': 'Cloud SIEM / Log Management'
            },
            {
                'id': 'UNM-05',
                'name': 'Zero-Friction Deployment & Self-Service Onboarding',
                'category': 'Operational Model',
                'description': 'Deploy in hours (not weeks/months). No professional services required. Self-service integration setup with pre-built connectors. Designed for IT teams without dedicated security staff.',
                'evidence': '"Easy - Blumira\'s intuitive design ensures a hassle-free setup. With guided configurations, you\'re secured in hours, not days." G2 recognition for Easiest Setup in IDPS category. "Deploy integrations in minutes, not weeks."',
                'maturity': 5,
                'schema_gap': 'All schemas assume enterprise deployment models. No schema assesses time-to-value, self-service enablement, or low-friction operational models.',
                'market_category': 'SMB Security Operations'
            },
            {
                'id': 'UNM-06',
                'name': 'SOAR-Lite Orchestration',
                'category': 'Security Orchestration',
                'description': 'Lightweight security orchestration with automated response playbooks, dynamic blocklist management, and cross-platform response actions (endpoint isolation, identity lockout, firewall rules). G2 High Performer in SOAR category.',
                'evidence': 'G2 badge: "Security Orchestration, Automation, and Response (SOAR) Best Estimated ROI." Automated playbooks with guided response steps and one-click containment actions across endpoint, identity, and network.',
                'maturity': 3,
                'schema_gap': 'MDR TDR-03 captures response orchestration but not the standalone SOAR market positioning or the SMB-specific automation model that replaces enterprise SOAR platforms.',
                'market_category': 'SOAR / Response Automation'
            },
            {
                'id': 'UNM-07',
                'name': 'Cyber Insurance Alignment & Readiness',
                'category': 'Risk Transfer',
                'description': 'Platform capabilities explicitly designed to meet cyber insurance requirements: 24/7 monitoring evidence, incident response documentation, compliance reporting, and log retention that satisfy underwriter questionnaires.',
                'evidence': '"Compliance and cyber insurance requirements built in with 24/7 automated monitoring, 365-day retention, and ready-made reports."',
                'maturity': 3,
                'schema_gap': 'No schema evaluates cyber insurance readiness, underwriter alignment, or risk transfer enablement as a security capability dimension.',
                'market_category': 'Cyber Insurance Tech'
            },
            {
                'id': 'UNM-08',
                'name': 'Unified Security Visualization & Threat Scoping Console',
                'category': 'Security Analytics',
                'description': 'Single-pane visualization across events, networks, systems, and users. Enables non-specialist IT teams to scope threats like phishing and ransomware without raw log analysis. Customizable dashboards and pre-built views.',
                'evidence': '"Blumira visualizes your security data across events, networks, systems, and users in one console, making it easy to scope threats like phishing and ransomware without digging through raw logs."',
                'maturity': 3,
                'schema_gap': 'MDR TDR-02 covers investigation but not the democratized, non-specialist visualization layer that enables IT generalists to perform security triage.',
                'market_category': 'Security Analytics / Visualization'
            },
            {
                'id': 'UNM-09',
                'name': 'Cross-Signal Endpoint-Network-Cloud Correlation',
                'category': 'Detection Architecture',
                'description': 'Blumira correlates endpoint data with network, cloud, and identity signals to catch sophisticated multi-vector attacks that single-point solutions miss. This cross-domain correlation approach is architectural rather than capability-specific.',
                'evidence': '"Blumira correlates endpoint data with network, cloud, and identity signals to catch sophisticated attacks that single-point solutions miss."',
                'maturity': 4,
                'schema_gap': 'MDR TDR-01 assesses signal correlation but treats it as a sub-pillar. The architectural pattern of unified cross-domain telemetry correlation as a platform design principle is not captured.',
                'market_category': 'XDR Architecture'
            },
            {
                'id': 'UNM-10',
                'name': 'SMB-Optimized Security Economics',
                'category': 'Market Model',
                'description': 'Entire platform and pricing model optimized for organizations with 25-1000 employees. Flat per-employee pricing, no data volume penalties, no minimum commitments, no professional services requirements. Fundamentally different economic model than enterprise security.',
                'evidence': 'Pricing: $12/$16/$21 per employee/month for Detect/Respond/Automate. "30 min/week on average to manage Blumira." 99.4% faster average detection time vs industry average.',
                'maturity': 5,
                'schema_gap': 'All schemas evaluate capability maturity without assessing economic accessibility, operational burden, or market segment optimization as dimensions.',
                'market_category': 'SMB Security Platform'
            }
        ]
    }

    return jsonify(result)


@app.route('/api/asmf-framework', methods=['GET'])
def get_asmf_framework():
    """Return the selected framework schema data for framework views."""
    schema_file = request.args.get('schema', app_state.current_schema_file)

    # If the active schema is not framework-shaped, fall back to the default
    # ASMF schema so framework views can still render in a read-only mode.
    if schema_file not in SCHEMA_REGISTRY or SCHEMA_REGISTRY.get(schema_file, {}).get('structure') != 'asmf':
        schema_file = 'agentic_soc_framework_v1.json'

    if schema_file not in SCHEMA_REGISTRY or SCHEMA_REGISTRY.get(schema_file, {}).get('structure') != 'asmf':
        return jsonify({'error': 'No framework schema selected'}), 404
    try:
        data = read_dataset(schema_file)
        if isinstance(data, dict):
            data = dict(data)
            data['capabilities'] = _framework_capabilities(schema_file)
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/apef-report', methods=['GET'])
def get_apef_report():
    """Return the APEF framework and its source-native vendor profiles."""
    schema_file = 'AI_platform_ecosystem_framework_v1.json'
    vendor_file = 'ai_platform_ecosystem_vendors_v1.json'
    try:
        data = load_schema_data(schema_file)
        if not isinstance(data, dict):
            return jsonify({'error': 'APEF framework is unavailable'}), 404
        data = dict(data)
        profiles = data.get('vendor_role_profiles', {}) or {}
        vendors = []
        for record in load_vendor_data(vendor_file):
            vendor_name = record.get('vendor', '')
            vendor_key = record.get('key') or vendor_name.lower().replace(' ', '-')
            profile = dict(profiles.get(vendor_key, {}))
            profile.update(record)
            profile['key'] = vendor_key
            profile.setdefault('vendor', vendor_name or vendor_key)
            vendors.append(profile)
        data['vendors'] = vendors
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/apef-graph', methods=['GET'])
def get_apef_graph():
    """Build the APEF component graph from CouchDB-backed vendor profiles."""
    schema_file = 'AI_platform_ecosystem_framework_v1.json'
    vendor_file = 'ai_platform_ecosystem_vendors_v1.json'
    try:
        framework = load_schema_data(schema_file)
        if not isinstance(framework, dict):
            return jsonify({'error': 'APEF framework is unavailable'}), 404

        profiles = framework.get('vendor_role_profiles', {}) or {}
        vendors = []
        for record in load_vendor_data(vendor_file):
            vendor_name = record.get('vendor', '')
            vendor_key = record.get('key') or vendor_name.lower().replace(' ', '-')
            profile = dict(profiles.get(vendor_key, {}))
            profile.update(record)
            profile['key'] = vendor_key
            profile.setdefault('vendor', vendor_name or vendor_key)
            vendors.append(profile)

        nodes = []
        edges = []
        component_ids = set()
        component_meta = {}

        for vendor in vendors:
            vendor_key = vendor.get('key') or vendor.get('vendor', '').lower()
            if not vendor_key:
                continue
            vendor_id = f'vendor:{vendor_key}'
            nodes.append({
                'id': vendor_id,
                'kind': 'vendor',
                'vendor': vendor_key,
                'name': vendor.get('vendor') or vendor_key,
            })
            for component in vendor.get('components', []) or []:
                component_id = component.get('id')
                if not component_id:
                    continue
                component_ids.add(component_id)
                component_meta[component_id] = (vendor_key, component)
                nodes.append({
                    'id': component_id,
                    'kind': 'component',
                    'vendor': vendor_key,
                    'name': component.get('name') or component_id,
                    'type': component.get('type', ''),
                    'layer': component.get('layer', ''),
                })
                edges.append({
                    'source': vendor_id,
                    'target': component_id,
                    'kind': 'owns',
                    'integration_type': 'ownership',
                    'integration_label': 'Ownership / native component',
                })

        aliases = {
            'google-vertex': 'gcp-vertex',
            'microsoft-copilot-studio': 'copilot-studio',
        }
        integration_labels = {
            'nvidia_compute': 'NVIDIA compute / acceleration',
            'local_runtime': 'Local / on-prem runtime',
            'model_distribution': 'Model distribution / hosting',
            'agent_orchestration': 'Agent / orchestration',
            'data_grounding': 'Data / grounding / RAG',
            'governance_safety': 'Governance / safety',
            'platform_api': 'Platform / API control plane',
            'ecosystem_partner': 'Ecosystem / partner integration',
        }
        local_component_ids = {
            'nvidia-h100-local', 'nvidia-vllm', 'nvidia-triton-local',
            'nvidia-mlflow-local', 'nvidia-docker-compose-local',
        }

        def classify_integration(source_vendor, source, target_vendor, target):
            source_id = source.get('id', '')
            target_id = target.get('id', '')
            names = ' '.join([
                source.get('name', ''), target.get('name', ''), source_id, target_id,
            ]).lower()
            types = f"{source.get('type', '')} {target.get('type', '')}".lower()
            if source_id in local_component_ids or target_id in local_component_ids or 'local' in names:
                return 'local_runtime'
            if source_vendor == 'nvidia' or target_vendor == 'nvidia' or any(term in names for term in ('nvidia', 'cuda', 'triton', 'tensorrt')):
                return 'nvidia_compute'
            if source.get('layer') == 'L6' or target.get('layer') == 'L6' or any(term in names for term in ('guardrail', 'safety', 'governance', 'policy', 'iam', 'identity', 'model armor', 'saif', 'purview', 'entra')):
                return 'governance_safety'
            if any(term in names for term in ('bigquery', 'search', 'vector', 'rag', 'retrieval', 'grounding', 'knowledge', 'fabric', 'graph', 's3', 'redshift', 'opensearch', 'memory bank')):
                return 'data_grounding'
            if any(term in names for term in ('agent', 'copilot studio', 'assistant', 'mcp', 'tool use', 'runtime', 'orchestration', 'strands', 'semantic kernel', 'autogen')):
                return 'agent_orchestration'
            if source.get('layer') == 'L3' or target.get('layer') == 'L3' or 'foundation-model' in types or any(term in names for term in ('model garden', 'bedrock', 'azure openai', 'vertex', 'foundry', 'chatgpt', 'claude', 'gemini', 'gpt')):
                return 'model_distribution'
            if source.get('layer') == 'L4' or target.get('layer') == 'L4' or any(term in names for term in ('api', 'platform', 'pipeline', 'mlflow', 'sagemaker')):
                return 'platform_api'
            return 'ecosystem_partner' if source_vendor != target_vendor else 'platform_api'

        integration_types = framework.get('integration_type_taxonomy', []) or []
        for vendor_key, component in component_meta.values():
            source_id = component.get('id')
            for raw_target_id in component.get('integrates_with', []) or []:
                target_id = aliases.get(raw_target_id, raw_target_id)
                if source_id and target_id in component_ids:
                    target_vendor, target_component = component_meta[target_id]
                    integration_type = classify_integration(
                        vendor_key, component, target_vendor, target_component,
                    )
                    edges.append({
                        'source': source_id,
                        'target': target_id,
                        'kind': 'integrates',
                        'integration_type': integration_type,
                        'integration_label': integration_labels[integration_type],
                    })

        layers = ((framework.get('enterprise_stack_lens') or {}).get('layers') or [])
        return jsonify({
            'vendors': [
                {'key': vendor.get('key'), 'vendor': vendor.get('vendor')}
                for vendor in vendors
            ],
            'layers': layers,
            'integration_type_taxonomy': integration_types,
            'nodes': nodes,
            'edges': edges,
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/asmf-orbital-map', methods=['GET'])
def get_asmf_orbital_map():
    """Return the orbital map compatible with the selected framework."""
    try:
        schema_file = request.args.get('schema', app_state.current_schema_file)
        if schema_file == 'agentic_enterprise_operations_framework_v1.json':
            framework = load_schema_data(schema_file)
            dimensions = framework.get('dimensions', {}) if isinstance(framework, dict) else {}
            relationship_types = {
                'data_flow': {'color': '#06b6d4', 'label': 'Data Flow', 'abbr': 'DF', 'dash': []},
                'execution': {'color': '#8b5cf6', 'label': 'Execution', 'abbr': 'EX', 'dash': []},
                'feedback': {'color': '#10b981', 'label': 'Feedback', 'abbr': 'FB', 'dash': [4, 4]},
                'governance': {'color': '#ef4444', 'label': 'Governance', 'abbr': 'GV', 'dash': [6, 3]},
                'augmentation': {'color': '#f59e0b', 'label': 'Augmentation', 'abbr': 'AU', 'dash': [3, 6]},
            }
            relationships = [
                ('OBS', 'RPL', 'Operational telemetry and dependency context initiate reasoning and planning.', 'data_flow', 3),
                ('OBS', 'OKG', 'Live signals enrich the operational graph with current service, asset, and dependency state.', 'data_flow', 3),
                ('OBS', 'AMS', 'Signal coverage, alert quality, and service health feed operational assurance metrics.', 'feedback', 2),
                ('RPL', 'EXE', 'Causal reasoning and prioritized plans direct governed remediation and fulfillment actions.', 'execution', 3),
                ('RPL', 'OPM', 'Plans shape dynamic workflow branches, service-risk prioritization, and coordinated work.', 'execution', 3),
                ('RPL', 'OKG', 'Reasoning uses graph-linked service context, evidence, policy, and workflow history.', 'data_flow', 3),
                ('EXE', 'OBS', 'Execution outcomes change service state and generate new sensing signals.', 'feedback', 2),
                ('EXE', 'AMS', 'Remediation success, reversals, and operational impact feed effectiveness assurance.', 'feedback', 3),
                ('EXE', 'AGC', 'Specialized agents coordinate bounded execution across operational domains.', 'execution', 3),
                ('POL', 'EXE', 'Machine-interpretable authority, policy, and guardrails bound every execution action.', 'governance', 3),
                ('POL', 'AGC', 'Policy defines agent roles, authority scopes, escalation boundaries, and controls.', 'governance', 3),
                ('POL', 'HGI', 'Governance establishes human intent, approval, accountability, and exception authority.', 'governance', 3),
                ('POL', 'AMS', 'Policy defines audit evidence, compliance thresholds, and assurance requirements.', 'governance', 2),
                ('CIL', 'OKG', 'Validated outcomes and external knowledge continuously evolve the operational graph.', 'augmentation', 3),
                ('CIL', 'RPL', 'Learning from incidents and workflow outcomes improves planning and prioritization.', 'feedback', 3),
                ('CIL', 'AGC', 'Continuous improvement updates agent skills, coordination patterns, and novelty handling.', 'augmentation', 2),
                ('OPM', 'EXE', 'The interaction model orchestrates non-linear execution across incidents, changes, and requests.', 'execution', 3),
                ('OPM', 'HGI', 'Workflow design defines where human intent, collaboration, and exceptions enter operations.', 'augmentation', 2),
                ('HGI', 'POL', 'Human intent and operational accountability inform policy and authority design.', 'governance', 2),
                ('HGI', 'AGC', 'Humans supervise agent behavior, resolve novel conditions, and govern operating intent.', 'governance', 3),
                ('AGC', 'OBS', 'Agents adapt sensing coverage and operational signal priorities as conditions change.', 'data_flow', 2),
                ('AGC', 'OKG', 'Agents use and produce graph-linked operational evidence, plans, and workflow artifacts.', 'data_flow', 3),
                ('OKG', 'RPL', 'The operational graph provides the shared context required for causal reasoning and planning.', 'data_flow', 3),
                ('OKG', 'EXE', 'Graph-linked dependencies, authority, and evidence guide safe execution.', 'execution', 2),
                ('AMS', 'POL', 'Assurance and audit findings validate policy effectiveness and risk boundaries.', 'feedback', 2),
                ('AMS', 'TRF', 'Operational performance and assurance gaps identify transformation priorities.', 'feedback', 3),
                ('TRF', 'AGC', 'Transformation readiness enables expanded agent coordination and operating-model adoption.', 'augmentation', 2),
                ('TRF', 'OPM', 'Transformation investments reshape the operational interaction model and workflow architecture.', 'augmentation', 2),
            ]
            return jsonify({
                'relationship_types': relationship_types,
                'dim_config': {
                    code: {
                        'plane': dimension.get('plane', ''),
                        'short': code,
                    }
                    for code, dimension in dimensions.items()
                },
                'relationships': [
                    {'from': source, 'to': target, 'label': label, 'type': relation_type, 'strength': strength}
                    for source, target, label, relation_type, strength in relationships
                    if source in dimensions and target in dimensions
                ],
            })
        if schema_file != 'agentic_soc_framework_v1.json':
            framework = load_schema_data(schema_file)
            dimensions = framework.get('dimensions', {}) if isinstance(framework, dict) else {}
            return jsonify({
                'relationship_types': {},
                'dim_config': {
                    code: {
                        'plane': dimension.get('plane', ''),
                        'short': code,
                    }
                    for code, dimension in dimensions.items()
                },
                'relationships': [],
            })
        data = read_dataset('static/asmf_orbital_map.json')
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# Allowed doc IDs to prevent path traversal
_ALLOWED_DOCS = {'precyber_methodology', 'architecture'}

@app.route('/api/docs/<doc_id>', methods=['GET'])
def get_docs(doc_id):
    """Serve a docs JSON file from static/ by safe doc_id."""
    if doc_id not in _ALLOWED_DOCS:
        return jsonify({'error': 'Not found'}), 404
    relative_path = f'static/docs_{doc_id}.json'
    try:
        data = read_dataset(relative_path)
        return jsonify(data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


def create_app():
    """Return the configured Flask gateway application."""
    return app


if __name__ == '__main__':
    import socket, sys, subprocess, signal

    PORT = int(sys.argv[sys.argv.index('--port') + 1]) if '--port' in sys.argv else 5000

    # ── Graceful pre-launch: stop any existing server on this port ──
    def _kill_existing(port):
        """Ask the existing server to shut down, then force-kill stragglers."""
        import urllib.request
        # 1. Try the graceful /api/shutdown endpoint
        try:
            req = urllib.request.Request(
                f'http://127.0.0.1:{port}/api/shutdown', method='POST'
            )
            urllib.request.urlopen(req, timeout=3)
            print(f"[startup] Sent shutdown to existing server on :{port}")
            time.sleep(2)
        except Exception:
            pass  # Nothing running, or it already died

        # 2. Force-kill any remaining PIDs bound to this port (Windows)
        if sys.platform == 'win32':
            try:
                out = subprocess.check_output(
                    f'netstat -ano | findstr ":{port}"',
                    shell=True, text=True, stderr=subprocess.DEVNULL
                )
                pids = set()
                my_pid = os.getpid()
                for line in out.strip().splitlines():
                    parts = line.split()
                    if 'LISTENING' in parts:
                        pid = int(parts[-1])
                        if pid != my_pid and pid != 0:
                            pids.add(pid)
                for pid in pids:
                    try:
                        subprocess.run(
                            f'taskkill /F /PID {pid}',
                            shell=True, capture_output=True, timeout=5
                        )
                        print(f"[startup] Killed lingering PID {pid}")
                    except Exception:
                        pass
                if pids:
                    time.sleep(2)
            except Exception:
                pass

        # 3. Verify port is free
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            try:
                s.bind(('0.0.0.0', port))
            except OSError as e:
                print(f"[startup] WARNING: port {port} still in use — {e}")
                print(f"[startup] Try:  taskkill /F /IM python.exe  (run as Admin)")
                sys.exit(1)
        print(f"[startup] Port {port} is free — starting server.")

    _kill_existing(PORT)

    from werkzeug.serving import WSGIRequestHandler
    WSGIRequestHandler.protocol_version = "HTTP/1.1"
    app.run(host='0.0.0.0', debug=False, use_reloader=False, port=PORT)
